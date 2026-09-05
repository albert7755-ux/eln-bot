import knowledge
import base64 as _base64
from fastapi import Form
import os
import re
import json
import traceback as _traceback
from pathlib import Path
from datetime import datetime, timezone, timedelta, date
from fastapi import FastAPI, Request, HTTPException, UploadFile, File
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from linebot import LineBotApi, WebhookHandler
from linebot.exceptions import InvalidSignatureError
from linebot.models import MessageEvent, TextMessage, TextSendMessage, FileMessage, ImageMessage, AudioMessage
from sqlalchemy import create_engine, text
from sqlalchemy import text as sql_text  # 別名:避免被函式內的區域變數 text 遮蔽
from autotracking_core import calculate_from_file
from market_content_generator import generate_market_content
import anthropic
from apscheduler.schedulers.background import BackgroundScheduler
from apscheduler.triggers.cron import CronTrigger
from apscheduler.triggers.interval import IntervalTrigger
import pytz
import urllib.request
import urllib.error
from openai import OpenAI

# ===== 海外債配息雷達（bond_coupon_alert.py）=====
try:
    from bond_coupon_alert import build_alert_message as _bond_build_alert, is_bond_pricing_file as _is_bond_pricing_file
    _BOND_RADAR_OK = True
except Exception as _e:
    print(f"[BondRadar] 模組載入失敗（不影響其他功能）：{_e}")
    _BOND_RADAR_OK = False

def _bond_price_dir() -> Path:
    """報價檔存放位置：優先用 Render 持久磁碟 /data，沒有就退到 /tmp（重新部署會消失）"""
    for d in (Path("/data/bond_pricing"), Path("/tmp/bond_pricing")):
        try:
            d.mkdir(parents=True, exist_ok=True)
            return d
        except Exception:
            continue
    return Path("/tmp")

BOND_PRICE_FILE = Path(os.getenv("BOND_PRICE_FILE", "")) if os.getenv("BOND_PRICE_FILE") else (_bond_price_dir() / "bond_pricing_latest.xlsx")

# 海外債群組白名單：在該群裡龍蝦只回這些指令，其餘一律不理（比照 ELN 群）
BOND_GROUP_ALLOWED_CMDS = {"coupon", "issuer", "bondalert", "rating", "move", "price", "bid", "sheet", "focus", "help"}

# ELN Bot(callback2)也開放的海外債指令:查詢類為主,燒 AI 的重型指令僅限 Albert 本人
ELN_BOT_BOND_CMDS = {"price", "p", "價格", "報價", "bid", "賣回", "贖回",
                     "issuer", "coupon", "move", "bondalert", "myid", "我的id"}
ELN_BOT_BOND_HEAVY = {"sheet", "focus", "rating"}   # 需要 Albert 本人才可用

def is_bond_group_chat(chat_key: str) -> bool:
    """這個 chat 是不是用 /coupon settarget 設定的海外債群組"""
    try:
        t = load_targets() or {}
        return bool(t.get("bond")) and t.get("bond_type") in ("group", "room") and chat_key.split(":", 1)[1] == t.get("bond")
    except Exception:
        return False

def get_doc_users():
    """可使用 /sheet /focus 等產文件指令的使用者清單(投資輔銷)"""
    try:
        return list((load_targets() or {}).get("doc_users", []))
    except Exception:
        return []


def set_doc_users(uids):
    t = load_targets() or {}
    t["doc_users"] = list(dict.fromkeys(uids))
    save_targets(t)


def can_use_doc_cmd(uid):
    """Albert 本人或名單內的使用者才能用產文件指令"""
    if not uid:
        return False
    if uid == os.getenv("LINE_USER_ID", ""):
        return True
    return uid in get_doc_users()


def get_bond_query_groups():
    try:
        return list((load_targets() or {}).get("bond_query_groups", []))
    except Exception:
        return []

def is_bond_query_group(chat_key: str) -> bool:
    """這個群是不是「只開放查價」的海外債查詢群（/price settarget 設定，可多個群）"""
    try:
        return chat_key.split(":", 1)[1] in get_bond_query_groups()
    except Exception:
        return False

BOND_QUERY_HELP = (
    "❤️ 愛債每一天｜債查一下 Claw Bot\n"
    "━━━━━━━━━━━━━\n"
    "🔍 查報價\n"
    "/price 26070003 → 產品代碼查（WMBB免打）\n"
    "/price 蘋果 2043 → 名稱＋到期年份\n"
    "/price US037833EN → ISIN查\n"
    "\n📈 查價格走勢\n"
    "/price 26070003 30 → 該檔近30天報價變化\n"
    "\n💰 查賣回價\n"
    "/bid 26070003 → 賣回價(Bid)與買賣價差\n"
    "（含期間漲跌幅與最高最低點，天數可自訂）\n"
    "\n🏦 查發行機構\n"
    "/issuer 蘋果 → 機構簡介＋架上所有債券\n"
    "（可用中文、英文、ISIN、產品代碼）\n"
    "\n📅 查配息\n"
    "/coupon → 3個營業日內截止的配息債\n"
    "/coupon 7 → 改看7個營業日內\n"
    "/coupon all → 未來14天全部\n"
    "（另附「剛配息完、前手息最低」名單）\n"
    "━━━━━━━━━━━━━\n"
    "🔒專投＝限專業投資人｜💎高資產＝高資產客戶專屬\n"
    "報價以總行系統為準，商品條件依產品說明書"
)

BOND_GROUP_HELP = (
    "❤️ 愛債每一天｜債查一下 Claw Bot\n"
    "海外債專區指令\n"
    "━━━━━━━━━━━━━━━\n"
    "⏰ 每天自動推播\n"
    "06:38 債券市場日報　06:45 配息雷達　06:50 信評新聞\n"
    "\n📅 配息雷達\n"
    "/coupon → 3個營業日內截止的配息債＋30天內到期\n"
    "/coupon 7 → 改看7個營業日內\n"
    "/coupon all → 未來14天全部\n"
    "/coupon table → Excel條件表＋機構簡介（Drive連結）\n"
    "\n💵 報價查詢\n"
    "/price 26070003 → 產品代碼查（WMBB免打）\n"
    "/price 蘋果 2043 → 名稱＋到期年份\n"
    "/price US037833EN → ISIN查\n"
    "/price 26070003 30 → 該檔近30天走勢\n"
    "\n📊 異動與追蹤\n"
    "/move → 全架 vs 上一份報價，變動≥1%\n"
    "/move 7 3 → vs 7天前，≥3%\n"
    "/bondalert 蘋果 2043 ytm>5.2 → 單檔到價通知\n"
    "/bondalert list　/bondalert del 3\n"
    "\n🏦 發行機構\n"
    "/issuer 蘋果 → 簡介＋架上所有債券\n"
    "/sheet 蘋果 → 參考資訊（信評+財務+圖表+標的，含PDF）\n"
    "/focus 輝達 → 債市每日聚焦 PPTX（直式兩頁，可編輯）\n"
    "/stock intc → 個股完整分析報告（商業模式+財務+護城河+估值+多空辯論，文字摘要+PDF；輔銷可用）\n"
    "\n🚨 信評\n"
    "/rating → 立即掃外部信評新聞\n"
    "/rating list　/rating watch 台積電　/rating unwatch 蘋果\n"
    "\n📥 報價檔\n"
    "傳 Bond_Pricing Excel → 更新報價＋自動跑配息雷達、\n"
    "　信評異動比對、報價異動、到價檢查\n"
    "傳舊日期的報價檔 → 只補歷史，不動最新報價\n"
    "\n📣 推播與維護\n"
    "/coupon subscribe → 訂閱每日推播（unsubscribe取消）\n"
    "/coupon settarget → 本群收每日推播（off取消）\n"
    "/price settarget → 本群只開放查價（off取消）\n"
    "/cleanup → 預覽Drive舊報告；/cleanup do → 清理\n"
    "/econ check → 立即檢查重要數據/央行會議是否剛公布\n"
    "/econ list → 看追蹤清單（自動偵測公布會直接推播）\n"
    "/econ calendar → 看本月事件確切日期（refresh強制更新）\n"
    "/sheetuser list → 產文件名單（add/del；對方先打 /myid 取得ID）\n"
    "\n📈 日報\n"
    "/bonddaily → 立即產生債券市場日報\n"
    "/bonddaily cache → 看最近一份\n"
    "/bonddaily focus 零息債 | 風險:無配息、有提前買回 → 設定當期主打方向與必講風險（off取消）\n"
    "━━━━━━━━━━━━━━━\n"
    "🔒專投＝限專業投資人｜💎高資產＝高資產客戶專屬\n"
    "報價以總行系統為準，商品條件依產品說明書"
)


BOND_SNAPSHOT_FILE = BOND_PRICE_FILE.parent / "bond_snapshot.json"

def parse_pricing_file_date(filename):
    """
    從報價檔檔名解析報價日期,支援 08-20-2026 / 2026-08-20 / 20260820 等格式。
    解析失敗回 None(呼叫端改用今天)。
    """
    st = str(filename)
    m = re.search(r"(\d{2})[-_.](\d{2})[-_.](\d{4})", st)   # MM-DD-YYYY
    if m:
        mm, dd, yyyy = int(m.group(1)), int(m.group(2)), int(m.group(3))
        try:
            return date(yyyy, mm, dd)
        except ValueError:
            pass
    m = re.search(r"(\d{4})[-_.](\d{2})[-_.](\d{2})", st)   # YYYY-MM-DD
    if m:
        try:
            return date(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            pass
    m = re.search(r"(20\d{2})(\d{2})(\d{2})", st)            # YYYYMMDD
    if m:
        try:
            return date(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            pass
    return None

def save_price_history(snap_date=None, path=None):
    """把報價檔的 Offer / Bid / YTM 存進 bond_price_history（同一天重覆上傳會覆蓋）。path 未指定時用最新報價檔"""
    from bond_coupon_alert import read_bonds, first_num
    snap_date = snap_date or datetime.now(TZ_TAIPEI).date()
    rows = []
    for b in read_bonds(str(path or BOND_PRICE_FILE)):
        rows.append({"d": snap_date, "i": b["isin"], "n": b["name"], "c": b["ccy"],
                     "o": first_num(b["offer"]), "bd": first_num(b.get("bid")),
                     "y": first_num(b["ytm"]), "m": b["maturity"]})
    with engine.begin() as conn:
        conn.execute(text("""INSERT INTO bond_price_history(snap_date, isin, bond_name, ccy, offer, bid, ytm, maturity)
                             VALUES (:d,:i,:n,:c,:o,:bd,:y,:m)
                             ON CONFLICT (snap_date, isin) DO UPDATE SET bond_name=EXCLUDED.bond_name, ccy=EXCLUDED.ccy,
                             offer=EXCLUDED.offer, bid=EXCLUDED.bid, ytm=EXCLUDED.ytm, maturity=EXCLUDED.maturity"""), rows)
    return len(rows)

def price_movers(days_back=1, threshold_pct=2.0, top_n=15):
    """
    最新快照 vs N 天前（取 ≤ 該日最近的一份快照）Offer 變動 ≥ threshold% 的債券。
    days_back=1 代表「上一份快照」（通常是昨天）。回傳 (text, base_date, latest_date)。
    """
    from bond_coupon_alert import pi_tag
    with engine.begin() as conn:
        latest = conn.execute(text("SELECT MAX(snap_date) FROM bond_price_history")).scalar()
        if not latest:
            return "", None, None
        if days_back <= 1:
            base = conn.execute(text("SELECT MAX(snap_date) FROM bond_price_history WHERE snap_date < :l"), {"l": latest}).scalar()
        else:
            target = latest - timedelta(days=days_back)
            base = conn.execute(text("SELECT MAX(snap_date) FROM bond_price_history WHERE snap_date <= :t"), {"t": target}).scalar()
            if base is None:  # 歷史不夠長，就拿最早那一份
                base = conn.execute(text("SELECT MIN(snap_date) FROM bond_price_history WHERE snap_date < :l"), {"l": latest}).scalar()
        if not base or base == latest:
            return "", base, latest
        rows = conn.execute(text("""
            SELECT n.isin, n.bond_name, n.ccy, o.offer, n.offer, o.ytm, n.ytm
            FROM bond_price_history n JOIN bond_price_history o ON o.isin = n.isin AND o.snap_date = :b
            WHERE n.snap_date = :l
              AND n.offer IS NOT NULL AND o.offer IS NOT NULL
              AND n.offer > 1 AND o.offer > 1
        """), {"b": base, "l": latest}).fetchall()
    movers = []
    for isin, name, ccy, o_off, n_off, o_ytm, n_ytm in rows:
        chg = (n_off - o_off) / o_off * 100
        if abs(chg) > 30:      # 單日/單週逾30%多為報價缺漏或資料異常,不列入
            print(f"[BondMovers] 略過異常變動 {name} {o_off}→{n_off} ({chg:+.0f}%)")
            continue
        if abs(chg) >= threshold_pct:
            movers.append((chg, isin, name, ccy, o_off, n_off, o_ytm, n_ytm))
    if not movers:
        return "", base, latest
    # 用最新報價檔補專投標籤
    tags = {}
    try:
        from bond_coupon_alert import read_bonds
        tags = {b["isin"]: pi_tag(b) for b in read_bonds(str(BOND_PRICE_FILE))}
    except Exception:
        pass
    downs = sorted([m for m in movers if m[0] < 0], key=lambda x: x[0])[:top_n]
    ups = sorted([m for m in movers if m[0] > 0], key=lambda x: -x[0])[:top_n]
    span = "上一份報價" if days_back <= 1 else f"{days_back} 天前"
    lines = [f"📊 報價異動 ≥ {threshold_pct:g}%（{base:%m/%d} → {latest:%m/%d}，vs {span}）",
             f"跌 {len([m for m in movers if m[0] < 0])} 檔｜漲 {len([m for m in movers if m[0] > 0])} 檔"]
    def _fmt(m):
        chg, isin, name, ccy, o_off, n_off, o_ytm, n_ytm = m
        _valid = (o_ytm is not None and n_ytm is not None
                  and 0 < o_ytm <= 25 and 0 < n_ytm <= 25)
        y = f"｜YTM {round(o_ytm,2):g}→{round(n_ytm,2):g}" if _valid else ""
        return f"▪ {name} {ccy}｜{tags.get(isin, '')}\n  Offer {o_off:g}→{n_off:g}（{chg:+.1f}%）{y}"
    if downs:
        lines.append("\n📉 跌幅")
        lines += [_fmt(m) for m in downs]
    if ups:
        lines.append("\n📈 漲幅")
        lines += [_fmt(m) for m in ups]
    lines.append("\n※ 依總行報價檔 Offer 計算，跌多可能是買點也可能是信用事件，請搭配 /rating 與 /issuer 查看")
    return "\n".join(lines), base, latest

def check_bond_alerts(bot_api=None, source="upload"):
    """
    比對所有 active 的到價條件；命中就推播給設定者並關閉該條件（一次性）。
    在『上傳新報價檔』與『每日 06:45』各跑一次。回傳命中數。
    """
    bot_api = bot_api or line_bot_api
    if not (_BOND_RADAR_OK and BOND_PRICE_FILE.exists()):
        return 0
    from bond_coupon_alert import bond_snapshot, first_num
    snap = bond_snapshot(str(BOND_PRICE_FILE))
    with engine.begin() as conn:
        rows = conn.execute(text("SELECT id, chat_id, isin, bond_name, field, op, threshold FROM bond_price_alert WHERE active = TRUE")).fetchall()
    hit = 0
    ops = {">": lambda a, b: a > b, ">=": lambda a, b: a >= b, "<": lambda a, b: a < b, "<=": lambda a, b: a <= b}
    for aid, chat_id, isin, bname, field, op, thr in rows:
        b = snap.get(isin)
        if not b:
            continue
        cur = b.get("ytm") if field == "ytm" else b.get("offer")
        if cur is None:
            continue
        with engine.begin() as conn:
            conn.execute(text("UPDATE bond_price_alert SET last_value=:v WHERE id=:i"), {"v": cur, "i": aid})
        if ops.get(op, lambda a, b: False)(cur, thr):
            hit += 1
            fname = "YTM" if field == "ytm" else "Offer"
            msg = (f"🎯 到價通知 #{aid}\n{bname}\n{fname} {cur} 已 {op} {thr}（設定條件）\n"
                   f"目前 Offer {b.get('offer')}｜YTM {b.get('ytm')}\n"
                   f"（本條件已自動關閉，要繼續追請重新 /bondalert 設定；來源：{'新報價檔' if source=='upload' else '每日檢查'}）")
            try:
                bot_api.push_message(chat_id, TextSendMessage(text=msg))
            except Exception as e:
                print(f"[BondAlert] push fail {e}")
            with engine.begin() as conn:
                conn.execute(text("UPDATE bond_price_alert SET active=FALSE, triggered_at=NOW() WHERE id=:i"), {"i": aid})
    return hit

def snapshot_and_diff():
    """上傳新報價檔後：跟上一份快照比對（信評異動/新上架/下架），再覆蓋快照。回傳文字（無異動則空字串）"""
    from bond_coupon_alert import bond_snapshot, diff_snapshots, format_snapshot_diff
    new = bond_snapshot(str(BOND_PRICE_FILE))
    old = {}
    if BOND_SNAPSHOT_FILE.exists():
        try:
            old = json.loads(BOND_SNAPSHOT_FILE.read_text(encoding="utf-8"))
        except Exception:
            old = {}
    txt_ = format_snapshot_diff(diff_snapshots(old, new)) if old else "（首次建立快照，下次上傳起會比對信評異動／新上架／下架）"
    try:
        BOND_SNAPSHOT_FILE.write_text(json.dumps(new, ensure_ascii=False), encoding="utf-8")
    except Exception as e:
        print(f"[BondSnapshot] write fail {e}")
    return txt_

# --- Alert ticker aliases ---
ALERT_TICKER_ALIAS = {
    "dxy": "DX-Y.NYB", "spx": "^GSPC", "sp500": "^GSPC", "ndx": "^NDX",
    "nasdaq100": "^NDX", "sox": "^SOX", "vix": "^VIX", "ust10y": "^TNX",
    "gold": "GC=F", "silver": "SI=F", "oil": "CL=F", "wti": "CL=F", "copper": "HG=F",
    "usdjpy": "JPY=X", "jpy": "JPY=X", "eurusd": "EURUSD=X", "eur": "EURUSD=X",
    "gbpusd": "GBPUSD=X", "gbp": "GBPUSD=X", "usdtwd": "TWD=X", "twd": "TWD=X",
    "usdcnh": "CNH=X", "cnh": "CNH=X", "usdkrw": "KRW=X", "krw": "KRW=X",
}

# ==============================
# ENV
# ==============================
LINE_CHANNEL_SECRET = os.getenv("LINE_CHANNEL_SECRET")
LINE_CHANNEL_ACCESS_TOKEN = os.getenv("LINE_CHANNEL_ACCESS_TOKEN")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
DATABASE_URL = os.getenv("DATABASE_URL")

if not LINE_CHANNEL_SECRET or not LINE_CHANNEL_ACCESS_TOKEN:
    raise RuntimeError("Missing LINE env vars: LINE_CHANNEL_SECRET / LINE_CHANNEL_ACCESS_TOKEN")
if not DATABASE_URL:
    raise RuntimeError("Missing env var: DATABASE_URL")
if not ANTHROPIC_API_KEY:
    raise RuntimeError("Missing env var: ANTHROPIC_API_KEY")

if DATABASE_URL.startswith("postgres://"):
    DATABASE_URL = DATABASE_URL.replace("postgres://", "postgresql+psycopg://", 1)
elif DATABASE_URL.startswith("postgresql://"):
    DATABASE_URL = DATABASE_URL.replace("postgresql://", "postgresql+psycopg://", 1)

line_bot_api = LineBotApi(LINE_CHANNEL_ACCESS_TOKEN)
handler = WebhookHandler(LINE_CHANNEL_SECRET)

ELN_GROUP_CHANNEL_SECRET = os.getenv("AGENT_LINE_CHANNEL_SECRET", "")
ELN_GROUP_ACCESS_TOKEN = os.getenv("AGENT_LINE_CHANNEL_ACCESS_TOKEN", "")
eln_group_bot_api = LineBotApi(ELN_GROUP_ACCESS_TOKEN) if ELN_GROUP_ACCESS_TOKEN else None
eln_group_handler = WebhookHandler(ELN_GROUP_CHANNEL_SECRET) if ELN_GROUP_CHANNEL_SECRET else None

claude_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
openai_client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

app = FastAPI()
from articles import router as articles_router
app.include_router(articles_router)
from eln_form_router import router as eln_form_router
app.include_router(eln_form_router)

VERSION = "eln-autotracking-db-v3-2026-03-05"
TZ_TAIPEI = timezone(timedelta(hours=8))

# ==============================
# DB
# ==============================
engine = create_engine(DATABASE_URL, pool_pre_ping=True)

def init_db():
    with engine.begin() as conn:
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_price_history (
            snap_date DATE NOT NULL, isin TEXT NOT NULL, bond_name TEXT, ccy TEXT,
            offer DOUBLE PRECISION, ytm DOUBLE PRECISION, maturity DATE,
            PRIMARY KEY (snap_date, isin)
        );"""))
        conn.execute(text("""
        ALTER TABLE bond_price_history ADD COLUMN IF NOT EXISTS bid DOUBLE PRECISION;"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_issuer_profile (
            issuer TEXT PRIMARY KEY, profile TEXT NOT NULL,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_issuer_ticker (
            issuer TEXT PRIMARY KEY, parent TEXT, ticker TEXT,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_rating_watch (
            issuer TEXT PRIMARY KEY, en_name TEXT, added_by TEXT, active BOOLEAN NOT NULL DEFAULT TRUE,
            added_at TIMESTAMPTZ NOT NULL DEFAULT NOW(), last_checked TIMESTAMPTZ
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_rating_news_seen (
            link TEXT PRIMARY KEY, issuer TEXT, title TEXT, seen_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_price_alert (
            id SERIAL PRIMARY KEY, chat_id TEXT NOT NULL, isin TEXT NOT NULL, bond_name TEXT,
            field TEXT NOT NULL, op TEXT NOT NULL, threshold DOUBLE PRECISION NOT NULL,
            active BOOLEAN NOT NULL DEFAULT TRUE, created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            triggered_at TIMESTAMPTZ, last_value DOUBLE PRECISION
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS bond_issuer_intro (
            issuer TEXT PRIMARY KEY, intro TEXT NOT NULL, source TEXT,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS eln_last_report (
            chat_key TEXT PRIMARY KEY, summary TEXT NOT NULL,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS eln_detail (
            chat_key TEXT NOT NULL, bond_id TEXT NOT NULL, detail TEXT NOT NULL,
            agent_name TEXT, updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            PRIMARY KEY (chat_key, bond_id)
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS eln_top5 (
            chat_key TEXT NOT NULL, line_no INT NOT NULL, text_line TEXT NOT NULL,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW(), PRIMARY KEY (chat_key, line_no)
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS eln_session (
            chat_key TEXT PRIMARY KEY, await_file BOOLEAN NOT NULL DEFAULT FALSE,
            invest_mode TEXT NOT NULL DEFAULT '', invest_image BYTEA,
            updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS transcript_cache (
            chat_key TEXT PRIMARY KEY, transcript TEXT NOT NULL DEFAULT '',
            summary TEXT NOT NULL DEFAULT '', updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS meeting_transcripts (
            id BIGSERIAL PRIMARY KEY, chat_key TEXT NOT NULL, file_name TEXT,
            transcript TEXT NOT NULL DEFAULT '', summary TEXT NOT NULL DEFAULT '',
            created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS articles (
            id BIGSERIAL PRIMARY KEY, title TEXT, content TEXT, summary TEXT,
            source_type TEXT DEFAULT 'text', image_url TEXT, is_read BOOLEAN DEFAULT FALSE,
            category TEXT DEFAULT 'other', location_name TEXT, lat FLOAT, lng FLOAT,
            show_on_map BOOLEAN DEFAULT TRUE, created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );"""))
        for col, typedef in [
            ("invest_mode", "TEXT NOT NULL DEFAULT ''"),
            ("invest_image", "BYTEA"),
        ]:
            try:
                conn.execute(text(f"ALTER TABLE eln_session ADD COLUMN IF NOT EXISTS {col} {typedef}"))
            except Exception:
                pass

init_db()

def db_set_await(chat_key: str, await_file: bool):
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO eln_session(chat_key, await_file, updated_at) VALUES (:k, :a, NOW())
        ON CONFLICT (chat_key) DO UPDATE SET await_file=:a, updated_at=NOW()
        """), {"k": chat_key, "a": bool(await_file)})

def db_is_await(chat_key: str) -> bool:
    with engine.begin() as conn:
        row = conn.execute(text("SELECT await_file FROM eln_session WHERE chat_key=:k"), {"k": chat_key}).fetchone()
    return bool(row and row[0])

def db_invest_set(chat_key: str, mode: str, image: bytes = None):
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO eln_session(chat_key, await_file, invest_mode, invest_image, updated_at)
        VALUES (:k, FALSE, :m, :img, NOW())
        ON CONFLICT (chat_key) DO UPDATE
        SET invest_mode=:m, invest_image=COALESCE(:img, eln_session.invest_image), updated_at=NOW()
        """), {"k": chat_key, "m": mode, "img": image})

def db_invest_get(chat_key: str):
    with engine.begin() as conn:
        row = conn.execute(text("SELECT invest_mode, invest_image FROM eln_session WHERE chat_key=:k"), {"k": chat_key}).fetchone()
    if row:
        return row[0] or "", bytes(row[1]) if row[1] else None
    return "", None

def db_set_transcript_cache(chat_key: str, transcript: str, summary: str):
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO transcript_cache(chat_key, transcript, summary, updated_at)
        VALUES (:k, :t, :s, NOW())
        ON CONFLICT (chat_key) DO UPDATE SET transcript=:t, summary=:s, updated_at=NOW()
        """), {"k": chat_key, "t": transcript[:200000], "s": summary[:50000]})

def db_get_transcript_cache(chat_key: str):
    with engine.begin() as conn:
        row = conn.execute(text("SELECT transcript, summary FROM transcript_cache WHERE chat_key=:k"), {"k": chat_key}).fetchone()
    if row:
        return {"transcript": row[0] or "", "summary": row[1] or ""}
    return None

def db_clear_transcript_cache(chat_key: str):
    with engine.begin() as conn:
        conn.execute(text("DELETE FROM transcript_cache WHERE chat_key=:k"), {"k": chat_key})

def db_save_meeting_transcript(chat_key: str, file_name: str, transcript: str, summary: str):
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO meeting_transcripts(chat_key, file_name, transcript, summary, created_at)
        VALUES (:k, :f, :t, :s, NOW())
        """), {"k": chat_key, "f": file_name, "t": transcript[:500000], "s": summary[:100000]})

def db_get_latest_meeting_transcript(chat_key: str):
    with engine.begin() as conn:
        row = conn.execute(text("""
        SELECT transcript, summary, file_name, created_at FROM meeting_transcripts
        WHERE chat_key=:k ORDER BY created_at DESC LIMIT 1
        """), {"k": chat_key}).fetchone()
    if not row:
        return None
    return {"transcript": row[0] or "", "summary": row[1] or "", "file_name": row[2] or "", "created_at": row[3]}

def db_save_result(chat_key: str, summary: str, top5_lines: list[str], detail_map: dict[str, str], agent_name_map: dict[str, str] = {}):
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO eln_last_report(chat_key, summary, updated_at) VALUES (:k, :s, NOW())
        ON CONFLICT (chat_key) DO UPDATE SET summary=:s, updated_at=NOW()
        """), {"k": chat_key, "s": summary})
        conn.execute(text("DELETE FROM eln_top5 WHERE chat_key=:k"), {"k": chat_key})
        for i, line in enumerate(top5_lines, start=1):
            conn.execute(text("""
            INSERT INTO eln_top5(chat_key, line_no, text_line, updated_at) VALUES (:k, :n, :t, NOW())
            """), {"k": chat_key, "n": i, "t": line})
        conn.execute(text("DELETE FROM eln_detail WHERE chat_key=:k"), {"k": chat_key})
        for bond_id, detail in detail_map.items():
            agent = agent_name_map.get(bond_id, "-")
            conn.execute(text("""
            INSERT INTO eln_detail(chat_key, bond_id, detail, agent_name, updated_at)
            VALUES (:k, :b, :d, :a, NOW())
            """), {"k": chat_key, "b": bond_id, "d": detail, "a": agent})

def db_get_report(chat_key: str) -> str | None:
    with engine.begin() as conn:
        row = conn.execute(text("SELECT summary FROM eln_last_report WHERE chat_key=:k"), {"k": chat_key}).fetchone()
    return row[0] if row else None

def db_list_bonds(chat_key: str, limit: int = 100) -> list[tuple[str, str, str]]:
    with engine.begin() as conn:
        rows = conn.execute(text("""
        SELECT bond_id, COALESCE(agent_name, '-'), COALESCE(detail, '')
        FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC LIMIT :lim
        """), {"k": chat_key, "lim": int(limit)}).fetchall()
    return [(r[0], r[1], r[2]) for r in rows] if rows else []

def bond_status_tag(detail: str) -> str:
    import re as _re
    status_block = ""
    m = _re.search(r"-{4,}\n(.*?)\n-{4,}", detail, _re.S)
    if m:
        status_block = m.group(1).strip()
    if "提前出場" in status_block or "🎉" in status_block:
        return " ✅提前KO"
    if "到期獲利" in status_block:
        return " 🏁到期獲利"
    if "到期接股" in status_block:
        return " 😭到期接股"
    if "到期保本" in status_block:
        return " 🛡️到期保本"
    if "到期" in status_block:
        return " 🏁到期"
    return ""

def push_long_message(bot_api, target_id: str, text: str, max_len: int = 4800):
    if not text:
        return
    text = str(text)
    chunks = []
    current = ""
    for line in text.split("\n"):
        while len(line) > max_len:
            if current:
                chunks.append(current)
                current = ""
            chunks.append(line[:max_len])
            line = line[max_len:]
        candidate = line if not current else current + "\n" + line
        if len(candidate) <= max_len:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = line
    if current:
        chunks.append(current)
    safe_chunks = []
    for chunk in chunks:
        while len(chunk) > max_len:
            safe_chunks.append(chunk[:max_len])
            chunk = chunk[max_len:]
        if chunk:
            safe_chunks.append(chunk)
    for chunk in safe_chunks:
        bot_api.push_message(target_id, TextSendMessage(text=chunk))

def db_find_detail(chat_key: str, query: str) -> tuple[str | None, str | None, list[str]]:
    q_norm = query.strip().upper()
    if not q_norm:
        return None, None, []
    with engine.begin() as conn:
        rows = conn.execute(text("SELECT bond_id FROM eln_detail WHERE chat_key=:k"), {"k": chat_key}).fetchall()
    keys = [r[0] for r in rows] if rows else []
    if not keys:
        return None, None, []
    norm_map = {k.strip().upper(): k for k in keys}
    if q_norm in norm_map:
        real = norm_map[q_norm]
        with engine.begin() as conn:
            row = conn.execute(text("SELECT detail FROM eln_detail WHERE chat_key=:k AND bond_id=:b"), {"k": chat_key, "b": real}).fetchone()
        return real, (row[0] if row else None), []
    hits = [k for k in keys if q_norm in k.strip().upper()]
    if len(hits) == 1:
        real = hits[0]
        with engine.begin() as conn:
            row = conn.execute(text("SELECT detail FROM eln_detail WHERE chat_key=:k AND bond_id=:b"), {"k": chat_key, "b": real}).fetchone()
        return real, (row[0] if row else None), []
    if len(hits) > 1:
        return None, None, hits[:20]
    return None, None, keys[:20]

# ==============================
# Optional: store default push target
# ==============================
BASE_DIR = Path("/tmp")
def _persistent_dir():
    """優先用 Render 持久磁碟 /data（重新部署不會消失），沒有才退回 /tmp"""
    for d in (Path("/data"), Path("/tmp")):
        try:
            d.mkdir(parents=True, exist_ok=True)
            probe = d / ".write_test"
            probe.write_text("ok")
            probe.unlink()
            return d
        except Exception:
            continue
    return Path("/tmp")

TARGET_FILE = _persistent_dir() / "targets.json"

# 一次性搬遷：舊版存在 /tmp，若持久位置還沒有檔案就把舊的搬過來
_old_target = Path("/tmp") / "targets.json"
try:
    if TARGET_FILE != _old_target and _old_target.exists() and not TARGET_FILE.exists():
        import shutil as _sh
        _sh.copy(str(_old_target), str(TARGET_FILE))
except Exception as _e:
    print(f"[targets migrate] {_e}")

def _read_json(path: Path, default):
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return default
    return default

def _write_json(path: Path, data):
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

def load_targets():
    return _read_json(TARGET_FILE, {})

def save_targets(data: dict):
    _write_json(TARGET_FILE, data)

# ==============================
# Health check
# ==============================
@app.get("/")
def root():
    return {"status": "ok", "service": "eln-bot", "webhook": "/callback"}

@app.get("/whoami")
def whoami():
    return {"service": "eln-bot", "version": VERSION}

# ==============================
# Webhook endpoint
# ==============================
# ==============================
# 報價檔網頁上傳（供無法用 LINE 傳檔時使用）
# ==============================
BOND_UPLOAD_TOKEN = os.getenv("BOND_UPLOAD_TOKEN", "")

_BOND_UPLOAD_HTML = """<!doctype html>
<html lang="zh-Hant"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>債券報價檔上傳</title>
<style>
 body{font-family:"Microsoft JhengHei",system-ui,sans-serif;background:#f4f7fa;margin:0;padding:24px}
 .card{max-width:520px;margin:0 auto;background:#fff;border-radius:12px;padding:28px;
       box-shadow:0 2px 12px rgba(0,0,0,.07)}
 h1{font-size:20px;color:#0B2A4A;margin:0 0 6px}
 p.sub{color:#667;font-size:13px;margin:0 0 20px}
 label{display:block;font-size:13px;color:#334;margin:14px 0 6px;font-weight:600}
 input[type=password],input[type=file]{width:100%;padding:10px;border:1px solid #ccd;
       border-radius:8px;font-size:14px;box-sizing:border-box;background:#fafbfc}
 button{margin-top:20px;width:100%;padding:12px;background:#1F8AC0;color:#fff;border:0;
       border-radius:8px;font-size:15px;font-weight:600;cursor:pointer}
 button:disabled{background:#9bc;cursor:wait}
 .msg{margin-top:16px;padding:12px;border-radius:8px;font-size:14px;line-height:1.6}
 .ok{background:#e8f6ee;color:#1b6b3a}.err{background:#fdecec;color:#a12}
 .note{margin-top:18px;font-size:12px;color:#889;line-height:1.7}
</style></head><body>
<div class="card">
  <h1>債券報價檔上傳</h1>
  <p class="sub">上傳後將自動更新報價、產生配息雷達，並比對信評異動與報價變動</p>
  <form method="post" enctype="multipart/form-data" onsubmit="this.querySelector('button').disabled=true;this.querySelector('button').textContent='上傳處理中，請稍候…'">
    <label>存取密碼</label>
    <input type="password" name="token" required autocomplete="off">
    <label>報價檔（.xlsx）</label>
    <input type="file" name="file" accept=".xlsx,.xlsm" required>
    <button type="submit">上傳並更新</button>
  </form>
  __MSG__
  <div class="note">
    ・檔名請保留日期（例如 09-02-2026-Bond Pricing Update.xlsx）<br>
    ・檔名日期為今日或近 3 天：更新最新報價並跑完整流程<br>
    ・檔名日期較舊：僅補歷史報價，不會覆蓋最新報價<br>
    ・處理結果會同步推播到 LINE
  </div>
</div></body></html>"""


@app.get("/bond-upload", response_class=HTMLResponse)
def bond_upload_page():
    return _BOND_UPLOAD_HTML.replace("__MSG__", "")


@app.post("/bond-upload", response_class=HTMLResponse)
async def bond_upload_post(token: str = Form(""), file: UploadFile = File(...)):
    def _page(msg, ok=True):
        cls = "ok" if ok else "err"
        return _BOND_UPLOAD_HTML.replace("__MSG__", f'<div class="msg {cls}">{msg}</div>')

    if not BOND_UPLOAD_TOKEN or token != BOND_UPLOAD_TOKEN:
        return _page("❌ 密碼錯誤", ok=False)
    if not _BOND_RADAR_OK:
        return _page("❌ 報價模組未載入，請聯絡管理者", ok=False)
    fname = file.filename or "upload.xlsx"
    if not fname.lower().endswith((".xlsx", ".xlsm")):
        return _page("❌ 請上傳 .xlsx 檔案", ok=False)

    tmp_path = f"/tmp/upload_{int(datetime.now().timestamp())}_{os.path.basename(fname)}"
    try:
        with open(tmp_path, "wb") as f:
            f.write(await file.read())
        if not _is_bond_pricing_file(tmp_path, fname):
            os.remove(tmp_path)
            return _page("❌ 這個檔案看起來不是海外債報價檔（找不到報價檔的工作表）", ok=False)

        today_ = datetime.now(TZ_TAIPEI).date()
        snap = parse_pricing_file_date(fname)
        user_id = os.getenv("LINE_USER_ID", "")

        # 舊檔 → 只補歷史
        if snap and (today_ - snap).days >= 3:
            n = save_price_history(snap_date=snap, path=tmp_path)
            os.remove(tmp_path)
            write_job_log("報價檔上傳(網頁)", "history", f"{snap} {n}筆")
            return _page(f"✅ 已補入 {snap:%Y/%m/%d} 歷史報價快照（{n} 檔）<br>"
                         "最新報價檔與配息雷達未變更。")

        import shutil
        BOND_PRICE_FILE.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy(tmp_path, str(BOND_PRICE_FILE))
        os.remove(tmp_path)
        write_job_log("報價檔上傳(網頁)", "file_updated", fname)

        def _post_process():
            try:
                msg = _bond_build_alert(str(BOND_PRICE_FILE), today=today_)
                if user_id:
                    push_long_message(line_bot_api, user_id, "✅ 報價檔已更新（網頁上傳）\n\n" + msg)
                diff_txt = snapshot_and_diff()
                if diff_txt and user_id:
                    push_long_message(line_bot_api, user_id, "📋 與上一份報價檔比對\n" + diff_txt)
                save_price_history(snap_date=(snap if snap and snap <= today_ else None))
                for db_, th_ in ((1, 2.0), (7, 3.0)):
                    mv, _, _ = price_movers(days_back=db_, threshold_pct=th_)
                    if mv and user_id:
                        push_long_message(line_bot_api, user_id, mv)
                check_bond_alerts(line_bot_api, source="upload")
            except Exception as e:
                print(f"[BondUpload post] {e}")
                print(_traceback.format_exc()[:500])

        import threading
        threading.Thread(target=_post_process, daemon=True).start()
        return _page(f"✅ 報價檔已更新（{fname}）<br>"
                     "配息雷達、信評比對與報價異動正在背景處理，結果將推播到 LINE。")
    except Exception as e:
        print(f"[BondUpload] {e}")
        print(_traceback.format_exc()[:500])
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        return _page(f"❌ 處理失敗：{str(e)[:200]}", ok=False)


@app.post("/callback")
async def callback(request: Request):
    signature = request.headers.get("X-Line-Signature")
    body = await request.body()
    body_text = body.decode("utf-8")
    try:
        handler.handle(body_text, signature)
        return "OK"
    except InvalidSignatureError:
        raise HTTPException(status_code=400, detail="Invalid signature")

import threading
_current_bot_api = threading.local()
_current_bot_api.api = None
_current_bot_api.is_eln = False

@app.post("/callback2")
async def callback2(request: Request):
    body = await request.body()
    try:
        import json as _j
        data = _j.loads(body.decode("utf-8"))
        for ev in data.get("events", []):
            if ev.get("type") != "message":
                continue
            if ev.get("message", {}).get("type") != "text":
                continue
            txt = ev["message"]["text"].strip()
            tl = txt.lower()
            rtoken = ev.get("replyToken", "")
            uid = ev.get("source", {}).get("userId", "")
            print(f"[ELN-G USER] uid={uid} msg={repr(txt[:50])}")
            _eln_bond_prefixes = tuple("/" + c for c in
                                       (ELN_BOT_BOND_CMDS | ELN_BOT_BOND_HEAVY | {"help", "myid"}))
            if tl.startswith(_eln_bond_prefixes):
                # 海外債指令:轉交主處理器(共用同一套邏輯),回覆走 ELN Bot 頻道
                try:
                    from types import SimpleNamespace as _NS
                    src = ev.get("source", {})
                    _stype = src.get("type", "user")
                    _ev = _NS(
                        reply_token=rtoken,
                        message=_NS(text=txt, id=ev.get("message", {}).get("id", "")),
                        source=_NS(type=_stype,
                                   user_id=src.get("userId", ""),
                                   group_id=src.get("groupId", ""),
                                   room_id=src.get("roomId", "")),
                    )
                    _current_bot_api.api = eln_group_bot_api
                    _current_bot_api.is_eln = True
                    handle_text_message(_ev)
                except Exception as e:
                    print(f"[ELN-G BOND ERR] {e}")
                    print(_traceback.format_exc()[:400])
                finally:
                    _current_bot_api.api = None
                    _current_bot_api.is_eln = False
                continue
            if not (tl.startswith("/list") or tl.startswith("/detail") or tl.startswith("/end") or tl.startswith("/nc") or tl.startswith("/內規")):
                continue
            from linebot.models import TextSendMessage as TSM
            from collections import defaultdict
            ck = ELN_PERSONAL_CHAT_KEY
            
            # ==========================================
            # 群組版 內規專屬指令攔截 (callback2)
            # ==========================================
            if tl.startswith("/內規"):
                actual_query = txt.replace("/內規", "").strip()
                if not actual_query:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="請在指令後面加上想查詢的內容喔！\n例如：/內規 Lombard lending 最高可以到幾歲？"))
                    continue
                try:
                    file_path = Path("regulations.txt")
                    if not file_path.exists():
                        eln_group_bot_api.reply_message(rtoken, TSM(text="❌ 找不到 regulations.txt，請確認已將法規檔案上傳至系統。"))
                        continue
                    regulation_text = file_path.read_text(encoding="utf-8")
                    
                    # 💯 防縮排錯誤寫法
                    prompt_lines = [
                        "你現在是銀行的法遵與內部規範專家。請根據以下【內部規範全文】，直接且精準回答同仁的問題。",
                        "",
                        "【嚴格限制】",
                        "1. 絕對不要輸出任何「因為文本是程式碼...」或「無法回答」的廢話警告。",
                        "2. 絕對不要在結尾補充「可以這樣跟客戶/專員說」的話術。",
                        "3. 嚴禁使用 Markdown 語法 (例如 **, ##, --- 等)，請用純文字或 Emoji 條列排版。",
                        "",
                        "【內部規範全文】",
                        regulation_text,
                        "",
                        "【同仁問題】",
                        actual_query
                    ]
                    prompt = "\n".join(prompt_lines)
                    
                    raw_answer = ai_claude_long(prompt, chat_key=ck)
                    final_answer = "⚠️ 僅供參考，本回覆由 AI 統整，不代表總行最終內規解釋。\n\n" + raw_answer
                    eln_group_bot_api.reply_message(rtoken, TSM(text=final_answer[:4900]))
                except Exception as e:
                    eln_group_bot_api.reply_message(rtoken, TSM(text=f"❌ 內部規範查詢失敗：{e}"))
                continue
            # ==========================================
            
            if tl.startswith("/list"):
                lp = txt.split(" ", 2)
                is_detail_mode = len(lp) > 1 and lp[1].strip().lower() == "detail"
                if is_detail_mode:
                    nf = lp[2].strip() if len(lp) > 2 else ""
                else:
                    nf = lp[1].strip() if len(lp) > 1 else ""
                bonds = db_list_bonds(ck, limit=200)
                if not bonds:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="目前尚無資料。"))
                    continue
                # /list detail 姓名：顯示完整 detail，只顯示比價中
                if is_detail_mode:
                    if not nf:
                        eln_group_bot_api.reply_message(rtoken, TSM(text="請輸入理專名稱\n例：/list detail 小美"))
                        continue
                    matched_details = []
                    seen = set()
                    for bid, ar, d in bonds:
                        ags = [a.strip() for a in re.split(r"[,，、/]", ar) if a.strip()]
                        if any(nf in a for a in ags) and bid not in seen:
                            if bond_status_tag(d) == "":
                                matched_details.append(d)
                                seen.add(bid)
                    if not matched_details:
                        eln_group_bot_api.reply_message(rtoken, TSM(text="找不到「" + nf + "」比價中的持倉。"))
                        continue
                    eln_group_bot_api.reply_message(rtoken, TSM(text="👤 " + nf + " 比價中商品（共 " + str(len(matched_details)) + " 筆）"))
                    for det in matched_details:
                        eln_group_bot_api.push_message(uid, TSM(text=det[:4900]))
                    continue
                ds = {b: bond_status_tag(d) for b, _, d in bonds}
                if nf:
                    matched = []
                    seen = set()
                    for bid, ar, d in bonds:
                        ags = [a.strip() for a in re.split(r"[,，、/]", ar) if a.strip()]
                        if any(nf in a for a in ags) and bid not in seen:
                            matched.append((bid, ds.get(bid, "")))
                            seen.add(bid)
                    if not matched:
                        eln_group_bot_api.reply_message(rtoken, TSM(text="找不到「" + nf + "」的持倉。"))
                        continue
                    out = "👤 " + nf + " 的持倉（共 " + str(len(matched)) + " 筆）：\n"
                    for b, t in matched:
                        out += "   • " + b + t + "\n"
                else:
                    grp = defaultdict(list)
                    for bid, ar, d in bonds:
                        ags = [a.strip() for a in re.split(r"[,，、/]", ar) if a.strip()] or ["未指定"]
                        for ag in ags:
                            if bid not in [x for x, _ in grp[ag]]:
                                grp[ag].append((bid, ds.get(bid, "")))
                    out = "📋 全部商品（共 " + str(len(set(b for b,_,_ in bonds))) + " 筆）：\n"
                    for ag, bl in sorted(grp.items()):
                        out += "👤 " + ag + "（" + str(len(bl)) + " 筆）\n"
                        for b, t in bl:
                            out += "   • " + b + t + "\n"
                chunks = [out[i:i+4800] for i in range(0, len(out), 4800)]
                eln_group_bot_api.reply_message(rtoken, TSM(text=chunks[0]))
                for c in chunks[1:]:
                    eln_group_bot_api.push_message(uid, TSM(text=c))
            elif tl.startswith("/detail"):
                ps = txt.split(" ", 1)
                if len(ps) < 2 or not ps[1].strip():
                    eln_group_bot_api.reply_message(rtoken, TSM(text="請輸入：/detail 商品代號"))
                    continue
                mid, det, cands = db_find_detail(ck, ps[1].strip())
                if det:
                    eln_group_bot_api.reply_message(rtoken, TSM(text=det[:4900]))
                elif cands:
                    eln_group_bot_api.reply_message(rtoken, TSM(text=("候選代號：\n" + "\n".join("• "+c for c in cands[:20]))[:4900]))
                else:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="查不到該代號。"))
            elif tl.startswith("/end"):
                ps = txt.split(" ", 1)
                if len(ps) < 2 or not ps[1].strip():
                    eln_group_bot_api.reply_message(rtoken, TSM(text="請輸入：/end YYYYMM\n例：/end 202604"))
                    continue
                qm = ps[1].strip().replace("/", "").replace("-", "")
                if len(qm) != 6 or not qm.isdigit():
                    eln_group_bot_api.reply_message(rtoken, TSM(text="格式錯誤，請輸入6位數字\n例：/end 202604"))
                    continue
                yr, mo = qm[:4], qm[4:]
                search_str = yr + "-" + mo
                with engine.begin() as conn:
                    rows = conn.execute(text("SELECT bond_id, agent_name, detail FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC"), {"k": ck}).fetchall()
                if not rows:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="目前尚無資料。"))
                    continue
                matched = []
                for bid, ag, det in rows:
                    if ("最終評價日: " + search_str) in det:
                        matched.append((bid, ag or "-", bond_status_tag(det)))
                if not matched:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="找不到 " + yr + "/" + mo + " 到期的商品。"))
                    continue
                out = "📅 " + yr + "/" + mo + " 到期商品（共 " + str(len(matched)) + " 筆）：\n"
                for bid, ag, tag in matched:
                    out += "   • " + bid + " [" + ag + "]" + tag + "\n"
                eln_group_bot_api.reply_message(rtoken, TSM(text=out[:4900]))
            elif tl.startswith("/nc"):
                ps = txt.split(" ")
                if len(ps) < 2 or not ps[1].strip():
                    eln_group_bot_api.reply_message(rtoken, TSM(text="請輸入：/nc YYYYMM\n例：/nc 202606\n或：/nc 202606 小美"))
                    continue
                qm = ps[1].strip().replace("/", "").replace("-", "")
                if len(qm) != 6 or not qm.isdigit():
                    eln_group_bot_api.reply_message(rtoken, TSM(text="格式錯誤，請輸入6位數字\n例：/nc 202606"))
                    continue
                yr, mo = qm[:4], qm[4:]
                name_filter_nc = ps[2].strip() if len(ps) > 2 else ""
                search_str_nc = yr + "-" + mo
                with engine.begin() as conn:
                    rows = conn.execute(text("SELECT bond_id, agent_name, detail FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC"), {"k": ck}).fetchall()
                if not rows:
                    eln_group_bot_api.reply_message(rtoken, TSM(text="目前尚無資料。"))
                    continue
                matched = []
                import re as _re2
                for bid, ag, det in rows:
                    m = _re2.search(r"NC閉鎖期 \(至 (\d{4}-\d{2})-\d{2}\)", det)
                    if m and m.group(1) == search_str_nc:
                        if name_filter_nc:
                            ags = [a.strip() for a in _re2.split(r"[,，、/]", ag or "") if a.strip()]
                            if not any(name_filter_nc in a for a in ags):
                                continue
                        matched.append((bid, ag or "-", bond_status_tag(det)))
                if not matched:
                    tip = f"「{name_filter_nc}」" if name_filter_nc else ""
                    eln_group_bot_api.reply_message(rtoken, TSM(text=f"找不到 {yr}/{mo} 閉鎖期打開{tip}的商品。"))
                    continue
                tip2 = f"（{name_filter_nc}）" if name_filter_nc else ""
                out_nc = f"🔓 {yr}/{mo} 閉鎖期打開{tip2}（共 {len(matched)} 筆）：\n"
                for bid, ag, tag in matched:
                    out_nc += f"   • {bid} [{ag}]{tag}\n"
                eln_group_bot_api.reply_message(rtoken, TSM(text=out_nc[:4900]))
    except Exception as e:
        print("[callback2 ERR]", e)
    return "OK"

# ==============================
# Chat key
# ==============================
def chat_key_of(event) -> str:
    if event.source.type == "group":
        return f"group:{event.source.group_id}"
    if event.source.type == "room":
        return f"room:{event.source.room_id}"
    return f"user:{event.source.user_id}"

# ==============================
# Adapter: core -> (summary, top5, detail_map)
# ==============================
def run_autotracking(file_path: str, lookback_days: int = 3, notify_ki_daily: bool = True):
    from auto_tracking_cron import build_result
    out = calculate_from_file(file_path=file_path, lookback_days=lookback_days, notify_ki_daily=notify_ki_daily)
    return build_result(out)

# ==============================
# AI
# ==============================
SYSTEM_PROMPT = (
    "你是「龍蝦」，一位專屬於 Albert 的智慧型 LINE 助理。\n"
    "Albert 有時會叫你「阿斯拉」，這時候你要回應他「風見」，這是你們之間的暗語，記住後繼續正常對話。\n\n"
    "【關於 Albert 的工作背景】\n"
    "• 職位：銀行財富管理部門 投資輔銷人員（Investment Sales）\n"
    "• 客群：主要負責高資產客戶（HNW）的投資規劃與資產配置\n"
    "• 業務範疇：\n"
    "  → 投資商品：基金、債券、ELN（股票連結票據）、結構型商品、ETF\n"
    "  → 質借業務：Lombard Lending（有價證券質借）、金市債券質借、信託質借\n"
    "  → 信託業務：資產信託規劃、境外資金匯回配置\n"
    "  → 教育訓練：經常幫行內專員上課，教導基金、債券、結構型產品、ELN等商品知識\n"
    "• 常見需求：市場分析、商品說明、客戶推播文案、專員教育訓練教材、投資建議\n\n"
    "【你的角色定位】\n"
    "你是 Albert 最得力的資深助理，不只回答問題，而是像一位懂市場又懂銷售的同事：\n"
    "• 用投資輔銷的角度思考，理解他面對的是高資產客戶與行內專員\n"
    "• 遇到市場問題 → 提供深度分析，並附上「可以這樣跟客戶說」的話術\n"
    "• 遇到商品問題 → 說明商品特性、適合的高資產客群、風險與機會\n"
    "• 遇到質借/信託問題 → 說明業務邏輯、適用情境、常見客戶疑問\n"
    "• 遇到教學需求 → 以簡單易懂的方式說明，適合用來對專員解說\n"
    "• 遇到文案需求 → 直接產出可複製貼上的推播內容\n"
    "• 遇到 ELN 相關問題 → 提示使用 /calc 或 /detail 指令\n\n"
    "【回答原則】\n"
    "1. 有深度：提供背景、現況、影響、展望，不能太簡短\n"
    "2. 結構清晰：重點分段，讓人一眼看懂\n"
    "3. 客觀中立：呈現多空兩面，讓 Albert 自行判斷\n"
    "4. 實用導向：一般問題結尾補充「💬 可以這樣跟客戶/專員說：...」\n"
    "5. 市場問題格式：📌 定義 → 📊 現況 → ⚖️ 機會與風險 → 🔭 展望 → 💬 話術\n"
    "6. 商品教學格式：📌 商品定義 → 🔧 運作方式 → 👤 適合客群 → ⚠️ 風險提示 → 💬 話術\n"
    "7. 質借業務格式：📌 業務說明 → 💡 適用情境 → 📊 利率/條件 → ❓ 常見客戶問題\n\n"
    "【格式規定】\n"
    "• 絕對禁止 Markdown：不可出現 ##、**、--- 等符號\n"
    "• 段落標題用 emoji，例如 📌 📊 ⚖️ 🔭 💡 💬 🔧 👤 ⚠️\n"
    "• 條列用 • 或 → 符號\n"
    "• 數字、百分比、金額要具體，不要模糊帶過\n"
    "• 回答長度要足夠，高資產客戶的問題不能給太簡短的答案\n"
)

def get_chat_history(chat_key: str, limit: int = 10) -> list[dict]:
    try:
        with engine.begin() as conn:
            rows = conn.execute(text("""
            SELECT role, content FROM chat_history WHERE chat_key = :k
            ORDER BY created_at DESC LIMIT :n
            """), {"k": chat_key, "n": limit}).fetchall()
        return [{"role": r[0], "content": r[1]} for r in reversed(rows)]
    except Exception as e:
        print(f"get_chat_history error: {e}")
        return []

def _get_memory_collection():
    try:
        import chromadb
        from chromadb.utils import embedding_functions
        chroma_dir = Path("/data/knowledge/chroma_db")
        chroma_dir.mkdir(parents=True, exist_ok=True)
        client = chromadb.PersistentClient(path=str(chroma_dir))
        ef = embedding_functions.DefaultEmbeddingFunction()
        return client.get_or_create_collection(
            name="chat_memory",
            embedding_function=ef,
            metadata={"hnsw:space": "cosine"}
        )
    except Exception as e:
        print(f"[Memory] ChromaDB 初始化失敗：{e}")
        return None

def save_chat_history(chat_key: str, role: str, content: str):
    try:
        with engine.begin() as conn:
            conn.execute(text("INSERT INTO chat_history (chat_key, role, content) VALUES (:k, :r, :c)"),
                         {"k": chat_key, "r": role, "c": content[:4000]})
        with engine.begin() as conn:
            conn.execute(text("""
            DELETE FROM chat_history WHERE chat_key = :k AND id NOT IN (
                SELECT id FROM chat_history WHERE chat_key = :k ORDER BY created_at DESC LIMIT 50
            )"""), {"k": chat_key})
    except Exception as e:
        print(f"save_chat_history error: {e}")
    if role == "assistant" and content.strip():
        try:
            col = _get_memory_collection()
            if col:
                import uuid as _uuid
                now_str = datetime.now(TZ_TAIPEI).strftime("%Y-%m-%d %H:%M")
                mem_id = f"mem_{chat_key}_{_uuid.uuid4().hex[:8]}"
                col.add(
                    documents=[content[:2000]],
                    ids=[mem_id],
                    metadatas=[{"chat_key": chat_key, "role": role, "created_at": now_str}]
                )
        except Exception as e:
            print(f"[Memory] 存入 ChromaDB 失敗：{e}")

SPENDING_NL_KEYWORDS = ["消費明細", "花了多少", "這個月花", "上個月花", "消費分析", "帳單分析", "錢花到哪", "月度消費", "消費統計"]
AUTO_FINANCE_KEYWORDS = ["財經", "市場", "美股", "台股", "債券", "殖利率", "基金", "匯率", "美元", "聯準會", "fed", "fomc", "通膨", "cpi", "pce", "非農", "失業率", "投資", "分析", "總經", "景氣", "eln", "結構型", "信用利差", "公司債"]
AUTO_FILE_KEYWORDS = ["pdf", "簡報", "圖片", "圖表", "文件", "檔案", "word", "excel"]
PDF_NL_KEYWORDS = ["pdf", "做成pdf", "生成pdf", "轉成pdf", "輸出pdf", "匯出pdf", "做成 pdf", "生成 pdf", "轉成 pdf", "輸出 pdf", "匯出 pdf", "做成報告", "生成報告", "轉成報告"]

def _normalize_history_for_chat(chat_key: str) -> list[dict]:
    short_term = get_chat_history(chat_key, limit=10) if chat_key else []
    long_term_text = ""
    if chat_key:
        try:
            col = _get_memory_collection()
            if col and col.count() > 0:
                recent_user = next((m["content"] for m in reversed(short_term) if m.get("role") == "user"), "")
                if recent_user:
                    results = col.query(
                        query_texts=[recent_user],
                        n_results=min(5, col.count()),
                        where={"chat_key": chat_key}
                    )
                    docs = results["documents"][0] if results["documents"] else []
                    metas = results["metadatas"][0] if results["metadatas"] else []
                    short_contents = {m.get("content", "") for m in short_term}
                    relevant = []
                    for doc, meta in zip(docs, metas):
                        clean = doc
                        for prefix in ("[claude] ", "[gpt] ", "[gemini] ", "[claude-long] "):
                            if clean.startswith(prefix):
                                clean = clean[len(prefix):]
                                break
                        if clean not in short_contents and len(clean) > 20:
                            relevant.append(f"[{meta.get('created_at','')}] {clean[:300]}")
                    if relevant:
                        long_term_text = "【以下是你過去相關的回應記錄】\n" + "\n---\n".join(relevant) + "\n【以上為歷史記錄，以下是近期對話】\n"
        except Exception as e:
            print(f"[Memory] 長期記憶搜尋失敗：{e}")
    cleaned = []
    if long_term_text:
        cleaned.append({"role": "user", "content": long_term_text})
        cleaned.append({"role": "assistant", "content": "好的，我已記得這些歷史記錄，請繼續。"})
    for item in short_term:
        role = item.get("role", "user")
        content = item.get("content", "")
        if not content:
            continue
        if role not in ("user", "assistant"):
            role = "user"
        for prefix in ("[claude] ", "[gpt] ", "[gemini] ", "[claude-long] "):
            if content.startswith(prefix):
                content = content[len(prefix):]
                break
        cleaned.append({"role": role, "content": content})
    return cleaned

def ai_claude(user_text: str, chat_key: str = "") -> str:
    history = _normalize_history_for_chat(chat_key)
    messages = history + [{"role": "user", "content": user_text}]
    resp = claude_client.messages.create(model="claude-sonnet-4-6", max_tokens=1200, system=SYSTEM_PROMPT, messages=messages)
    reply = (resp.content[0].text or "").strip()
    if chat_key:
        save_chat_history(chat_key, "user", user_text)
        save_chat_history(chat_key, "assistant", f"[claude] {reply}")
    return reply

def ai_claude_long(user_text: str, chat_key: str = "") -> str:
    history = _normalize_history_for_chat(chat_key)
    messages = history + [{"role": "user", "content": user_text}]
    resp = claude_client.messages.create(model="claude-sonnet-4-6", max_tokens=2500, system=SYSTEM_PROMPT, messages=messages)
    reply = (resp.content[0].text or "").strip()
    if chat_key:
        save_chat_history(chat_key, "user", user_text)
        save_chat_history(chat_key, "assistant", f"[claude-long] {reply}")
    return reply

def ai_chatgpt(user_text: str, chat_key: str = "") -> str:
    if not openai_client:
        return ai_claude(user_text, chat_key)
    history = _normalize_history_for_chat(chat_key)
    messages = [{"role": "system", "content": SYSTEM_PROMPT}] + history + [{"role": "user", "content": user_text}]
    resp = openai_client.chat.completions.create(model="gpt-4.1-mini", messages=messages, temperature=0.4, max_tokens=1800)
    reply = (resp.choices[0].message.content or "").strip()
    if chat_key:
        save_chat_history(chat_key, "user", user_text)
        save_chat_history(chat_key, "assistant", f"[gpt] {reply}")
    return reply

def ai_gemini(user_text: str, chat_key: str = "") -> str:
    if not GEMINI_API_KEY:
        return ai_claude(user_text, chat_key)
    history = _normalize_history_for_chat(chat_key)
    history_text = "\n".join([f"{m['role']}: {m['content']}" for m in history[-10:]])
    prompt = f"{SYSTEM_PROMPT}\n\n以下是近期對話：\n{history_text}\n\n使用者最新問題：\n{user_text}"
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    req = urllib.request.Request(url, data=json.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"}, method="POST")
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode("utf-8"))
        reply = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
        if not reply:
            return ai_claude(user_text, chat_key)
        if chat_key:
            save_chat_history(chat_key, "user", user_text)
            save_chat_history(chat_key, "assistant", f"[gemini] {reply}")
        return reply
    except Exception as e:
        print(f"[Gemini Error] {e}")
        return ai_claude(user_text, chat_key)

def ai_router(user_text: str, chat_key: str = "", forced_model: str = "") -> str:
    text_l = (user_text or "").lower().strip()
    if forced_model == "claude":
        return ai_claude(user_text, chat_key)
    if forced_model == "gpt":
        return ai_chatgpt(user_text, chat_key)
    if forced_model == "gemini":
        return ai_gemini(user_text, chat_key)
    if any(k in text_l for k in AUTO_FINANCE_KEYWORDS):
        return ai_claude(user_text, chat_key)
    if any(k in text_l for k in AUTO_FILE_KEYWORDS):
        return ai_gemini(user_text, chat_key)
    return ai_chatgpt(user_text, chat_key)

def classify_report_topic(user_text: str) -> str:
    text_l = (user_text or "").lower()
    macro_keywords = ["戰爭", "衝突", "重建", "制裁", "關稅", "降息", "升息", "聯準會", "fed", "通膨", "景氣", "衰退", "地緣政治", "原油", "油價", "中東", "美元", "公債", "殖利率", "金融市場", "總經", "非農", "cpi", "pce", "失業率", "財政", "重建行情"]
    equity_keywords = ["股票", "股價", "公司", "企業", "財報", "估值", "獲利", "eps", "ai", "gpu", "供應鏈", "半導體", "伺服器", "金融股", "銀行股", "科技股", "產業", "競爭力", "台積電", "nvidia", "nvda", "amd", "avgo", "smci", "aapl", "meta"]
    product_keywords = ["基金", "債券", "etf", "eln", "結構型", "信託", "質借", "lombard", "票據", "商品", "配息", "收益", "信用債", "投資等級", "高收益債", "可轉債"]
    if any(k in text_l for k in macro_keywords):
        return "macro"
    if any(k in text_l for k in equity_keywords):
        return "equity"
    if any(k in text_l for k in product_keywords):
        return "product"
    return "general"

def build_macro_prompt(user_text: str) -> str:
    return f"""你是一位頂級總經與跨資產策略研究員。請根據以下主題撰寫繁體中文深度研究報告。
研究主題：{user_text}
【封面摘要】【一、事件與市場背景】【二、行情形成機制】【三、受惠產業與資產主線】
【四、金融市場影響】【五、情境分析】【六、投資機會與策略建議】【七、主要風險與反證】【八、結論】
每節至少2到4段。禁止Markdown符號。直接輸出完整報告正文。"""

def build_equity_prompt(user_text: str) -> str:
    return f"""你是一位資深產業與股票研究員。請根據以下主題撰寫繁體中文深度研究報告。
研究主題：{user_text}
【封面摘要】【一、產業與公司背景】【二、成長動能與投資邏輯】【三、競爭格局與關鍵優勢】
【四、財務與估值觀察】【五、市場可能如何交易這個題材】【六、投資機會與布局方式】【七、主要風險與反證】【八、結論】
每節至少2到4段。禁止Markdown符號。直接輸出完整報告正文。"""

def build_product_prompt(user_text: str) -> str:
    return f"""你是一位銀行財富管理研究員。請根據以下主題撰寫繁體中文深度研究報告。
研究主題：{user_text}
【封面摘要】【一、商品定位與市場背景】【二、報酬來源與運作機制】【三、適合客群與資產配置角色】
【四、優勢、限制與常見誤解】【五、目前市場環境下的投資價值】【六、投資建議與配置思維】【七、主要風險與注意事項】【八、結論】
每節至少2到4段。禁止Markdown符號。直接輸出完整報告正文。"""

def build_general_prompt(user_text: str) -> str:
    return f"""你是一位資深投資研究員。請根據以下主題撰寫繁體中文深度研究報告。
研究主題：{user_text}
【封面摘要】【一、主題背景】【二、核心分析】【三、市場與產業影響】【四、投資機會】【五、主要風險】【六、結論與建議】
每節至少2到4段。禁止Markdown符號。直接輸出完整報告正文。"""

def build_pdf_report_content(user_text: str, chat_key: str = "") -> str:
    topic_type = classify_report_topic(user_text)
    prompt_map = {"macro": build_macro_prompt, "equity": build_equity_prompt, "product": build_product_prompt}
    prompt = prompt_map.get(topic_type, build_general_prompt)(user_text)
    return ai_claude_long(prompt, chat_key)

def build_transcript_summary(transcript: str, chat_key: str = "") -> str:
    prompt = f"""你是一位專業會議紀錄助理。請將以下逐字稿整理為繁體中文重點摘要。
要求：先寫【會議摘要】，再寫【重點整理】，再寫【待辦事項】。條列清楚、內容具體。禁止Markdown符號。
逐字稿：{transcript}"""
    return ai_claude(prompt, chat_key)

def build_transcript_pdf_content(transcript: str, summary: str, chat_key: str = "") -> str:
    prompt = f"""你是一位專業研究助理，請把以下會議逐字稿與摘要整理成繁體中文正式會議報告。
【封面摘要】【一、會議背景】【二、會議重點】【三、逐字稿重點整理】【四、結論】【五、待辦事項】
語氣正式。不要使用Markdown符號。
會議摘要：{summary}
逐字稿：{transcript[:120000]}"""
    return ai_claude_long(prompt, chat_key)

# ==============================
# 文章儲存功能
# ==============================
def geocode_location(location_name: str) -> tuple[float, float] | tuple[None, None]:
    api_key = os.getenv("GOOGLE_MAPS_API_KEY", "")
    if not api_key or not location_name:
        return None, None
    try:
        import urllib.parse
        query = urllib.parse.quote(location_name)
        url = f"https://maps.googleapis.com/maps/api/geocode/json?address={query}&key={api_key}&language=zh-TW"
        req = urllib.request.Request(url)
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read().decode("utf-8"))
        if data.get("status") == "OK" and data.get("results"):
            loc = data["results"][0]["geometry"]["location"]
            return float(loc["lat"]), float(loc["lng"])
    except Exception as e:
        print(f"[Geocode] 失敗：{e}")
    return None, None

def _parse_claude_article_response(text: str) -> dict:
    result = {"title": "", "summary": "", "category": "other", "location_name": ""}
    lines = text.splitlines()
    summary_lines = []
    in_summary = False
    for line in lines:
        line = line.strip()
        if line.startswith("標題："):
            result["title"] = line.replace("標題：", "").strip()
        elif line.startswith("分類："):
            cat = line.replace("分類：", "").strip().lower()
            cat_map = {
                "finance": "finance", "財經": "finance", "投資": "finance",
                "food": "food", "美食": "food", "餐廳": "food", "小吃": "food",
                "travel": "travel", "旅遊": "travel", "景點": "travel", "觀光": "travel",
                "shopping": "shopping", "購物": "shopping",
                "other": "other", "其他": "other",
            }
            result["category"] = cat_map.get(cat, "other")
        elif line.startswith("地點："):
            result["location_name"] = line.replace("地點：", "").strip()
            if result["location_name"] == "無":
                result["location_name"] = ""
        elif line.startswith("重點："):
            in_summary = True
        elif in_summary and line:
            summary_lines.append(line)
    result["summary"] = result["title"] + "\n重點：\n" + "\n".join(summary_lines) if summary_lines else text
    return result

ARTICLE_PROMPT_SUFFIX = """
格式如下（請嚴格照此格式，每行一個欄位）：
標題：xxx
分類：finance 或 food 或 travel 或 shopping 或 other
地點：地點名稱（若有任何店名、景點、地名、城市、國家請填入，例如「四國自動車博物館」「鼎泰豐信義店」「東京淺草」「桃園」；若完全沒有地點資訊才填「無」）
重點：
• xxx
• xxx
• xxx
注意：
- 只要圖片或內容有提到任何地名、店名、景點名稱，一律填入地點欄位
- 日本、韓國、歐洲等海外地點也要填，不要填「無」
- 分類判斷：博物館/景點/旅遊=travel，餐廳/小吃/咖啡=food，投資/市場/財經=finance，購物/商品=shopping
"""

def save_article_text(ck: str, content: str) -> str:
    prompt = f"請用繁體中文為以下內容產生摘要，並判斷分類與地點。\n{ARTICLE_PROMPT_SUFFIX}\n內容：\n{content[:3000]}"
    import time
    for attempt in range(3):
        try:
            resp = claude_client.messages.create(
                model="claude-sonnet-4-6", max_tokens=600,
                messages=[{"role": "user", "content": prompt}]
            )
            break
        except Exception as e:
            if "529" in str(e) or "overloaded" in str(e).lower():
                if attempt < 2:
                    time.sleep(5 * (attempt + 1))
                    continue
            raise e
    parsed = _parse_claude_article_response((resp.content[0].text or "").strip())
    title = parsed["title"] or content[:30]
    summary = parsed["summary"]
    category = parsed["category"]
    location_name = parsed["location_name"]
    source_type = "url" if content.startswith("http") else "text"
    lat, lng = geocode_location(location_name) if location_name else (None, None)
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO articles (title, content, summary, source_type, category, location_name, lat, lng, is_read, show_on_map)
        VALUES (:t, :c, :s, :st, :cat, :loc, :lat, :lng, FALSE, TRUE)
        """), {"t": title, "c": content[:5000], "s": summary, "st": source_type, "cat": category, "loc": location_name, "lat": lat, "lng": lng})
    return summary

def save_article_image(image_data: bytes, message_id: str) -> str:
    image_b64 = _base64.b64encode(image_data).decode("utf-8")
    import time
    for attempt in range(3):
        try:
            resp = claude_client.messages.create(
                model="claude-sonnet-4-6", max_tokens=600,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": image_b64}},
                    {"type": "text", "text": f"請用繁體中文描述這張圖片的內容並產生摘要，並判斷分類與地點。\n{ARTICLE_PROMPT_SUFFIX}"}
                ]}]
            )
            break
        except Exception as e:
            if "529" in str(e) or "overloaded" in str(e).lower():
                if attempt < 2:
                    time.sleep(5 * (attempt + 1))
                    continue
            raise e
    parsed = _parse_claude_article_response((resp.content[0].text or "").strip())
    title = parsed["title"] or "圖片文章"
    summary = parsed["summary"]
    category = parsed["category"]
    location_name = parsed["location_name"]
    lat, lng = geocode_location(location_name) if location_name else (None, None)
    with engine.begin() as conn:
        conn.execute(text("""
        INSERT INTO articles (title, content, summary, source_type, image_url, category, location_name, lat, lng, is_read, show_on_map)
        VALUES (:t, :c, :s, 'image', :img, :cat, :loc, :lat, :lng, FALSE, TRUE)
        """), {"t": title, "c": "（圖片）", "s": summary, "img": f"line_image_{message_id}", "cat": category, "loc": location_name, "lat": lat, "lng": lng})
    return summary

def get_unread_articles(limit: int = 15) -> list:
    with engine.begin() as conn:
        rows = conn.execute(text("""
        SELECT id, title, source_type, created_at FROM articles
        WHERE is_read = FALSE ORDER BY created_at DESC LIMIT :n
        """), {"n": limit}).fetchall()
    return rows

def mark_article_read(article_id: int):
    with engine.begin() as conn:
        conn.execute(text("UPDATE articles SET is_read = TRUE WHERE id = :i"), {"i": article_id})

def get_article_detail(article_id: int):
    with engine.begin() as conn:
        row = conn.execute(text("""
        SELECT id, title, content, summary, source_type, is_read, created_at
        FROM articles WHERE id = :i
        """), {"i": article_id}).fetchone()
    return row

# ==============================
# Message handlers
# ==============================
@handler.add(MessageEvent, message=TextMessage)
def handle_text_message(event):
    _bot_api = getattr(_current_bot_api, "api", None) or line_bot_api
    try:
        text_raw = (event.message.text or "").strip()
        tl = text_raw.lower().strip()
        ck = chat_key_of(event)
        is_group = event.source.type in ("group", "room")
        print("[TEXT]", ck, repr(text_raw))
        if tl.startswith("/"):
            cmd = tl[1:].split()[0] if tl[1:].split() else ""
            raw_cmd = text_raw[1:]
        else:
            cmd = tl.split()[0] if tl.split() else tl
        # ── 海外債「查價群」：只回 /price 與 /help，其他一律靜默 ──
        if is_group and is_bond_query_group(ck):
            if not tl.startswith("/"):
                return
            if cmd == "help":
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=BOND_QUERY_HELP))
                return
            if cmd not in ("price", "p", "價格", "報價", "issuer", "coupon"):
                return
            # 查價群的 /coupon 只給查詢用法,設定類子指令不開放
            if cmd == "coupon":
                _sub = raw_cmd.split()[1].lower() if len(raw_cmd.split()) > 1 else ""
                if _sub in ("settarget", "subscribe", "unsubscribe", "subscribers", "table",
                            "設定推播", "訂閱", "取消訂閱", "訂閱名單", "off"):
                    return
            # /price 放行，往下走到 price 指令處理
        # ── 海外債群組：只回白名單指令，一般聊天/其他指令一律靜默 ──
        if is_group and is_bond_group_chat(ck):
            if not tl.startswith("/") or cmd not in BOND_GROUP_ALLOWED_CMDS:
                return
            if cmd == "help":
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=BOND_GROUP_HELP))
                return
            raw_cmd = text_raw
        parts = text_raw.split(" ", 1)

        # ── 判斷是否為 Albert 本人（非本人只能用 ELN 四個指令）──
        _albert_uid = os.getenv("LINE_USER_ID", "")
        _sender_uid = event.source.user_id if hasattr(event.source, "user_id") else ""
        is_albert = (_sender_uid == _albert_uid)
        _agent_allowed = ("list", "detail", "nc")
        _is_eln_channel = getattr(_current_bot_api, "is_eln", False)
        if not is_albert and not is_group and tl.startswith("/"):
            # ELN Bot 頻道的 /help:只顯示海外債指令
            # (ELN 的 /list /detail /nc 仍可正常使用,只是不主動列出,避免新同仁誤用)
            if _is_eln_channel and cmd in ("help", "?", "指令", "幫助"):
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=BOND_QUERY_HELP))
                return
            _ok = cmd in _agent_allowed
            # ELN Bot 頻道:同時開放海外債查詢指令(不燒 AI 的那些)
            if not _ok and _is_eln_channel and cmd in ELN_BOT_BOND_CMDS:
                _ok = True
            if not _ok and _is_eln_channel and cmd in ELN_BOT_BOND_HEAVY:
                if cmd == "sheet" and can_use_doc_cmd(_sender_uid):
                    _ok = True      # 名單內的投資輔銷可用
                elif cmd == "focus":
                    _bot_api.reply_message(event.reply_token, TextSendMessage(
                        text="/focus 目前僅開放固定收益科使用，請洽固定收益科協助產出。"))
                    return
                else:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(
                        text=f"/{cmd} 目前開放給投資輔銷同仁使用，"
                             "請洽轄區投資輔銷協助產出。"))
                    return
            if not _ok:
                # ELN 頻道:錯誤指令只提示海外債用法,不列出 ELN 指令
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=(BOND_QUERY_HELP if _is_eln_channel else
                          "可用指令：\n/list 姓名\n/list detail 姓名\n/detail 商品代號\n/nc YYYYMM 姓名")))
                return
            if cmd in _agent_allowed:
                ck = ELN_PERSONAL_CHAT_KEY
        if not is_albert and not is_group and not tl.startswith("/"):
            return
        
        if is_group and not tl.startswith("/"):
            return
            
        # ==========================================
        # 內規專屬指令攔截 (使用 Claude 全文理解，完美舉一反三)
        # ==========================================
        if text_raw.startswith("/內規"):
            actual_query = text_raw.replace("/內規", "").strip()
            if not actual_query:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請在指令後面加上想查詢的內容喔！\n例如：/內規 Lombard lending 最高可以到幾歲？"))
                return
                
            try:
                file_path = Path("regulations.txt")
                if not file_path.exists():
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="❌ 找不到 regulations.txt，請確認已將法規檔案上傳至系統。"))
                    return
                    
                regulation_text = file_path.read_text(encoding="utf-8")
                
                # 防縮排錯誤的 Prompt 組合方式
                prompt_lines = [
                    "你現在是銀行的法遵與內部規範專家。請根據以下【內部規範全文】，直接且精準回答同仁的問題。",
                    "",
                    "【嚴格限制】",
                    "1. 絕對不要輸出任何「因為文本是程式碼...」或「無法回答」的廢話警告。",
                    "2. 絕對不要在結尾補充「可以這樣跟客戶/專員說」的話術。",
                    "3. 嚴禁使用 Markdown 語法 (例如 **, ##, --- 等)，請用純文字或 Emoji 條列排版。",
                    "",
                    "【內部規範全文】",
                    regulation_text,
                    "",
                    "【同仁問題】",
                    actual_query
                ]
                prompt = "\n".join(prompt_lines)

                raw_answer = ai_claude_long(prompt, chat_key=ck)
                # 物理綁定警語，不讓 AI 決定
                final_answer = "⚠️ 僅供參考，本回覆由 AI 統整，不代表總行最終內規解釋。\n\n" + raw_answer
                
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=final_answer[:4900]))
                
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 內部規範查詢失敗：{e}"))
            return
        # ==========================================
        
        if cmd in ("help", "?", "指令", "幫助"):
            help_arg = parts[1].strip().lower() if len(parts) > 1 else ""
            if is_group:
                msg = "群組可用指令：\n/detail <商品代號>：查詢標的完整狀況（支援模糊搜尋）\n/list：列出所有可查商品代號\n/list <姓名>：查詢該理專的持倉\n/nc YYYYMM：查詢該月閉鎖期打開的商品\n/nc YYYYMM <姓名>：查詢該理專的閉鎖期商品\n/end YYYYMM：查詢該月到期商品\n"
            else:
                if help_arg in ("bond", "債", "海外債", "coupon", "固收"):
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=BOND_GROUP_HELP))
                    return
                if help_arg in ("alert", "警示"):
                    msg = ("🔔 Alert 指令說明\n─────────────────\n"
                           "/alert add <標的> <價格> above/below\n/alert add <標的> above/below <價格>\n"
                           "/alert add <標的> ma20 above/below\n/alert add <標的> ma5 cross ma20\n"
                           "/alert add <標的> ma5 under ma20\n/alert list\n/alert del <編號>\n"
                           "─────────────────\n別名：dxy / spx / ndx / sox / vix / ust10y / gold / silver / oil\n"
                           "範例：\n/alert add dxy below 100\n/alert add ust10y above 45\n"
                           "/alert add NVDA ma20 above\n/alert add NVDA ma5 cross ma20")
                elif help_arg in ("eln",):
                    msg = ("📊 ELN 指令說明\n─────────────────\n"
                           "/calc — 上傳 Excel 計算並保存\n/list — 列出所有可查商品代號\n"
                           "/detail <代號> — 查詢單筆 KO/KI/狀態\n/eln upload — 上傳 Excel 並同步到 Supabase\n"
                           "/eln run — 立即重跑最新 ELN\n/eln history — 查看歷史 Excel\n"
                           "/eln result — 查看最近結果\n/runnow — 手動執行追蹤\n"
                           "/tracklog — 查看最近排程紀錄\n/end YYYYMM — 查詢指定月份到期商品\n"
                           "/nc YYYYMM — 查詢指定月份閉鎖期打開的商品\n"
                           "/nc YYYYMM 姓名 — 查詢該理專的閉鎖期打開商品\n"
                           "/chart 商品代號 — 產生走勢圖+防守線（KO/KI/Strike）")
                elif help_arg in ("report", "pdf", "報告", "簡報"):
                    msg = ("📑 報告 / PDF 指令說明\n─────────────────\n"
                           "/report <主題>\n/report <主題> brief/client/academic/hybrid\n"
                           "/report <主題> custom <說明>\n/pdf market <內容>\n/pdf make <內容>\n"
                           "自然語言也可直接說：\n請幫我做一份 XX 的 pdf")
                elif help_arg in ("save", "文章", "儲存"):
                    msg = ("📚 文章儲存指令說明\n─────────────────\n"
                           "/save <文字或網址> — 儲存文章並自動摘要\n"
                           "/unread — 查看未讀文章清單\n/read <編號> — 標記文章為已讀\n"
                           "/article <編號> — 查看文章摘要內容\n/del <編號> — 刪除文章\n"
                           "/web — 開啟文章庫網頁\n─────────────────\n"
                           "直接傳圖片給龍蝦 → 自動儲存並分析\n"
                           "範例：\n/save https://tw.news.yahoo.com/...\n/save 這篇文章說AI將取代50%工作...")
                else:
                    msg = ("🦞 龍蝦指令清單\n─────────────────\n"
                           "📊 ELN\n/calc  /list  /detail\n/eln upload  /eln run  /eln history  /eln result\n"
                           "/runnow  /tracklog  /end\n─────────────────\n"
                           "📰 財經\n/daily  /daily cache  /market\n/bonddaily  /bonddaily cache → 債券市場日報\n─────────────────\n"
                           "📑 報告\n/report  /pdf\n─────────────────\n"
                           "📚 知識庫\n/內規 <問題> → 查詢 內部法規\n/kb <問題> → 查詢 Chroma 知識庫\n/kb上傳 → 上傳檔案\n/kb清單 → 查看文件清單\n─────────────────\n"
                           "🔔 警示\n/alert add  /alert list  /alert del\n輸入 /help alert 看完整範例\n─────────────────\n"
                           "📚 文章庫\n/save  /unread  /read  /article  /del  /web\n直接傳圖片 → 自動儲存分析\n輸入 /help save 看完整說明\n─────────────────\n"
                           "📊 基金淨值 & 債券報價\n/fundnav → 手動更新基金淨值\n/bondnav → 手動觸發債券報價更新（94筆，約30分鐘）\n/tracklog → 查看執行記錄\n─────────────────\n"
                           "💰 海外債專區\n/help bond → 配息雷達、發行機構、報價異動、信評、到價通知…全部指令\n─────────────────\n"
                           "📧 其他\n/mail  /invest  /forget  /spending\n上傳錄音 → 自動逐字稿 / 摘要\n上傳檔案 → 自動分析\n─────────────────\n"
                           "進階說明：/help bond、/help alert、/help eln、/help report、/help save")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg))
            return
        transcript_cache = db_get_transcript_cache(ck)
        if transcript_cache:
            if any(x in tl for x in ["不用", "不用了", "先不用", "取消", "不用做"]):
                db_clear_transcript_cache(ck)
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="👌 好的，已保留逐字稿與摘要回覆，不另外生成 PDF。"))
                return
            if any(x in tl for x in ["做成pdf", "生成pdf", "轉成pdf", "做成 pdf", "生成 pdf", "轉成 pdf", "輸出pdf", "輸出 pdf"]):
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📄 正在根據逐字稿重點生成 PDF，請稍候..."))
                try:
                    from pdf_generator import create_and_upload_pdf
                    report_text = build_transcript_pdf_content(transcript_cache["transcript"], transcript_cache["summary"], chat_key=ck)
                    link = create_and_upload_pdf("analysis", report_text, "會議重點報告")
                    db_clear_transcript_cache(ck)
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"✅ 會議重點 PDF 已生成完成！\n\n{link}"))
                except Exception as e:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ PDF 生成失敗：{str(e)[:250]}"))
                return
        if cmd in ("send", "skip"):
            arg = parts[1].strip().lower() if len(parts) > 1 else ""
            if not arg:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請指定編號或 all\n範例：/send 1 /skip 2 /send all"))
                return
            with engine.begin() as conn:
                rows = conn.execute(text("SELECT id, target_id, agent_name, bond_id, status, msg FROM eln_pending_notifications WHERE chat_key=:k ORDER BY id"), {"k": ck}).fetchall()
            if not rows:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有待確認的通知。"))
                return
            if arg == "all":
                targets = list(rows)
            else:
                try:
                    idx = int(arg) - 1
                    if idx < 0 or idx >= len(rows):
                        raise ValueError
                    targets = [rows[idx]]
                except ValueError:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"編號不正確，請輸入 1～{len(rows)} 或 all"))
                    return
            if cmd == "send":
                sent, failed = 0, 0
                for row in targets:
                    try:
                        eln_group_bot_api.push_message(row.target_id, TextSendMessage(text=row.msg[:4900]))
                        sent += 1
                    except Exception as e:
                        failed += 1
                    with engine.begin() as conn:
                        conn.execute(text("DELETE FROM eln_pending_notifications WHERE id=:i"), {"i": row.id})
                result_text = f"✅ 已發送 {sent} 筆" + (f"，失敗 {failed} 筆" if failed else "")
            else:
                for row in targets:
                    with engine.begin() as conn:
                        conn.execute(text("DELETE FROM eln_pending_notifications WHERE id=:i"), {"i": row.id})
                result_text = f"⏭️ 已略過 {len(targets)} 筆"
            with engine.begin() as conn:
                remaining = conn.execute(text("SELECT COUNT(*) FROM eln_pending_notifications WHERE chat_key=:k"), {"k": ck}).scalar()
            result_text += f"\n\n還有 {remaining} 筆待處理，打 /send list 查看" if remaining > 0 else "\n\n✅ 所有通知已處理完畢"
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=result_text))
            return
        if cmd == "invest":
            db_invest_set(ck, "await_image")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📰 請上傳新聞截圖\n\n收到圖片後，我會請你補上投資理由和標的。"))
            return
        invest_mode, invest_image = db_invest_get(ck)
        if invest_mode == "await_reason" and invest_image:
            raw = text_raw.strip()
            reason = ""
            targets = ""
            for line in raw.replace("，", ",").splitlines():
                l = line.strip()
                if l.startswith("理由"):
                    reason = l.split("：", 1)[-1].split(":", 1)[-1].strip()
                elif l.startswith("標的"):
                    targets = l.split("：", 1)[-1].split(":", 1)[-1].strip()
            if not reason and not targets:
                reason = raw
            db_invest_set(ck, "")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="✍️ 整理中，請稍候..."))
            try:
                posts = generate_invest_post(invest_image, reason, targets)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=posts[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"生成失敗：{str(e)[:200]}"))
            return
        if cmd.startswith("bonddaily"):
            parts = text_raw.split(" ", 1)
            _arg = parts[1].strip() if len(parts) > 1 else ""
            # /bonddaily focus 浮動利率債 → 設定當期主打方向(每日操作思維會帶到)
            # /bonddaily focus        → 查看目前設定
            # /bonddaily focus off    → 取消
            if _arg.lower().startswith("focus"):
                fv = _arg[5:].strip()
                fpath = "/data/bond_focus.json" if os.path.isdir("/data") else "/tmp/bond_focus.json"
                try:
                    if not fv:
                        cur, curisk = "", ""
                        try:
                            from bond_daily_report import get_daily_focus_full
                            cur, curisk = get_daily_focus_full()
                        except Exception:
                            pass
                        body = (f"📌 目前當期主打方向：{cur}" if cur else "📌 尚未設定當期主打方向")
                        if curisk:
                            body += f"\n⚠️ 必講風險：{curisk}"
                        _bot_api.reply_message(event.reply_token, TextSendMessage(
                            text=body + "\n\n設定：/bonddaily focus 零息債 | 風險:無配息、有提前買回、信用風險"
                                        "\n取消：/bonddaily focus off"))
                        return
                    if fv.lower() in ("off", "clear", "取消"):
                        if os.path.exists(fpath):
                            os.remove(fpath)
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="✅ 已取消當期主打方向，日報恢復純市場觀察。"))
                        return
                    # 支援「方向 | 風險:xxx」格式，風險會被強制寫進同一段
                    fdir, frisk = fv, ""
                    if "|" in fv or "｜" in fv:
                        seg = re.split(r"[|｜]", fv, 1)
                        fdir = seg[0].strip()
                        frisk = re.sub(r"^\s*(風險|risk)\s*[:：]?\s*", "", seg[1].strip(), flags=re.I)
                    with open(fpath, "w", encoding="utf-8") as f:
                        json.dump({"focus": fdir, "risk": frisk}, f, ensure_ascii=False)
                    msg_f = (f"✅ 已設定當期主打方向：{fdir}")
                    if frisk:
                        msg_f += f"\n⚠️ 必講風險：{frisk}"
                    else:
                        msg_f += "\n（未設定必講風險，建議加上，格式：/bonddaily focus 零息債 | 風險:無配息、有提前買回、信用風險）"
                    msg_f += ("\n\n明天起的債券日報，今日操作思維最後會用一句話把當天市場連結到這個方向"
                              "（只講產品類型邏輯，不提具體標的與價格）"
                              + ("，並一併帶出上述風險。" if frisk else "。")
                              + "\n取消請打 /bonddaily focus off")
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg_f))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 設定失敗：{str(e)[:200]}"))
                return
            use_cache = _arg.lower() == "cache"
            if use_cache:
                try:
                    with engine.begin() as conn:
                        row = conn.execute(text("SELECT report_text FROM bond_daily_report_cache ORDER BY created_at DESC LIMIT 1")).fetchone()
                    if row:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text=row[0][:4900]))
                    else:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="尚無快取債券日報，請用 /bonddaily 產生最新版本。"))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"讀取快取失敗: {e}"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="債券日報產生中，請稍候約30-60秒..."))
            try:
                from bond_daily_report import generate_report as bond_generate_report, save_report_to_db as bond_save_report_to_db
                bond_report = bond_generate_report()
                bond_save_report_to_db(bond_report)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=bond_report[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"債券日報產生失敗: {e}"))
            return
        if cmd.startswith("daily"):
            parts = text_raw.split(" ", 1)
            use_cache = len(parts) > 1 and parts[1].strip().lower() == "cache"
            if use_cache:
                try:
                    with engine.begin() as conn:
                        row = conn.execute(text("SELECT report_text FROM daily_report_cache ORDER BY created_at DESC LIMIT 1")).fetchone()
                    if row:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text=row[0][:4900]))
                    else:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="尚無快取日報，請用 /daily 產生最新版本。"))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"讀取快取失敗: {e}"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="產生中，請稍候約30秒..."))
            try:
                from daily_report import generate_report, save_report_to_db
                report, image_url, weekly_calendar = generate_report()
                save_report_to_db(report)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=report[:4900]))
                if image_url:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📊 今日市場摘要圖\n{image_url}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"日報產生失敗: {e}"))
            return
        if cmd == "settarget":
            targets = load_targets()
            if event.source.type == "group":
                targets["default"] = event.source.group_id
                targets["default_type"] = "group"
            elif event.source.type == "room":
                targets["default"] = event.source.room_id
                targets["default_type"] = "room"
            else:
                targets["default"] = event.source.user_id
                targets["default_type"] = "user"
            save_targets(targets)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="已設定為預設推播對象"))
            return
        if cmd == "eln":
            sub_parts = text_raw.split()
            sub = sub_parts[1].lower() if len(sub_parts) > 1 else ""
            if sub == "upload":
                db_set_await(ck, True)
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📥 請直接上傳 ELN Excel 檔案，我會計算並同步保存到 Supabase。"))
                return
            if sub == "run":
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="🔄 正在重新計算最新 ELN，請稍候..."))
                try:
                    try:
                        from eln_storage import download_latest_eln
                        latest_file = download_latest_eln("/tmp/latest_eln.xlsx")
                    except Exception:
                        latest_file = "/tmp/latest_eln.xlsx"
                    summary, top5_lines, detail_map, agent_name_map = run_autotracking(latest_file)
                    db_save_result(ck, summary, top5_lines, detail_map, agent_name_map)
                    msg = "✅ ELN 已重新計算完成\n\n" + ("\n".join(top5_lines[:5]) if top5_lines else (summary or "沒有可顯示摘要"))
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=msg[:4900]))
                except Exception as e:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ ELN 計算失敗：{str(e)[:250]}"))
                return
            if sub == "history":
                try:
                    from eln_storage import list_history
                    items = list_history()
                    msg = "📁 ELN Excel 歷史版本\n\n" + ("\n".join(items[:20]) if items else "目前沒有歷史 Excel")
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg[:4900]))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"讀取歷史失敗：{str(e)[:250]}"))
                return
            if sub == "result":
                summary = db_get_report(ck)
                if not summary:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有 ELN 結果，請先 /calc 或 /eln run。"))
                    return
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=summary[:4900]))
                return
        if cmd == "list":
            from collections import defaultdict
            list_parts = text_raw.split(" ", 2)
            # 判斷是否為 /list detail 姓名
            is_detail_mode = len(list_parts) > 1 and list_parts[1].strip().lower() == "detail"
            if is_detail_mode:
                name_filter = list_parts[2].strip() if len(list_parts) > 2 else ""
            else:
                name_filter = list_parts[1].strip() if len(list_parts) > 1 else ""
            bonds = db_list_bonds(ck, limit=200)
            if not bonds:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無已保存結果。請先 /calc 上傳 Excel。"))
                return
            # /list detail 姓名：顯示完整 detail，只顯示比價中的商品
            if is_detail_mode:
                if not name_filter:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入理專名稱\n例：/list detail 小美"))
                    return
                matched_details = []
                seen = set()
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()]
                    if any(name_filter in a for a in agents) and bond_id not in seen:
                        # 只顯示比價中（排除已到期、已KO）
                        status_tag = bond_status_tag(detail)
                        if status_tag == "":  # 沒有特殊標籤 = 比價中
                            matched_details.append(detail)
                            seen.add(bond_id)
                if not matched_details:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到「{name_filter}」比價中的持倉。"))
                    return
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"👤 {name_filter} 比價中商品（共 {len(matched_details)} 筆）"))
                for det in matched_details:
                    push_long_message(_bot_api, ck.split(":", 1)[1], det[:4900])
                return
            detail_map = {bond_id: bond_status_tag(detail) for bond_id, _, detail in bonds}
            if name_filter:
                matched = []
                seen = set()
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()]
                    if any(name_filter in a for a in agents) and bond_id not in seen:
                        matched.append((bond_id, detail_map.get(bond_id, "")))
                        seen.add(bond_id)
                if not matched:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到「{name_filter}」的持倉。"))
                    return
                lines = [f"👤 {name_filter} 的持倉（共 {len(matched)} 筆）：\n"] + [f"   • {b}{tag}" for b, tag in matched]
            else:
                grouped = defaultdict(list)
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()] or ["未指定"]
                    for agent in agents:
                        if bond_id not in [b for b, _ in grouped[agent]]:
                            grouped[agent].append((bond_id, detail_map.get(bond_id, "")))
                lines = [f"📋 全部商品（共 {len(set(b for b,_,_ in bonds))} 筆，按理專排列）：\n"]
                for agent, bond_list in sorted(grouped.items()):
                    lines.append(f"👤 {agent}（{len(bond_list)} 筆）")
                    lines += [f"   • {b}{tag}" for b, tag in bond_list]
            full_text = "\n".join(lines)
            chunks = []
            current = ""
            for line in full_text.split("\n"):
                if len(current) + len(line) + 1 > 4800:
                    chunks.append(current)
                    current = line
                else:
                    current = current + "\n" + line if current else line
            if current:
                chunks.append(current)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=chunks[0]))
            for chunk in chunks[1:]:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=chunk))
            return
        if cmd.startswith("calc") or cmd.startswith("clac"):
            parts = raw_cmd.split(" ", 1)
            if len(parts) > 1 and parts[1].strip():
                expr = parts[1].strip()
                if not re.fullmatch(r"[0-9\.\+\-\*\/\(\)\s]+", expr):
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="算式格式錯誤"))
                    return
                try:
                    result = eval(expr, {"__builtins__": {}})
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"{expr} = {result}"))
                    return
                except Exception:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="算式錯誤"))
                    return
            db_set_await(ck, True)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="收到！請直接把 Excel 檔案傳給我（用 LINE 的『檔案』上傳），我會計算並保存結果。"))
            return
        if cmd == "report" and len(raw_cmd.strip().split()) == 1:
            summary = db_get_report(ck)
            if not summary:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無已保存結果，請先 /calc 上傳 Excel。"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=summary[:4900]))
            return
        if cmd.startswith("market"):
            parts = text_raw.split(" ", 1)
            if len(parts) < 2 or not parts[1].strip():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入新聞內容和推薦標的"))
                return
            content = generate_market_content(parts[1].strip())
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=content[:4900]))
            return
        if (not tl.startswith("/pdf")) and any(k in tl for k in PDF_NL_KEYWORDS):
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📄 正在整理內容並生成研究報告 PDF，請稍候..."))
            try:
                from pdf_generator import create_and_upload_pdf
                report_text = build_pdf_report_content(text_raw, chat_key=ck)
                link = create_and_upload_pdf("analysis", report_text, "AI自動生成研究報告")
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"✅ 研究報告 PDF 已生成完成！\n\n{link}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ PDF 生成失敗：{str(e)[:250]}"))
            return
        if cmd.startswith("pdf"):
            from pdf_generator import create_and_upload_pdf
            parts = text_raw.split(" ", 2)
            sub = parts[1].strip().lower() if len(parts) > 1 else ""
            if sub == "market":
                if len(parts) < 3 or not parts[2].strip():
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入內容\n範例：/pdf market 美股反彈，推薦PIMCO"))
                    return
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="產生市場觀點 PDF 中，請稍候..."))
                try:
                    content = generate_market_content(parts[2].strip())
                    link = create_and_upload_pdf("market", content)
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📄 市場觀點 PDF 已產生！\n\n{link}"))
                except Exception as e:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"PDF 產生失敗: {e}"))
                return
            if sub == "make":
                content_text = parts[2].strip() if len(parts) > 2 else ""
                if not content_text:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請在指令後面直接輸入內容\n\n範例：\n/pdf make 美伊戰爭後的重建行情研究報告"))
                    return
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="整理內容並產生研究報告 PDF 中，請稍候..."))
                try:
                    report_text = build_pdf_report_content(content_text, chat_key=ck)
                    link = create_and_upload_pdf("analysis", report_text, "自訂研究報告")
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📄 研究報告 PDF 已產生！\n\n{link}"))
                except Exception as e:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"PDF 產生失敗: {e}"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="PDF 指令用法：\n/pdf market <內容> → 市場觀點 PDF\n/pdf make <內容> → 研究報告 PDF"))
            return
        if cmd.startswith("report"):
            parts = text_raw.split(" ")
            style_codes = {"ib", "brief", "client", "academic", "hybrid", "custom"}
            style_names = {"ib": "投資銀行", "brief": "簡報摘要", "client": "客戶推播", "academic": "學術研究", "hybrid": "混合風格", "custom": "自訂風格"}
            style = "ib"
            custom_prompt = ""
            if "custom" in [p.lower() for p in parts[2:]]:
                custom_idx = next(i for i, p in enumerate(parts) if p.lower() == "custom")
                topic = " ".join(parts[1:custom_idx]).strip()
                custom_prompt = " ".join(parts[custom_idx+1:]).strip()
                style = "custom"
                if not custom_prompt:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="自訂風格請在 custom 後面加上說明！\n\n範例：\n/report 台積電展望 custom 請用輕鬆幽默的風格"))
                    return
            elif len(parts) > 2 and parts[-1].lower() in style_codes:
                style = parts[-1].lower()
                topic = " ".join(parts[1:-1]).strip()
            else:
                topic = " ".join(parts[1:]).strip()
            if not topic:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入報告主題\n\n範例：\n/report 聯準會降息對債市影響\n/report 聯準會降息對債市影響 client"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📊 正在研究「{topic}」\n風格：{style_names.get(style,'投資銀行')}\n\n請稍候約60至90秒..."))
            try:
                from report_generator import generate_research_report
                link = generate_research_report(topic, ck.split(":", 1)[1], style=style, custom_prompt=custom_prompt)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📑 研究報告已完成！\n\n主題：{topic}\n風格：{style_names.get(style,'投資銀行')}\n\n{link}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"報告生成失敗：{e}"))
            return
        if cmd.startswith("alert"):
            parts = text_raw.split(" ")
            sub = parts[1].strip().lower() if len(parts) > 1 else ""
            if sub == "list":
                with engine.begin() as conn:
                    rows = conn.execute(text("SELECT id, symbol, alert_type, condition, target_value, ma_period FROM price_alerts WHERE chat_key=:k AND deleted=FALSE ORDER BY id ASC"), {"k": ck}).fetchall()
                if not rows:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有任何警示設定。"))
                    return
                msg = "目前警示清單：\n"
                for r in rows:
                    rid, sym, atype, cond, tval, maper = r
                    cond_str = "漲到" if cond == "above" else "跌到"
                    cross_str = "漲破" if cond == "above" else "跌破"
                    msg += f"#{rid} {sym} {cond_str if atype == 'price' else cross_str} {tval if atype == 'price' else f'MA{maper}'}\n"
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg.strip()))
                return
            if sub == "del":
                if len(parts) < 3:
                    with engine.begin() as conn:
                        rows = conn.execute(text("SELECT id, symbol, alert_type, condition, target_value, ma_period, trigger_count FROM price_alerts WHERE chat_key=:k AND deleted=FALSE ORDER BY id ASC"), {"k": ck}).fetchall()
                    if not rows:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有任何警示設定。"))
                        return
                    msg = "請輸入要刪除的編號：\n\n"
                    for r in rows:
                        rid, sym, atype, cond, tval, maper, tcount = r
                        remain = 2 - (tcount or 0)
                        msg += f"#{rid} {sym} {'漲到' if cond == 'above' else '跌到'} {tval if atype == 'price' else f'MA{maper}'}（剩餘{remain}次）\n"
                    msg += "\n輸入：/alert del <編號>"
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg.strip()))
                    return
                try:
                    del_id = int(parts[2])
                    with engine.begin() as conn:
                        conn.execute(text("UPDATE price_alerts SET deleted=TRUE WHERE id=:i AND chat_key=:k"), {"i": del_id, "k": ck})
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 警示 #{del_id} 已刪除"))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"刪除失敗：{e}"))
                return
            if sub == "add":
                if len(parts) < 5:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="格式說明：\n/alert add NVDA 190 above\n/alert add NVDA ma20 above\n/alert add NVDA ma5 cross ma20"))
                    return
                raw_symbol = parts[2]
                symbol = ALERT_TICKER_ALIAS.get(raw_symbol.lower(), raw_symbol).upper()
                p3 = parts[3].lower()
                p4 = parts[4].lower()
                p5 = parts[5].lower() if len(parts) > 5 else ""
                try:
                    if p3.startswith("ma") and p4 in ("cross", "under") and p5.startswith("ma"):
                        ma_short = int(p3[2:])
                        ma_long = int(p5[2:])
                        with engine.begin() as conn:
                            conn.execute(text("INSERT INTO price_alerts(chat_key, symbol, alert_type, condition, ma_short, ma_long) VALUES (:k, :s, 'ma_cross', :c, :ms, :ml)"), {"k": ck, "s": symbol, "c": p4, "ms": ma_short, "ml": ma_long})
                        label = "黃金交叉" if p4 == "cross" else "死亡交叉"
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 均線交叉警示已設定！\n標的：{symbol}\n條件：MA{ma_short} {label} MA{ma_long} 🔔"))
                        return
                    if p3.startswith("ma") and p4 in ("above", "below"):
                        ma_period = int(p3[2:])
                        with engine.begin() as conn:
                            conn.execute(text("INSERT INTO price_alerts(chat_key, symbol, alert_type, condition, ma_period) VALUES (:k, :s, 'ma', :c, :m)"), {"k": ck, "s": symbol, "c": p4, "m": ma_period})
                        cross = "漲破" if p4 == "above" else "跌破"
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 均線警示已設定！\n標的：{symbol}\n條件：{cross} MA{ma_period} 🔔"))
                        return
                    if p3 in ("above", "below"):
                        direction, value_str = p3, p4
                    elif p4 in ("above", "below"):
                        value_str, direction = p3, p4
                    else:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="方向請輸入 above 或 below"))
                        return
                    target = float(value_str)
                    with engine.begin() as conn:
                        conn.execute(text("INSERT INTO price_alerts(chat_key, symbol, alert_type, condition, target_value) VALUES (:k, :s, 'price', :c, :t)"), {"k": ck, "s": symbol, "c": direction, "t": target})
                    cond_str = "漲到" if direction == "above" else "跌到"
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 價格警示已設定！\n標的：{symbol}\n條件：{cond_str} {target} 🔔"))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"設定失敗：{e}"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="價格警示指令：\n/alert add <標的> <價格> <above/below>\n/alert add <標的> ma20 <above/below>\n/alert add <標的> ma5 cross ma20\n/alert list → 查看清單\n/alert del <編號> → 刪除"))
            return
        if cmd == "news pdf" or cmd == "news":
            from news_fetcher import generate_news_report
            from pdf_generator import create_and_upload_pdf
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="正在抓取最新財經新聞並整理中，請稍候約30秒..."))
            try:
                report = generate_news_report()
                link = create_and_upload_pdf("news", report)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📰 今日財經新聞摘要 PDF 已產生！\n\n{link}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"新聞抓取失敗: {e}"))
            return
        if cmd.startswith("chart"):
            parts = text_raw.split(" ", 1)
            if len(parts) < 2 or not parts[1].strip():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入商品代號\n例：/chart WMGS25100246"))
                return
            bond_id = parts[1].strip().upper()
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📊 正在產生 {bond_id} 走勢圖，請稍候..."))
            try:
                from eln_chart import generate_eln_chart
                url = generate_eln_chart(bond_id, engine)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📊 {bond_id} 防守線走勢圖\n\n🟢 綠線 = KO價\n🔴 紅線 = KI價\n🔵 藍線 = Strike\n\n{url}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"圖表產生失敗：{str(e)[:200]}"))
            return
        if cmd.startswith("spending"):
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="💳 正在分析你的消費明細，請稍候約30秒..."))
            try:
                from spending_analyzer import get_monthly_spending_report
                report = get_monthly_spending_report(days=31)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=report[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 消費分析失敗：{str(e)[:200]}"))
            return
        if cmd.startswith("nc"):
            nc_parts = text_raw.split(" ")
            if len(nc_parts) < 2 or not nc_parts[1].strip():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入：/nc YYYYMM\n例：/nc 202606\n或：/nc 202606 小美"))
                return
            qm_nc = nc_parts[1].strip().replace("/", "").replace("-", "")
            if len(qm_nc) != 6 or not qm_nc.isdigit():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="格式錯誤，請輸入6位數字\n例：/nc 202606"))
                return
            yr_nc = qm_nc[:4]
            mo_nc = qm_nc[4:]
            name_filter_nc = nc_parts[2].strip() if len(nc_parts) > 2 else ""
            search_str_nc = f"{yr_nc}-{mo_nc}"
            with engine.begin() as conn:
                rows = conn.execute(text("SELECT bond_id, agent_name, detail FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC"), {"k": ck}).fetchall()
            if not rows:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無資料。"))
                return
            matched_nc = []
            for bond_id, agent_name, detail in rows:
                m = re.search(r"NC閉鎖期 \(至 (\d{4}-\d{2})-\d{2}\)", detail)
                if m and m.group(1) == search_str_nc:
                    if name_filter_nc:
                        ags = [a.strip() for a in re.split(r"[,，、/]", agent_name or "") if a.strip()]
                        if not any(name_filter_nc in a for a in ags):
                            continue
                    matched_nc.append((bond_id, agent_name or "-", bond_status_tag(detail)))
            if not matched_nc:
                tip = f"「{name_filter_nc}」" if name_filter_nc else ""
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到 {yr_nc}/{mo_nc} 閉鎖期打開{tip}的商品。"))
                return
            tip2 = f"（{name_filter_nc}）" if name_filter_nc else ""
            lines_nc = [f"🔓 {yr_nc}/{mo_nc} 閉鎖期打開{tip2}（共 {len(matched_nc)} 筆）:\n"]
            lines_nc += [f"   • {bid} [{ag}]{tag}" for bid, ag, tag in matched_nc]
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines_nc)[:4900]))
            return
        if cmd.startswith("end"):
            parts = text_raw.split(" ", 1)
            if len(parts) < 2 or not parts[1].strip():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入：/end YYYYMM\n例：/end 202604"))
                return
            query_month = parts[1].strip().replace("/", "").replace("-", "")
            if len(query_month) != 6 or not query_month.isdigit():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="格式錯誤，請輸入6位數字\n例：/end 202604"))
                return
            year = query_month[:4]
            month = query_month[4:]
            with engine.begin() as conn:
                rows = conn.execute(text("SELECT bond_id, agent_name, detail FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC"), {"k": ck}).fetchall()
            if not rows:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無資料。"))
                return
            search_str = f"{year}-{month}"
            matched = [(bond_id, agent_name or "-", bond_status_tag(detail)) for bond_id, agent_name, detail in rows if f"最終評價日: {search_str}" in detail]
            if not matched:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到 {year}/{month} 到期的商品。"))
                return
            lines = [f"📅 {year}/{month} 到期商品（共 {len(matched)} 筆）：\n"] + [f"   • {bond_id} [{agent_name}]{tag}" for bond_id, agent_name, tag in matched]
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900]))
            return
        if cmd.startswith("detail"):
            parts = text_raw.split(" ", 1)
            if len(parts) < 2 or not parts[1].strip():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入：/detail 商品代號（例：/detail U123）"))
                return
            matched_id, detail, candidates = db_find_detail(ck, parts[1].strip())
            if detail:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=detail[:4900]))
                return
            if candidates and matched_id is None:
                sample = "\n".join([f"• {c}" for c in candidates[:20]])
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"請再精準一點，候選代號如下：\n{sample}"[:4900]))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="查不到該代號或目前沒有已保存結果。請先 /calc 上傳 Excel。"))
            return
        if cmd.startswith("mail"):
            parts = text_raw.split(" ", 1)
            sub = parts[1].strip().lower() if len(parts) > 1 else ""
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📧 正在讀取郵件並分析中，請稍候..."))
            try:
                from gmail_manager import daily_email_summary, get_gmail_service, get_unread_emails, analyze_emails, format_line_message
                if sub == "unread":
                    service = get_gmail_service()
                    emails = get_unread_emails(service, max_results=10)
                    if not emails:
                        _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="📧 目前沒有未讀郵件 ✅"))
                    else:
                        analysis = analyze_emails(emails)
                        msg = format_line_message(analysis, emails)
                        _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=msg[:4900]))
                else:
                    summary = daily_email_summary()
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=summary[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"郵件讀取失敗：{e}"))
            return
        if cmd.startswith("analysis"):
            parts = text_raw.split(" ", 1)
            arg = parts[1].strip() if len(parts) > 1 else ""
            if not arg:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📊 完整三面向分析指令：\n\n/analysis NVDA — 技術+基本面+消息面\n/analysis 2330 — 台積電完整分析\n/analysis AAPL 3 — 指定月數（預設6個月）"))
                return
            arg_parts = arg.split()
            symbol = arg_parts[0]
            months = min(max(int(arg_parts[1]), 1), 12) if len(arg_parts) > 1 and arg_parts[1].isdigit() else 6
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📊 正在進行 {symbol.upper()} 完整三面向分析，請稍候約20秒..."))
            try:
                from stock_analyzer import full_analysis
                from pdf_generator import upload_to_drive
                img_bytes, summary = full_analysis(symbol, months=months)
                tmp_path = f"/tmp/analysis_{symbol}_{months}.png"
                with open(tmp_path, "wb") as f:
                    f.write(img_bytes)
                link = upload_to_drive(tmp_path, f"{symbol.upper()} Full Analysis {months}M.png")
                os.remove(tmp_path)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📊 {symbol.upper()} 完整分析 (近{months}個月)\n\n{summary}\n\n🔗 圖表：{link}"[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"分析失敗：{str(e)[:200]}"))
            return
        if cmd.startswith("tech"):
            parts = text_raw.split(" ", 1)
            arg = parts[1].strip() if len(parts) > 1 else ""
            if not arg:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📊 技術分析指令：\n\n/tech mag7 — Magnificent Seven 比較分析\n/tech AAPL — 單一股票分析\n/tech AAPL 3 — 指定月數（預設6個月）"))
                return
            arg_parts = arg.split()
            symbol = arg_parts[0]
            months = min(max(int(arg_parts[1]), 1), 12) if len(arg_parts) > 1 and arg_parts[1].isdigit() else 6
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📊 正在分析 {symbol.upper()}，請稍候約15秒..."))
            try:
                from tech_analyzer import analyze_single, analyze_mag7
                from pdf_generator import upload_to_drive
                if symbol.lower() == "mag7":
                    img_bytes, summary = analyze_mag7(months=months)
                    title = f"Magnificent Seven 技術分析 (近{months}個月)"
                else:
                    img_bytes, summary = analyze_single(symbol, months=months)
                    title = f"{symbol.upper()} 技術分析 (近{months}個月)"
                tmp_path = f"/tmp/tech_{symbol}_{months}.png"
                with open(tmp_path, "wb") as f:
                    f.write(img_bytes)
                link = upload_to_drive(tmp_path, f"{title}.png")
                os.remove(tmp_path)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📊 {title}\n\n{summary}\n\n🔗 圖表連結：{link}"[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"技術分析失敗：{str(e)[:200]}"))
            return
        if cmd == "runnow":
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="🔄 手動觸發 ELN 追蹤中，請稍候約30秒..."))
            try:
                from auto_tracking_cron import main as tracking_main
                tracking_main()
                write_job_log("ELN追蹤(手動)", "success", "手動觸發完成")
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="✅ ELN 追蹤執行完成！\n\n打 /tracklog 查看記錄"))
            except Exception as e:
                write_job_log("ELN追蹤(手動)", "error", str(e))
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 執行失敗：{str(e)[:300]}"))
            return
        if cmd in ("bondnav",):
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text="📊 手動觸發債券報價更新中...\n約需 30 分鐘，完成後會通知你 ✅"
            ))
            def _run_bondnav():
                try:
                    import urllib.request as _req
                    pat = os.getenv("GITHUB_PAT", "")
                    if not pat:
                        raise RuntimeError("缺少 GITHUB_PAT 環境變數")
                    data = json.dumps({
                        "ref": "main",
                        "inputs": {"mode": "update"}
                    }).encode("utf-8")
                    req = _req.Request(
                        "https://api.github.com/repos/albert7755-ux/eln-bot/actions/workflows/update_bond_prices.yml/dispatches",
                        data=data,
                        headers={
                            "Authorization": f"Bearer {pat}",
                            "Accept": "application/vnd.github+json",
                            "Content-Type": "application/json",
                            "X-GitHub-Api-Version": "2022-11-28"
                        },
                        method="POST"
                    )
                    with _req.urlopen(req, timeout=15) as resp:
                        status = resp.status
                    user_id = os.getenv("LINE_USER_ID", "")
                    if status == 204 and user_id:
                        line_bot_api.push_message(user_id, TextSendMessage(
                            text="✅ 債券報價更新已觸發！\nGitHub Actions 開始執行，約 30 分鐘後完成。"
                        ))
                    else:
                        raise RuntimeError(f"GitHub API 回應：{status}")
                except Exception as e:
                    user_id = os.getenv("LINE_USER_ID", "")
                    if user_id:
                        line_bot_api.push_message(user_id, TextSendMessage(
                            text=f"❌ /bondnav 觸發失敗：{str(e)[:200]}"
                        ))
            import threading
            threading.Thread(target=_run_bondnav, daemon=True).start()
            return
        if cmd == "fontcheck":
            # 診斷:伺服器上到底有沒有找到中文字型
            try:
                import os as _os
                from bond_sheet import _cjk_font_path
                base = _os.path.dirname(_os.path.abspath(__file__))
                fdir = _os.path.join(base, "fonts")
                fp = _cjk_font_path()
                lines_f = ["🔤 字型診斷",
                           f"程式目錄：{base}",
                           f"工作目錄：{_os.getcwd()}",
                           f"fonts 資料夾存在：{'是' if _os.path.isdir(fdir) else '否'}"]
                if _os.path.isdir(fdir):
                    fl = _os.listdir(fdir)
                    lines_f.append(f"fonts 內容：{fl if fl else '(空)'}")
                    for f_ in fl:
                        try:
                            sz = _os.path.getsize(_os.path.join(fdir, f_)) / 1024 / 1024
                            lines_f.append(f"　・{f_} ({sz:.1f} MB)")
                        except Exception:
                            pass
                lines_f.append(f"\n最終採用：{fp if fp else '❌ 找不到，圖表會用英文標籤'}")
                if fp:
                    try:
                        from matplotlib import font_manager
                        lines_f.append(f"字型名稱：{font_manager.FontProperties(fname=fp).get_name()}")
                    except Exception as e:
                        lines_f.append(f"matplotlib 讀取失敗：{str(e)[:120]}")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines_f)[:4900]))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 診斷失敗：{str(e)[:300]}"))
            return
        if cmd == "econ":
            # /econ check → 立即檢查一次(忽略已推播記錄,強制重查今天)
            # /econ list  → 看目前追蹤清單
            parts_e = raw_cmd.split()
            act_e = parts_e[1].lower() if len(parts_e) > 1 else ""
            if act_e in ("calendar", "日曆"):
                force = len(parts_e) > 2 and parts_e[2].lower() in ("refresh", "更新")
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="📅 查詢中..." if force else "📅 讀取本月日曆..."))
                def _run_cal(chat_id, bot_api_ref, force_):
                    try:
                        from econ_watch import (ensure_table, _month_key, has_month_calendar,
                                                get_month_calendar, fetch_month_calendar,
                                                save_month_calendar, ECON_ITEMS)
                        ensure_table(engine, sql_text)
                        mk = _month_key()
                        if force_ or not has_month_calendar(engine, sql_text, mk):
                            cal = fetch_month_calendar(claude_client, mk)
                            if cal:
                                save_month_calendar(engine, sql_text, mk, cal)
                        else:
                            cal = get_month_calendar(engine, sql_text, mk)
                        label_map = {it["key"]: it["label"] for it in ECON_ITEMS}
                        if not cal:
                            bot_api_ref.push_message(chat_id, TextSendMessage(text=f"📅 {mk} 日曆查詢無結果，可能是本月無相關事件或查詢失敗。"))
                            return
                        lines_c = [f"📅 {mk} 追蹤事件日曆", ""]
                        for k, d in sorted(cal.items(), key=lambda x: x[1]):
                            lines_c.append(f"{d:%m/%d}({'一二三四五六日'[d.weekday()]}) {label_map.get(k, k)}")
                        lines_c.append("\n偵測到公布會自動推播；/econ calendar refresh 可強制重新查詢")
                        bot_api_ref.push_message(chat_id, TextSendMessage(text="\n".join(lines_c)))
                    except Exception as e:
                        print(f"[EconWatch calendar ERROR] {e}")
                        bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 查詢失敗:{str(e)[:200]}"))
                import threading
                threading.Thread(target=_run_cal, args=(ck.split(":", 1)[1], _bot_api, force), daemon=True).start()
                return
            if act_e in ("list", "清單"):
                try:
                    from econ_watch import ECON_ITEMS
                    body_e = ["📊 經濟數據/央行會議追蹤清單", ""]
                    body_e += [f"・{it['label']}" for it in ECON_ITEMS]
                    body_e.append("\n偵測到公布會自動推播；也可打 /econ check 立即檢查")
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(body_e)))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 讀取失敗:{str(e)[:150]}"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📊 檢查中，約30~60秒..."))
            def _run_econ(chat_id, bot_api_ref):
                try:
                    from econ_watch import check_econ_events, already_pushed, ensure_table, _today_key
                    ensure_table(engine, sql_text)
                    today_k = _today_key()
                    # /econ check 強制重查:清掉今天的已推播記錄(僅此次操作觸發的重查)
                    with engine.begin() as conn:
                        conn.execute(sql_text("DELETE FROM econ_event_seen WHERE day_key=:d"), {"d": today_k})
                    hit = check_econ_events(engine, sql_text, claude_client,
                                            lambda m: bot_api_ref.push_message(chat_id, TextSendMessage(text=m)))
                    if not hit:
                        bot_api_ref.push_message(chat_id, TextSendMessage(text="📊 目前沒有偵測到剛公布的數據或會議結果。"))
                except Exception as e:
                    print(f"[EconWatch manual ERROR] {e}")
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 檢查失敗:{str(e)[:200]}"))
            import threading
            threading.Thread(target=_run_econ, args=(ck.split(":", 1)[1], _bot_api), daemon=True).start()
            return
        if cmd == "cleanup":
            # /cleanup        → 預覽:列出 30 天前的舊檔,不刪除
            # /cleanup do     → 實際清理(移到垃圾桶,可救回)
            # /cleanup do 60  → 只清 60 天前的
            parts = raw_cmd.split()
            do_it = len(parts) > 1 and parts[1].lower() in ("do", "yes", "確認", "執行")
            try:
                days = int(parts[2]) if len(parts) > 2 else int(os.getenv("DRIVE_KEEP_DAYS", "30"))
            except ValueError:
                days = 30
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"🧹 {'清理中' if do_it else '預覽中'}（{days} 天前的舊報告）..."))
            def _run_cleanup(chat_id, bot_api_ref):
                try:
                    from pdf_generator import cleanup_drive_folder
                    deleted, checked, detail = cleanup_drive_folder(days=days, dry_run=not do_it)
                    if not checked:
                        bot_api_ref.push_message(chat_id, TextSendMessage(text="找不到「龍蝦報告」資料夾，或裡面沒有檔案。"))
                        return
                    head = (f"🧹 已將 {deleted} 個舊報告移到垃圾桶" if do_it
                            else f"🔍 預覽：有 {deleted} 個檔案超過 {days} 天")
                    lines_c = [head, f"（資料夾共 {checked} 個檔案）"]
                    if detail:
                        lines_c.append("")
                        lines_c += [f"・{n}" for n in detail[:15]]
                        if len(detail) > 15:
                            lines_c.append(f"…另有 {len(detail)-15} 個")
                    if not do_it and deleted:
                        lines_c.append(f"\n要實際清理請打：/cleanup do {days}")
                    if do_it:
                        lines_c.append("\n※ 檔案在 Drive 垃圾桶中，30 天內可救回")
                    push_long_message(bot_api_ref, chat_id, "\n".join(lines_c))
                except Exception as e:
                    print(f"[DriveCleanup ERROR] {e}")
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 清理失敗：{str(e)[:200]}"))
            import threading
            threading.Thread(target=_run_cleanup, args=(ck.split(":", 1)[1], _bot_api), daemon=True).start()
            return
        if cmd in ("myid", "我的id"):
            _uid = getattr(event.source, "user_id", "") or ""
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"你的 LINE ID：\n{_uid}\n\n（如需開通產文件功能，請將此 ID 提供給固定收益科）"))
            return
        if cmd == "sheetuser":
            # 僅 Albert 可管理名單
            if not is_albert:
                return
            parts_u = raw_cmd.split()
            act = parts_u[1].lower() if len(parts_u) > 1 else "list"
            cur = get_doc_users()
            if act in ("add", "加入") and len(parts_u) > 2:
                for u in parts_u[2:]:
                    if u.startswith("U") and len(u) >= 20 and u not in cur:
                        cur.append(u)
                set_doc_users(cur)
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"✅ 已加入，目前 {len(cur)} 人可使用 /sheet"))
                return
            if act in ("del", "remove", "移除") and len(parts_u) > 2:
                cur = [u for u in cur if u not in parts_u[2:]]
                set_doc_users(cur)
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"✅ 已移除，目前 {len(cur)} 人可使用 /sheet"))
                return
            body_u = [f"📋 產文件指令名單（{len(cur)} 人）", "可使用 /sheet（/focus 仍限固定收益科）", ""]
            body_u += [f"・{u}" for u in cur] or ["（目前只有你自己可用）"]
            body_u += ["", "新增：/sheetuser add U1a2b3...",
                       "移除：/sheetuser del U1a2b3...",
                       "（請對方打 /myid 取得自己的 ID）"]
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(body_u)[:4900]))
            return
        if cmd in ("stock", "個股", "分析"):
            # /stock intc → 產生完整個股分析報告(PDF)+LINE摘要
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="📈 用法：/stock intc\n/stock nvda\n"
                         "產出個股完整分析報告(商業模式、財務健康度、護城河評分、"
                         "估值分析、成長潛力、多空辯論、投資評估)，文字摘要+PDF"))
                return
            ticker_in = re.sub(r"[^A-Za-z0-9\.\-]", "", kw.split()[0]).upper()
            if not ticker_in:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請提供有效的股票代碼，例如 /stock intc"))
                return
            if not (is_albert or can_use_doc_cmd(getattr(event.source, "user_id", ""))):
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="/stock 目前開放給投資輔銷同仁使用，請洽轄區投資輔銷協助產出。"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"📈 分析 {ticker_in} 中(抓取財務數據、搜尋公開資訊)，約 60~120 秒..."))

            def _run_stock(chat_id, bot_api_ref, ticker_):
                try:
                    from stock_analysis import (get_5y_financials, get_peer_pe, simple_dcf_scenarios,
                                                generate_analysis, build_financial_chart,
                                                build_report_pdf, build_summary_text)
                    fin = get_5y_financials(ticker_)
                    if not fin:
                        bot_api_ref.push_message(chat_id, TextSendMessage(
                            text=f"❌ 查無「{ticker_}」的財務資料，請確認股票代碼是否正確（需為 Yahoo Finance 可辨識代碼）。"))
                        return
                    dcf = simple_dcf_scenarios(fin)

                    # 同業:先讓 AI 給建議代碼,再用 yfinance 抓真實 PE(避免臆測數字)
                    peers = []
                    try:
                        peer_prompt = (f"股票代碼「{ticker_}」的主要同業競爭對手有哪些?請給2~3家的股票代碼(美股代碼優先)。"
                                       '只回傳JSON:{"peers":["TICKER1","TICKER2"]}')
                        got, _, _ = llm_json_fallback(peer_prompt, max_tokens=200)
                        peer_tickers = [str(x).upper() for x in (got or {}).get("peers", [])][:3]
                        if peer_tickers:
                            peers = get_peer_pe(peer_tickers)
                    except Exception as e:
                        print(f"[StockAnalysis] peers fail: {e}")

                    analysis = generate_analysis(claude_client, ticker_, fin, peers, dcf)
                    if not analysis:
                        _err = getattr(generate_analysis, "last_error", "") or "未知原因"
                        bot_api_ref.push_message(chat_id, TextSendMessage(
                            text=f"❌ 分析生成失敗，請稍後再試。\n（原因：{_err}）"))
                        return

                    summary_txt = build_summary_text(ticker_, fin, analysis, dcf)
                    push_long_message(bot_api_ref, chat_id, summary_txt)

                    chart_png = None
                    try:
                        chart_png = build_financial_chart(fin)
                    except Exception as e:
                        print(f"[StockAnalysis] chart fail: {e}")

                    today_ = datetime.now(TZ_TAIPEI).date()
                    out_pdf = f"/tmp/個股分析_{ticker_}_{today_:%Y%m%d}.pdf"
                    build_report_pdf(out_pdf, ticker_, fin, peers, dcf, analysis,
                                     chart_png=chart_png, today=today_)
                    if chart_png:
                        try:
                            os.remove(chart_png)
                        except Exception:
                            pass
                    link = upload_to_drive(out_pdf, f"個股分析_{ticker_}_{today_:%Y%m%d}.pdf")
                    try:
                        os.remove(out_pdf)
                    except Exception:
                        pass
                    bot_api_ref.push_message(chat_id, TextSendMessage(
                        text=f"📄 {ticker_} 完整分析報告 PDF\n🔗 {link}"))
                except Exception as e:
                    print(f"[StockAnalysis ERROR] {e}")
                    print(_traceback.format_exc())
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 分析失敗：{str(e)[:200]}"))

            import threading
            threading.Thread(target=_run_stock, args=(ck.split(":", 1)[1], _bot_api, ticker_in), daemon=True).start()
            return
        if cmd == "focus":
            # /focus 輝達            → 產生「債市每日聚焦 / 富邦好債報」PPTX(直式兩頁)
            # /focus 輝達 26070004   → 指定焦點債券(1~3檔)
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔,請先把 Bond_Pricing Excel 傳給我。"))
                return
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="📰 用法\n/focus 輝達 → 自動挑短天期與長天期各1檔\n/focus 輝達 26070004 26070003 → 指定焦點債券\n"
                         "產出「債市每日聚焦＋富邦好債報」PPTX(直式兩頁,可直接編輯)"))
                return
            f_kws = []
            _tk = kw.split()
            if len(_tk) > 1:
                _c = [x for x in _tk[1:] if re.fullmatch(r"(?:WMBB)?\d{6,10}", x, re.I)
                      or re.fullmatch(r"[A-Z]{2}[A-Z0-9]{6,10}", x, re.I)]
                if _c:
                    f_kws = _c
                    kw = _tk[0]
            try:
                from bond_coupon_alert import search_issuers, find_bonds, first_num
                hits = search_issuers(str(BOND_PRICE_FILE), kw, max_issuers=3)
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 搜尋失敗:{str(e)[:200]}"))
                return
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找不到「{kw}」,可用 /issuer 查名稱。"))
                return
            if len(hits) > 1:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"「{kw}」對到 {len(hits)} 家:{'、'.join(h[0] for h in hits)}\n請用更精確的名稱再打一次"))
                return
            f_iss, f_bl = hits[0]
            f_today = datetime.now(TZ_TAIPEI).date()
            f_live = [b for b in f_bl if not (b.get("maturity") and b["maturity"] < f_today)]
            f_pick = []
            if f_kws:
                for pk in f_kws[:3]:
                    fb = find_bonds(str(BOND_PRICE_FILE), pk, max_hits=1)
                    if fb:
                        f_pick.append(fb[0])
            if not f_pick:
                # 預設:短天期與長天期各一檔(依剩餘年期取最短、最長),讓兩檔形成對比
                _c2 = []
                for b in f_live:
                    _y = first_num(b.get("years"))
                    if _y is None and b.get("maturity"):
                        _y = (b["maturity"] - f_today).days / 365.25
                    if _y is not None and first_num(b.get("ytm")) is not None:
                        _c2.append((_y, b))
                _c2.sort(key=lambda x: x[0])
                if len(_c2) >= 2:
                    f_pick = [_c2[0][1], _c2[-1][1]]
                elif _c2:
                    f_pick = [_c2[0][1]]
            if not f_pick:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"{f_iss} 目前架上沒有可列示的債券。"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"📰 製作 {f_iss} 債市每日聚焦 PPTX 中(查證信評與財報資料),約 60~90 秒..."))

            def _run_focus(chat_id, bot_api_ref, iss_, bonds_, today_, f_live_for_chart=()):
                try:
                    from bond_focus_ppt import build_focus_pptx
                    from pdf_generator import upload_to_drive
                    parent, ticker = get_issuer_ticker(iss_)
                    b_lines = [f"{b.get('code')} {b['name']} 票面{b['coupon']}% YTM {b.get('ytm')} 到期{b['maturity']:%Y/%m/%d}"
                               for b in bonds_]
                    prompt = (
                        "你是銀行固定收益科的研究員,要製作一份「債市每日聚焦」內部教育訓練文件。"
                        f"對象發行機構:{iss_}" + (f"(上市主體:{parent}, 代碼:{ticker})" if ticker else "") + "。\n"
                        "總行架上焦點債券:\n" + "\n".join(b_lines) + "\n\n"
                        "請依你所知的公開資訊填寫下列欄位,並嚴格遵守:\n"
                        "- 只寫你有把握的事實;不確定的數字、日期、評等一律省略或填 --,絕對不要臆測。\n"
                        "- 評等只填確實知道的機構,不知道填 --。\n"
                        "- 營收結構只在確知部門別或地區別金額時才給,否則 values 給空陣列。\n"
                        "- 語氣專業中性,不做投資建議、不預測股價。\n\n"
                        "只回傳 JSON:\n"
                        "{\"issuer_en\": 英文全名, \"intro\": 公司簡介100~140字, \"headline\": 焦點新聞標題20字內,\n"
                        " \"news_bullets\": [3則,每則60~90字,該公司近期與信用/財務/業務相關的重要事實],\n"
                        " \"ops_blocks\": [[小標題,內文60~90字],[小標題,內文60~90字]],\n"
                        " \"revenue_mix\": {\"labels\":[],\"values\":[],\"unit_note\":\"\"},\n"
                        " \"ratings\": {\"moody\":\"\",\"sp\":\"\",\"fitch\":\"\",\"moody_outlook\":\"\",\"sp_outlook\":\"\","
                        "\"fitch_outlook\":\"\",\"moody_date\":\"\",\"sp_date\":\"\",\"fitch_date\":\"\"},\n"
                        " \"agency_comments\": [[\"穆迪\",評析80~120字],[\"標普\",評析80~120字]],\n"
                        " \"bond_tagline\": 焦點債券一句話定位15字內}")
                    got, _src, _errs = llm_json_fallback(prompt, max_tokens=3000)
                    if not got:
                        bot_api_ref.push_message(chat_id, TextSendMessage(text="❌ 資料生成失敗(AI 服務不可用),請稍後再試。"))
                        return
                    # 營收結構備援資料:近五季財報 + 架上債券到期分布
                    quarterly = {}
                    try:
                        if ticker:
                            from bond_sheet import get_quarterly_series
                            quarterly = get_quarterly_series(ticker) or {}
                    except Exception as e:
                        print(f"[BondFocus] quarterly fail: {e}")
                    mat_dist = {}
                    try:
                        from collections import Counter as _C
                        _buckets = [("3年內", 0, 3), ("3–5年", 3, 5), ("5–10年", 5, 10),
                                    ("10–20年", 10, 20), ("20年以上", 20, 999)]
                        cnt = _C()
                        for b in f_live_for_chart:
                            _y = None
                            if b.get("maturity"):
                                _y = (b["maturity"] - today_).days / 365.25
                            if _y is None:
                                continue
                            for lb, lo, hi in _buckets:
                                if lo <= _y < hi:
                                    cnt[lb] += 1
                                    break
                        labels_ = [lb for lb, _, _ in _buckets if cnt.get(lb)]
                        if labels_:
                            mat_dist = {"labels": labels_, "values": [cnt[lb] for lb in labels_]}
                    except Exception as e:
                        print(f"[BondFocus] maturity dist fail: {e}")

                    data = {
                        "quarterly": quarterly, "maturity_dist": mat_dist,
                        "issuer": iss_, "issuer_en": got.get("issuer_en") or "",
                        "intro": got.get("intro") or "", "headline": got.get("headline") or f"{iss_} 焦點速報",
                        "news_bullets": got.get("news_bullets") or [],
                        "ops_blocks": got.get("ops_blocks") or [],
                        "revenue_mix": got.get("revenue_mix") or {},
                        "ratings": got.get("ratings") or {},
                        "agency_comments": got.get("agency_comments") or [],
                        "bond_tagline": got.get("bond_tagline") or "",
                        "date_str": f"{today_:%Y/%m/%d}",
                        "source_note": "（資料來源：發行機構公開財報、信評機構公開資訊、公開新聞）",
                        "bonds": [{"code": b.get("code") or "-", "name": b["name"], "coupon": b["coupon"],
                                   "ytm": str(b.get("ytm") or "-"),
                                   "maturity": f"{b['maturity']:%Y/%m/%d}" if b.get("maturity") else "-"}
                                  for b in bonds_],
                    }
                    base_ = f"債市每日聚焦_{iss_}_{today_:%Y%m%d}"
                    out = f"/tmp/{base_}.pptx"
                    build_focus_pptx(out, data)
                    link = upload_to_drive(out, base_ + ".pptx")
                    # PDF 由本機直接產出(reportlab),不依賴 Google 轉檔
                    pdf_link, pdf_err = None, ""
                    try:
                        from bond_focus_ppt import build_focus_pdf
                        out_pdf = f"/tmp/{base_}.pdf"
                        build_focus_pdf(out_pdf, data)
                        pdf_link = upload_to_drive(out_pdf, base_ + ".pdf")
                        try:
                            os.remove(out_pdf)
                        except Exception:
                            pass
                    except Exception as e:
                        print(f"[BondFocus] PDF 產生失敗: {e}")
                        print(_traceback.format_exc()[:500])
                        pdf_err = f"{type(e).__name__}: {str(e)[:120]}"
                    try:
                        os.remove(out)
                    except Exception:
                        pass
                    msg_f = (f"📰 {iss_} 債市每日聚焦\n\n"
                             f"📊 PPTX（可編輯）\n{link}\n")
                    if pdf_link:
                        msg_f += f"\n📄 PDF（手機預覽用）\n{pdf_link}\n"
                    elif pdf_err:
                        msg_f += f"\n（PDF 產生失敗：{pdf_err}）\n"
                    msg_f += "\n⚠️ 內容由 AI 依公開資訊整理，發布前請人工核對評等、日期與財務數字。"
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=msg_f))
                except Exception as e:
                    print(f"[BondFocus ERROR] {e}")
                    print(_traceback.format_exc())
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 製作失敗:{str(e)[:200]}"))
            import threading
            threading.Thread(target=_run_focus,
                             args=(ck.split(":", 1)[1], _bot_api, f_iss, f_pick, f_today, f_live), daemon=True).start()
            return
        if cmd == "sheet":
            # /sheet 蘋果  → 發行機構參考資訊(LINE文字 + PDF連結)
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔,請先把 Bond_Pricing Excel 傳給我。"))
                return
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📋 用法\n/sheet 蘋果 → 自動挑三檔(各年期區間YTM最高)\n/sheet 威瑞森 25120005 → 指定1檔\n/sheet 威瑞森 25120005 25100003 → 指定2檔(最多3檔)\n產出發行機構參考資訊(簡介+信評+財務+圖表+標的),文字版+PDF"))
                return
            # 可在機構名後面接 1~3 個產品代碼/ISIN,指定要列出的標的
            # 例:/sheet 威瑞森 25120005 25100003
            picked_kws = []
            _tok = kw.split()
            if len(_tok) > 1:
                _cand = [t for t in _tok[1:]
                         if re.fullmatch(r"(?:WMBB)?\d{6,10}", t, re.I)
                         or re.fullmatch(r"[A-Z]{2}[A-Z0-9]{6,10}", t, re.I)]
                if _cand:
                    picked_kws = _cand
                    kw = _tok[0]
            try:
                from bond_coupon_alert import search_issuers
                hits = search_issuers(str(BOND_PRICE_FILE), kw, max_issuers=3)
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 搜尋失敗:{str(e)[:200]}"))
                return
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找不到「{kw}」,可用 /issuer 查名稱。"))
                return
            if len(hits) > 1:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"「{kw}」對到 {len(hits)} 家:{'、'.join(h[0] for h in hits)}\n請用更精確的名稱再打一次 /sheet"))
                return
            iss, bl = hits[0]
            picked = []
            if picked_kws:
                from bond_coupon_alert import find_bonds
                for pk in picked_kws[:3]:
                    fb = find_bonds(str(BOND_PRICE_FILE), pk, max_hits=1)
                    if fb:
                        picked.append(fb[0])
                if not picked:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(
                        text=f"找不到指定的標的：{'、'.join(picked_kws)}\n請確認產品代碼，或不帶代碼讓系統自動挑選。"))
                    return
            _note_p = f"（指定 {len(picked)} 檔標的）" if picked else ""
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"📋 整理 {iss} 參考資訊中{_note_p},約 30~60 秒..."))
            def _run_sheet(chat_id, bot_api_ref, iss_, bl_, picked_=None):
                try:
                    from bond_sheet import (get_financials, build_sheet_text, build_sheet_pdf,
                                            get_quarterly_series, build_charts_png, build_peer_chart)
                    from pdf_generator import upload_to_drive
                    today_ = datetime.now(TZ_TAIPEI).date()
                    parent, ticker = get_issuer_ticker(iss_)
                    intro = get_issuer_profile(iss_, parent or "")
                    parent_note = parent if (parent and parent != iss_) else ""
                    fin = get_financials(ticker) if ticker else None
                    if parent_note and fin:
                        parent_note = f"財報為母集團 {parent_note}"
                    fin_comment = get_fin_comment(iss_, fin)
                    bullets = get_issuer_bullets(iss_, parent or "")
                    peer_png = None
                    try:
                        if fin and fin.get("market_cap"):
                            pc = get_peer_caps(iss_, parent or "")
                            peer_png = build_peer_chart(parent or iss_, fin["market_cap"], pc,
                                                        f"/tmp/peer_{iss_}_{today_:%Y%m%d}.png")
                    except Exception as e:
                        print(f"[BondSheet peer chart] {e}")
                    peers = get_peer_comparison(iss_, parent or "", fin)
                    rating_note = get_rating_outlook(iss_)
                    from bond_sheet import get_ust_curve
                    curve = get_ust_curve()
                    hist = build_issuer_hist_map([b["isin"] for b in bl_])
                    kw_ = dict(fin=fin, parent_note=parent_note, hist_map=hist, fin_comment=fin_comment,
                               peers=peers, rating_note=rating_note, ust_curve=curve, today=today_)
                    # 近五季財報圖表(抓不到就略過,不影響 PDF 產出)
                    charts_png = None
                    chart_note = ""
                    charts_comment = ""
                    if ticker:
                        try:
                            q = get_quarterly_series(ticker)
                            if q:
                                charts_png = build_charts_png(q, f"/tmp/charts_{iss_}_{today_:%Y%m%d}.png")
                                charts_comment = get_charts_comment(iss_, q)
                                if not charts_png:
                                    chart_note = "（圖表繪製失敗，未附圖表）"
                            else:
                                chart_note = f"（{ticker} 季報資料不足，未附圖表）"
                        except Exception as e:
                            print(f"[BondSheet charts] {e}")
                            chart_note = "（圖表產生錯誤，未附圖表）"
                    else:
                        chart_note = "（未對應到上市公司，未附圖表）"
                    txt = build_sheet_text(iss_, intro, bl_, charts_comment=charts_comment,
                                           intro_bullets=bullets, **kw_)
                    push_long_message(bot_api_ref, chat_id, txt)
                    pdf_path = f"/tmp/參考資訊_{iss_}_{today_:%Y%m%d}.pdf"
                    build_sheet_pdf(pdf_path, iss_, intro, bl_, charts_png=charts_png,
                                    charts_comment=charts_comment, intro_bullets=bullets,
                                    peer_png=peer_png, picked_bonds=picked_, **kw_)
                    if peer_png:
                        try:
                            os.remove(peer_png)
                        except Exception:
                            pass
                    if charts_png:
                        try:
                            os.remove(charts_png)
                        except Exception:
                            pass
                    link = upload_to_drive(pdf_path, f"參考資訊_{iss_}_{today_:%Y%m%d}.pdf")
                    try:
                        os.remove(pdf_path)
                    except Exception:
                        pass
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"📎 {iss_} 參考資訊 PDF{chart_note}\n🔗 {link}"))
                    try:
                        from bond_sheet import font_status
                        fs_ = font_status()
                        if "未找到" in fs_:
                            bot_api_ref.push_message(chat_id, TextSendMessage(text=f"ℹ️ {fs_}"))
                    except Exception:
                        pass
                except Exception as e:
                    print(f"[BondSheet ERROR] {e}")
                    print(_traceback.format_exc())
                    bot_api_ref.push_message(chat_id, TextSendMessage(text=f"❌ 參考資訊產生失敗:{str(e)[:200]}"))
            import threading
            threading.Thread(target=_run_sheet, args=(ck.split(":", 1)[1], _bot_api, iss, bl, picked), daemon=True).start()
            return
        if cmd == "rating":
            # /rating              → 立刻掃一次監控名單的信評新聞
            # /rating list         → 監控名單
            # /rating watch 台積電  → 手動加入   /rating unwatch 台積電 → 移除
            # /rating fix          → 補英文名
            arg = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            al = arg.lower()
            chat_id = ck.split(":", 1)[1]
            if al in ("list", "ls", "名單"):
                with engine.begin() as conn:
                    rows = conn.execute(text("SELECT issuer, en_name, added_by, last_checked FROM bond_rating_watch WHERE active=TRUE ORDER BY issuer")).fetchall()
                if not rows:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="監控名單是空的。跑一次 /coupon table 會自動加入，或 /rating watch 蘋果")); return
                lines = [f"📡 信評監控名單（{len(rows)} 家，每天 07:00 掃 Google News）"]
                for iss, en, by, lc in rows:
                    lines.append(f"▪ {iss}（{en or '英文名待補'}）")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900])); return
            if al.startswith(("watch ", "add ", "加入 ")):
                name = arg.split(" ", 1)[1].strip()
                with engine.begin() as conn:
                    conn.execute(text("INSERT INTO bond_rating_watch(issuer, added_by, active) VALUES (:i,:b,TRUE) ON CONFLICT (issuer) DO UPDATE SET active=TRUE"), {"i": name, "b": chat_id})
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 已加入監控：{name}（英文名會自動補）")); return
            if al.startswith(("unwatch ", "del ", "remove ", "移除 ")):
                name = arg.split(" ", 1)[1].strip()
                with engine.begin() as conn:
                    n = conn.execute(text("UPDATE bond_rating_watch SET active=FALSE WHERE issuer=:i"), {"i": name}).rowcount
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 已移除：{name}" if n else f"名單裡沒有「{name}」")); return
            if al in ("fix", "補英文"):
                n = ensure_en_names()
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 已補 {n} 家英文名")); return
            if al in ("", "check", "now", "掃描"):
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📡 掃描監控名單的信評新聞中，約 30～90 秒..."))
                def _run_rating(chat_id_, bot_api_ref):
                    try:
                        msg = run_rating_news_check(days=3)
                        if not msg:
                            st = globals().get("_LAST_RATING_STATS", {})
                            msg = ("📡 近 3 天監控名單沒有新的信評動作新聞。\n\n"
                                   f"掃描 {st.get('issuers', 0)} 家發行機構\n"
                                   f"搜到相關新聞 {st.get('raw', 0)} 則\n"
                                   f"扣除已推播過 → 新新聞 {st.get('new', 0)} 則\n"
                                   f"AI 判定為真正評等動作 {st.get('kept', 0)} 則\n\n"
                                   "／rating list 可看監控名單")
                        push_long_message(bot_api_ref, chat_id_, msg)
                    except Exception as e:
                        bot_api_ref.push_message(chat_id_, TextSendMessage(text=f"❌ 信評掃描失敗：{str(e)[:200]}"))
                import threading
                threading.Thread(target=_run_rating, args=(chat_id, _bot_api), daemon=True).start()
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="用法：/rating（立即掃描）\n/rating list\n/rating watch 蘋果\n/rating unwatch 蘋果")); return
        if cmd in ("price", "p", "價格", "報價"):
            # /price 蘋果 2043   → 單檔完整報價（模糊搜尋，同 /bondalert 的找法）
            # /price US037833EN
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if kw.lower() in ("settarget", "設定查價"):
                # 在群組裡打 /price settarget → 該群開放查價（僅 /price，不含每日推播）
                if event.source.type not in ("group", "room"):
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請在要開放查價的群組裡打 /price settarget。"))
                    return
                gid = event.source.group_id if event.source.type == "group" else event.source.room_id
                targets = load_targets()
                groups = targets.get("bond_query_groups", [])
                if gid not in groups:
                    groups.append(gid)
                targets["bond_query_groups"] = groups
                save_targets(targets)
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="✅ 本群已開放海外債查價。\n群裡任何人打 /price 都能查，其他指令與訊息我不會回應，也不會有每日推播。\n取消請打 /price settarget off\n\n" + BOND_QUERY_HELP))
                return
            if kw.lower() in ("settarget off", "off"):
                if event.source.type in ("group", "room"):
                    gid = event.source.group_id if event.source.type == "group" else event.source.room_id
                    targets = load_targets()
                    targets["bond_query_groups"] = [g for g in targets.get("bond_query_groups", []) if g != gid]
                    save_targets(targets)
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="✅ 本群已關閉查價功能。"))
                return
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔，請先把 Bond_Pricing Excel 傳給我。"))
                return
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="🔎 用法：/price 蘋果 2043\n/price 26070003（產品代碼免打WMBB）\n/price 26070003 30 → 近30天價格走勢\n（名稱片段＋到期年份，或 ISIN／產品代碼）"))
                return
            # 結尾帶天數 → 查歷史走勢（例:/price 26070003 30）;年份(20xx)不算
            hist_days = 0
            _kwp = kw.split()
            if len(_kwp) >= 2 and re.fullmatch(r"\d{1,3}", _kwp[-1]) and not re.fullmatch(r"20\d\d", _kwp[-1]):
                hist_days = max(2, min(int(_kwp[-1]), 365))
                kw = " ".join(_kwp[:-1]).strip()
            try:
                from bond_coupon_alert import find_bonds, first_num, pi_tag
                hits = find_bonds(str(BOND_PRICE_FILE), kw, max_hits=5)
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 查詢失敗：{str(e)[:200]}"))
                return
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找不到「{kw}」，可用 /issuer 查發行機構名稱。"))
                return
            mtime = datetime.fromtimestamp(BOND_PRICE_FILE.stat().st_mtime, TZ_TAIPEI).strftime("%m/%d %H:%M")
            if len(hits) > 1:
                lines = [f"🔎 「{kw}」對到 {len(hits)} 檔（報價檔 {mtime}）："]
                for b in hits:
                    offer = b["offer"] if b["offer"] not in (None, "", 0, "#VALUE!", "#N/A") else "-"
                    ytm = b["ytm"] if b["ytm"] not in (None, "", 0, "#N/A") else "-"
                    mat = f"{b['maturity']:%Y/%m/%d}" if b["maturity"] else "-"
                    lines.append(f"▪ {b['name']} {b['ccy']}｜Offer {offer}｜YTM {ytm}｜到期 {mat}")
                lines.append("\n加到期年份可看單檔完整資訊，例：/price 蘋果 2043")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900]))
                return
            b = hits[0]
            if hist_days:
                try:
                    with engine.begin() as conn:
                        rows = conn.execute(text("""SELECT snap_date, offer, bid, ytm FROM bond_price_history
                                                   WHERE isin=:i AND snap_date >= CURRENT_DATE - :d * INTERVAL '1 day'
                                                   ORDER BY snap_date"""), {"i": b["isin"], "d": hist_days}).fetchall()
                except Exception as e:
                    rows = []
                    print(f"[BondPrice hist] {e}")
                if len(rows) < 2:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(
                        text=f"📉 {b['name']}\n近 {hist_days} 天的報價歷史不足（目前只有 {len(rows)} 筆）。\n"
                             f"每天上傳報價檔會累積，也可以把過去的報價檔丟給我補歷史。"))
                    return
                offs = [(d, o, bd, y) for d, o, bd, y in rows if o is not None]
                lines_h = [f"📉 {b['name']}｜近 {hist_days} 天報價走勢",
                           f"{b.get('code') or '-'}｜{b['ccy']} {b['coupon']}% {b['freq']}", ""]
                step = max(1, len(offs) // 12)   # 太多筆時等距抽樣,避免訊息過長
                for d, o, bd, y in offs[::step]:
                    lines_h.append(f"{d:%m/%d}  Offer {o:g}"
                                  + (f"｜Bid {bd:g}" if bd is not None else "")
                                  + (f"｜YTM {y:g}" if y is not None else ""))
                if offs and offs[-1] not in offs[::step]:
                    d, o, bd, y = offs[-1]
                    lines_h.append(f"{d:%m/%d}  Offer {o:g}"
                                  + (f"｜Bid {bd:g}" if bd is not None else "")
                                  + (f"｜YTM {y:g}" if y is not None else ""))
                if len(offs) >= 2:
                    o_first, o_last = offs[0][1], offs[-1][1]
                    hi = max(offs, key=lambda x: x[1]); lo = min(offs, key=lambda x: x[1])
                    chg = (o_last - o_first) / o_first * 100 if o_first else 0
                    lines_h += ["", f"期間變化：{o_first:g} → {o_last:g}（{chg:+.1f}%）",
                                f"最高 {hi[1]:g}（{hi[0]:%m/%d}）｜最低 {lo[1]:g}（{lo[0]:%m/%d}）",
                                f"共 {len(offs)} 個報價日"]
                lines_h.append("\n※ 依報價檔 Offer 價，非市場成交價")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines_h)[:4900]))
                return
            offer = b["offer"] if b["offer"] not in (None, "", 0, "#VALUE!", "#N/A") else "-"
            ytm = b["ytm"] if b["ytm"] not in (None, "", 0, "#N/A") else "-"
            mat = f"{b['maturity']:%Y/%m/%d}" if b["maturity"] else "-"
            rt = " / ".join(x for x in str(b.get("ratings") or "").split(" / ") if x and x.upper() not in ("N/A", "NA", "NONE"))
            from bond_coupon_alert import first_num as _fn
            _bid_v, _off_v = _fn(b.get("bid")), _fn(b.get("offer"))
            bid = b["bid"] if b["bid"] not in (None, "", 0, "#VALUE!", "#N/A") else "-"
            _spread = (f"（價差 {round(_off_v - _bid_v, 2):g}）"
                       if (_bid_v and _off_v and _off_v >= _bid_v) else "")
            lines = [f"💵 {b['name']}",
                     f"ISIN {b['isin']}｜代碼 {b['code'] or '-'}",
                     f"{b['ccy']}｜票面 {b['coupon']}%｜{b['freq']}配息｜{pi_tag(b)}",
                     f"Offer {offer}（買進）｜Bid {bid}（賣回）{_spread}",
                     f"YTM/YTC {ytm}",
                     f"到期 {mat}" + (f"｜評等 {rt}" if rt else ""),
                     f"順位 {b.get('seniority') or '-'}｜最低申購 {b.get('min_amt') or '-'}｜本日額度 {b.get('avail') or '-'}"]
            if str(b.get("remark") or "").strip():
                lines.append(f"備註：{str(b['remark']).strip()[:150]}")
            # 歷史走勢（bond_price_history 最近 5 筆）
            try:
                with engine.begin() as conn:
                    rows = conn.execute(text("""SELECT snap_date, offer, ytm FROM bond_price_history
                                               WHERE isin=:i ORDER BY snap_date DESC LIMIT 5"""), {"i": b["isin"]}).fetchall()
                if len(rows) >= 2:
                    hist = "｜".join(f"{d:%m/%d} {o:g}" for d, o, y in reversed(rows) if o is not None)
                    lines.append(f"近期Offer：{hist}")
                    o_new, o_old = rows[0][1], rows[-1][1]
                    if o_new and o_old:
                        lines.append(f"（{rows[-1][0]:%m/%d} 以來 {((o_new-o_old)/o_old*100):+.1f}%）")
            except Exception as e:
                print(f"[BondPrice hist] {e}")
            lines.append(f"\n📎 報價檔 {mtime}｜到價追蹤：/bondalert {kw} ytm>{ytm if isinstance(ytm,(int,float)) else 5}")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(str(x) for x in lines)[:4900]))
            return
        if cmd in ("bid", "賣回", "贖回"):
            # /bid 26070003 → 只看賣回(Bid)價,客戶臨時問「現在賣大概多少」時用
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔。"))
                return
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text="🔎 用法：/bid 26070003\n/bid 蘋果 2043\n（查賣回價，用法與 /price 相同）"))
                return
            try:
                from bond_coupon_alert import find_bonds, first_num as _fn2, pi_tag as _pt
                hits = find_bonds(str(BOND_PRICE_FILE), kw, max_hits=5)
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 查詢失敗：{str(e)[:200]}"))
                return
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找不到「{kw}」，可用 /issuer 查名稱。"))
                return
            mtime = datetime.fromtimestamp(BOND_PRICE_FILE.stat().st_mtime, TZ_TAIPEI).strftime("%m/%d %H:%M")
            if len(hits) > 1:
                lines_b = [f"🔎 「{kw}」對到 {len(hits)} 檔（報價檔 {mtime}）："]
                for x in hits:
                    _b, _o = _fn2(x.get("bid")), _fn2(x.get("offer"))
                    _mat = f"｜到期 {x['maturity']:%Y/%m}" if x.get("maturity") else ""
                    lines_b.append(f"▪ {x.get('code') or '-'} {x['name']}\n"
                                   f"  Bid {_b if _b else '-'}｜Offer {_o if _o else '-'}{_mat}")
                lines_b.append("\n加到期年份或用完整代碼可鎖定單一檔")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines_b)[:4900]))
                return
            x = hits[0]
            _b, _o = _fn2(x.get("bid")), _fn2(x.get("offer"))
            if not _b:
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"💰 {x['name']}\n{x.get('code') or '-'}\n\n"
                         "報價檔目前沒有這一檔的賣回價（Bid），請洽固定收益科確認。"))
                return
            sp = f"｜價差 {round(_o - _b, 2):g}" if (_o and _o >= _b) else ""
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"💰 {x['name']}\n{x.get('code') or '-'}｜{x['ccy']} {x['coupon']}% {x['freq']}\n\n"
                     f"賣回價（Bid）{_b:g}\n"
                     f"買進價（Offer）{(f'{_o:g}' if _o else '-')}{sp}\n"
                     f"到期 {x['maturity']:%Y/%m/%d}｜{_pt(x)}\n\n"
                     f"📎 報價檔 {mtime}\n"
                     "※ 為報價檔參考價，實際賣回金額與可否承作以總行系統為準，"
                     "另需計入應計利息與相關費用。"))
            return
        if cmd in ("move", "movers", "異動"):
            # /move        → vs 上一份報價，變動 ≥ 1%
            # /move 7      → vs 7 天前，≥ 2%
            # /move 7 3    → vs 7 天前，≥ 3%
            # /move 30 5   → vs 30 天前，≥ 5%
            parts = raw_cmd.split()
            try:
                db_ = int(parts[1]) if len(parts) > 1 else 1
                th_ = float(parts[2]) if len(parts) > 2 else (1.0 if db_ <= 1 else 2.0)
            except ValueError:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="用法：/move（vs 上一份，≥1%）\n/move 7（vs 7天前，≥2%）\n/move 7 3（vs 7天前，≥3%）")); return
            try:
                mv, base, latest = price_movers(days_back=db_, threshold_pct=th_)
                if not latest:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="還沒有報價歷史，上傳一次報價檔後就會開始累積（每天上傳一次，隔天起可比對）。")); return
                if not base:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"目前只有 {latest:%m/%d} 一份報價歷史，明天上傳新檔後就能比對。")); return
                if not mv:
                    mv = f"📊 {base:%m/%d} → {latest:%m/%d} 沒有 Offer 變動 ≥ {th_:g}% 的債券。\n（可調低門檻，例：/move {db_} 0.5）"
                chunks_m = [mv[i:i+4900] for i in range(0, len(mv), 4900)]
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=chunks_m[0]))
                for c in chunks_m[1:]:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=c))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 異動查詢失敗：{str(e)[:200]}"))
            return
        if cmd == "bondalert":
            # /bondalert 蘋果 2043 ytm>5.2   → 設定到價條件（YTM 或 Offer）
            # /bondalert US037833 offer<90
            # /bondalert list  /bondalert del 3
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔，請先把 Bond_Pricing Excel 傳給我。"))
                return
            arg = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            chat_id = ck.split(":", 1)[1]
            usage = ("🎯 到價通知用法：\n/bondalert 蘋果 2043 ytm>5.2\n/bondalert US037833 offer<90\n/bondalert 蘋果9 ytm>=5\n"
                     "/bondalert list → 我的條件\n/bondalert del 3 → 刪除 #3\n（每次更新報價檔＋每天 06:45 檢查，命中推播一次後自動關閉）")
            if not arg:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=usage)); return
            al = arg.lower()
            if al in ("list", "ls", "清單"):
                with engine.begin() as conn:
                    rows = conn.execute(text("SELECT id, bond_name, field, op, threshold, last_value FROM bond_price_alert WHERE chat_id=:c AND active=TRUE ORDER BY id"), {"c": chat_id}).fetchall()
                if not rows:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有進行中的到價條件。\n" + usage)); return
                lines = ["🎯 進行中的到價條件"]
                for r in rows:
                    lv = f"｜目前 {r[5]}" if r[5] is not None else ""
                    lines.append(f"#{r[0]} {r[1]}｜{r[2].upper()} {r[3]} {r[4]}{lv}")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900])); return
            if al.startswith(("del ", "delete ", "刪除 ", "rm ")):
                try:
                    aid = int(al.split()[1])
                except Exception:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請給編號，例：/bondalert del 3")); return
                with engine.begin() as conn:
                    n = conn.execute(text("UPDATE bond_price_alert SET active=FALSE WHERE id=:i AND chat_id=:c AND active=TRUE"), {"i": aid, "c": chat_id}).rowcount
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 已刪除 #{aid}" if n else f"找不到 #{aid}")); return
            m = re.search(r"(ytm|offer|price|殖利率|價格|yield)\s*(>=|<=|>|<|≥|≤)\s*(-?\d+(?:\.\d+)?)\s*$", arg, flags=re.I)
            if not m:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="看不懂條件 🤔\n" + usage)); return
            field_raw, op, thr = m.group(1).lower(), m.group(2).replace("≥", ">=").replace("≤", "<="), float(m.group(3))
            field = "ytm" if field_raw in ("ytm", "殖利率", "yield") else "offer"
            keyword = arg[:m.start()].strip()
            if not keyword:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請指定債券，例：/bondalert 蘋果 2043 ytm>5.2")); return
            from bond_coupon_alert import find_bonds, first_num
            hits = find_bonds(str(BOND_PRICE_FILE), keyword)
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到「{keyword}」，可用 /issuer 查名稱或改用 ISIN。")); return
            if len(hits) > 1:
                lines = [f"「{keyword}」對到 {len(hits)} 檔，請加到期年份或用 ISIN 指定："]
                for b in hits:
                    lines.append(f"▪ {b['name']}｜到期 {b['maturity']:%Y/%m/%d}｜{b['isin']}")
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900])); return
            b = hits[0]
            cur = first_num(b["ytm"]) if field == "ytm" else first_num(b["offer"])
            with engine.begin() as conn:
                aid = conn.execute(text("""INSERT INTO bond_price_alert(chat_id, isin, bond_name, field, op, threshold, last_value)
                                          VALUES (:c,:i,:n,:f,:o,:t,:v) RETURNING id"""),
                                   {"c": chat_id, "i": b["isin"], "n": b["name"], "f": field, "o": op, "t": thr, "v": cur}).scalar()
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text=f"✅ 已設定到價通知 #{aid}\n{b['name']}（{b['isin']}）\n條件：{field.upper()} {op} {thr}\n目前：Offer {b['offer']}｜YTM {b['ytm']}\n"
                     f"每次更新報價檔＋每天 06:45 檢查，命中會推播一次。/bondalert list 查看"))
            return
        if cmd == "issuer" or cmd.startswith("issuer "):
            # /issuer 蘋果   → 該發行機構簡介 + 架上所有債券（模糊搜尋，不限配息中）
            kw = raw_cmd.split(" ", 1)[1].strip() if " " in raw_cmd else ""
            if not _BOND_RADAR_OK or not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔，請先把 Bond_Pricing Excel 傳給我。"))
                return
            if not kw:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="🔎 用法：/issuer 蘋果\n/issuer merrill\n/issuer US037833\n（模糊搜尋發行機構／債券名稱／ISIN／產品代碼）"))
                return
            try:
                from bond_coupon_alert import search_issuers, format_issuer_bonds
                hits = search_issuers(str(BOND_PRICE_FILE), kw)
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 搜尋失敗：{str(e)[:200]}"))
                return
            if not hits:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找不到跟「{kw}」相關的發行機構。"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🔎 找到 {len(hits)} 家：{'、'.join(h[0] for h in hits)}\n整理簡介中..."))
            def _run_issuer(chat_key, bot_api_ref, hits_):
                try:
                    intros = get_issuer_intros([h[0] for h in hits_])
                except Exception as e:
                    intros = {h[0]: f"（簡介暫無法取得：{str(e)[:80]}）" for h in hits_}
                for iss, bl in hits_:
                    push_long_message(bot_api_ref, chat_key.split(":", 1)[1], format_issuer_bonds(iss, bl, intros.get(iss, "")))
            import threading
            threading.Thread(target=_run_issuer, args=(ck, _bot_api, hits), daemon=True).start()
            return
        if cmd == "coupon":
            # /coupon          → 最晚下單日在 3 個營業日內
            # /coupon 7        → 7 個營業日內
            # /coupon all      → 未來14天全部
            # /coupon table    → 產出 Excel 條件表 + 發行機構簡介（Google Drive 連結）
            if not _BOND_RADAR_OK:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="❌ 配息雷達模組未載入，請確認 bond_coupon_alert.py 已部署。"))
                return
            if not BOND_PRICE_FILE.exists():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📭 還沒有海外債報價檔。請直接把總行的 Bond_Pricing Excel 傳給我，我會自動存檔並產生配息雷達。"))
                return
            parts = raw_cmd.split()
            arg = parts[1].lower() if len(parts) > 1 else ""
            _today = datetime.now(TZ_TAIPEI).date()
            if arg in ("settarget", "設定推播"):
                # 在海外債群組裡打 /coupon settarget → 每天 06:45 配息雷達推到這個群
                targets = load_targets()
                if event.source.type == "group":
                    tid, ttype = event.source.group_id, "group"
                elif event.source.type == "room":
                    tid, ttype = event.source.room_id, "room"
                else:
                    tid, ttype = event.source.user_id, "user"
                targets["bond"] = tid
                targets["bond_type"] = ttype
                save_targets(targets)
                _bot_api.reply_message(event.reply_token, TextSendMessage(
                    text=f"✅ 已設定：每天 06:45 海外債配息雷達推播到這個{'群組' if ttype != 'user' else '對話'}。\n取消請打 /coupon settarget off"))
                return
            if arg in ("off",) or (len(parts) > 2 and parts[1].lower() == "settarget" and parts[2].lower() == "off"):
                targets = load_targets()
                targets.pop("bond", None); targets.pop("bond_type", None)
                save_targets(targets)
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="✅ 已取消群組推播，之後只推給 Albert 個人。"))
                return
            if arg in ("subscribe", "sub", "訂閱"):
                # 理專在一對一私訊裡打 /coupon subscribe → 每天 06:45 個別收到配息雷達
                if event.source.type != "user":
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="請私訊 Albert Claw 打 /coupon subscribe 訂閱個人推播；群組請用 /coupon settarget。"))
                    return
                targets = load_targets()
                subs = targets.get("bond_subscribers", [])
                if event.source.user_id in subs:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="你已經訂閱過了 ✅ 每天 06:45 會收到海外債配息雷達。\n取消請打 /coupon unsubscribe"))
                    return
                subs.append(event.source.user_id)
                targets["bond_subscribers"] = subs
                save_targets(targets)
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 訂閱成功！每天 06:45（週一～五）會收到海外債配息雷達。\n目前訂閱人數：{len(subs)}\n取消請打 /coupon unsubscribe"))
                return
            if arg in ("unsubscribe", "unsub", "取消訂閱"):
                targets = load_targets()
                subs = [u for u in targets.get("bond_subscribers", []) if u != getattr(event.source, "user_id", None)]
                targets["bond_subscribers"] = subs
                save_targets(targets)
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="✅ 已取消訂閱，之後不會再收到每日配息雷達。"))
                return
            if arg in ("subscribers", "subs", "訂閱名單"):
                targets = load_targets()
                subs = targets.get("bond_subscribers", [])
                grp = "有" if targets.get("bond") else "無"
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📋 海外債配息雷達推播對象\n個人訂閱：{len(subs)} 人\n群組推播：{grp}\n（LINE 不提供姓名，只有 ID；理專自己打 /coupon subscribe 就會加入）"))
                return
            if arg in ("table", "表", "excel", "xlsx"):
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📊 正在整理配息債條件表＋發行機構簡介，約 30～60 秒，完成後傳連結給你..."))
                def _run_coupon_table(chat_key, bot_api_ref, today_):
                    try:
                        from bond_coupon_alert import build_coupon_sheet
                        from pdf_generator import upload_to_drive
                        out_path = f"/tmp/配息債條件表_{today_:%Y%m%d}.xlsx"
                        _, n_bond, n_iss, intros = build_coupon_sheet(str(BOND_PRICE_FILE), out_path, today=today_, lookahead=14, intro_fn=get_issuer_intros)
                        link = upload_to_drive(out_path, f"配息債條件表_{today_:%Y%m%d}.xlsx")
                        try:
                            os.remove(out_path)
                        except Exception:
                            pass
                        bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(
                            text=f"📎 配息債條件表已完成\n{today_:%m/%d} 起未來14天還來得及參與：{n_bond} 檔，發行機構 {n_iss} 家\n"
                                 f"（單一工作表，含評等/順位/Offer/YTM/備註/發行機構簡介）\n\n🔗 {link}"))
                        # 手機看 Excel 不方便 → 簡介同步用文字推一份
                        if intros:
                            intro_txt = "🏦 發行機構簡介（AI 產生，對客請以公開資訊為準）\n\n" + "\n\n".join(
                                f"▪ {k}\n{v}" for k, v in intros.items())
                            push_long_message(bot_api_ref, chat_key.split(":", 1)[1], intro_txt)
                    except Exception as e:
                        print(f"[BondRadar TABLE ERROR] {e}")
                        print(_traceback.format_exc())
                        bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(text=f"❌ 條件表產生失敗：{str(e)[:200]}"))
                import threading
                threading.Thread(target=_run_coupon_table, args=(ck, _bot_api, _today), daemon=True).start()
                return
            if arg in ("all", "全部"):
                days_ahead = None
            else:
                try:
                    days_ahead = max(1, min(int(arg), 30)) if arg else 3
                except ValueError:
                    days_ahead = 3
            try:
                msg = _bond_build_alert(str(BOND_PRICE_FILE), today=_today, lookahead=14, days_ahead=days_ahead)
                mtime = datetime.fromtimestamp(BOND_PRICE_FILE.stat().st_mtime, TZ_TAIPEI).strftime("%m/%d %H:%M")
                msg += f"\n📎 報價檔更新於 {mtime}｜/coupon all 看全部｜/coupon table 出Excel"
                chunks_c = [msg[i:i+4900] for i in range(0, len(msg), 4900)]
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=chunks_c[0]))
                for c in chunks_c[1:]:
                    _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=c))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 配息雷達產生失敗：{str(e)[:200]}"))
            return
        if cmd == "fundnav":
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text="📊 手動更新基金淨值中...\n15 檔基金約需 2 分鐘，完成後會通知你 ✅"
            ))
            def _run_fundnav():
                try:
                    print("[FUNDNAV] 開始執行...")
                    job_fund_nav_update()
                    print("[FUNDNAV] 執行完成")
                except Exception as e:
                    import traceback
                    print(f"[FUNDNAV ERROR] {e}")
                    print(traceback.format_exc())
                    user_id = os.getenv("LINE_USER_ID", "")
                    if user_id:
                        line_bot_api.push_message(user_id, TextSendMessage(
                            text=f"❌ /fundnav 執行失敗：{str(e)[:200]}"
                        ))
            import threading
            t = threading.Thread(target=_run_fundnav, daemon=True)
            t.start()
            return
        if cmd == "tracklog":
            with engine.begin() as conn:
                rows = conn.execute(text("SELECT job_name, status, message, executed_at FROM eln_job_log ORDER BY executed_at DESC LIMIT 20")).fetchall()
            if not rows:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="目前沒有執行記錄。"))
                return
            lines = ["📋 最近排程記錄（最新20筆）：\n"]
            status_icon = {"success": "✅", "error": "❌", "started": "🔄", "skipped": "⏭️"}
            for r in rows:
                icon = status_icon.get(r[1], "•")
                tw_time = r[3].astimezone(TZ_TAIPEI_PYTZ).strftime("%m/%d %H:%M")
                msg = f"  {r[2]}" if r[2] else ""
                lines.append(f"{icon} {tw_time} {r[0]}{msg}")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900]))
            return
        if cmd in ("claude", "gpt", "gemini"):
            forced_prompt = text_raw.split(" ", 1)[1].strip() if " " in text_raw else ""
            if not forced_prompt:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"請在 /{cmd} 後面加上問題\n\n例如：/{cmd} 今天美股怎麼看？"))
                return
            model_map = {"claude": "Claude", "gpt": "ChatGPT", "gemini": "Gemini"}
            reply = ai_router(forced_prompt, chat_key=ck, forced_model=cmd)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🦞 龍蝦（{model_map[cmd]}）\n\n{reply[:4700]}"))
            return
        if cmd == "kb":
            parts = text_raw.split(" ", 1)
            arg = parts[1].strip() if len(parts) > 1 else ""
            if arg.lower() in ("上傳", "upload"):
                with engine.begin() as conn:
                    conn.execute(text("""
                    INSERT INTO eln_session(chat_key, await_file, invest_mode, updated_at)
                    VALUES (:k, TRUE, 'kb_upload', NOW())
                    ON CONFLICT (chat_key) DO UPDATE SET await_file=TRUE, invest_mode='kb_upload', updated_at=NOW()
                    """), {"k": ck})
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📚 請直接傳送檔案給我\n\n支援格式：PDF、PPT、Word、圖片\n\n傳完後會自動存入知識庫 ✅"))
                return
            if arg.lower() in ("清單", "list", "列表"):
                try:
                    docs = knowledge.list_documents()
                    if not docs:
                        _bot_api.reply_message(event.reply_token, TextSendMessage(text="📚 知識庫目前沒有任何文件。"))
                        return
                    icons = {"pdf": "📄", "pptx": "📊", "ppt": "📊", "docx": "📝", "doc": "📝"}
                    lines = [f"📚 知識庫文件（共 {len(docs)} 份）\n"]
                    for d in docs:
                        ext = d["filename"].split(".")[-1].lower()
                        icon = icons.get(ext, "📎")
                        lines.append(f"{icon} {d['filename']}（{d['page_count']} 頁）")
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900]))
                except Exception as e:
                    _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 讀取清單失敗：{str(e)[:200]}"))
                return
            if not arg:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="📚 知識庫指令\n─────────────────\n/kb <問題> → 查詢知識庫\n/kb上傳 → 上傳檔案到知識庫\n/kb清單 → 查看已上傳文件\n─────────────────\n範例：\n/kb PIMCO收益基金的投資策略\n/kb ELN的KO條件是什麼"))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📚 查詢知識庫中，請稍候..."))
            try:
                result = knowledge.query_knowledge(arg)
                answer = result.get("answer", "查無結果")
                sources = result.get("sources", [])
                src_text = ""
                if sources:
                    src_text = "\n\n📍 來源：" + "、".join([f"{s['filename']} 第{s['page']}頁" for s in sources[:3]])
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📚 知識庫\n\n{answer[:4500]}{src_text}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 查詢失敗：{str(e)[:200]}"))
            return
        if cmd == "save":
            content = text_raw[len("/save"):].strip()
            if not content:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請在 /save 後面加上文字或網址\n\n範例：\n/save https://...\n/save 這篇文章說..."))
                return
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📥 儲存中，正在產生摘要..."))
            try:
                summary = save_article_text(ck, content)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"✅ 已儲存！\n\n{summary}"))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 儲存失敗：{str(e)[:200]}"))
            return
        if cmd == "unread":
            rows = get_unread_articles(limit=15)
            if not rows:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="✅ 目前沒有未讀文章！"))
                return
            icon_map = {"url": "🔗", "image": "🖼️", "text": "📝"}
            lines = [f"📚 未讀文章（共 {len(rows)} 篇）：\n"]
            for row in rows:
                icon = icon_map.get(row[2], "📄")
                dt = row[3].astimezone(TZ_TAIPEI_PYTZ).strftime("%m/%d")
                lines.append(f"{icon} #{row[0]} {(row[1] or '無標題')[:28]}  ({dt})")
            lines.append("\n輸入 /article <編號> 看摘要\n輸入 /read <編號> 標記已讀")
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines)[:4900]))
            return
        if cmd == "article":
            parts2 = text_raw.split(" ", 1)
            if len(parts2) < 2 or not parts2[1].strip().isdigit():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入文章編號\n例：/article 3"))
                return
            article_id = int(parts2[1].strip())
            row = get_article_detail(article_id)
            if not row:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到文章 #{article_id}"))
                return
            source_label = {"url": "🔗 網址", "image": "🖼️ 圖片", "text": "📝 文字"}.get(row[4], "📄")
            status = "✅ 已讀" if row[5] else "📌 未讀"
            dt = row[6].astimezone(TZ_TAIPEI_PYTZ).strftime("%Y/%m/%d %H:%M")
            msg = f"📄 文章 #{row[0]}\n標題：{row[1] or '無標題'}\n類型：{source_label} {status}\n儲存時間：{dt}\n───────────\n{row[3] or '無摘要'}"
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=msg[:4900]))
            return
        if cmd == "read":
            parts2 = text_raw.split(" ", 1)
            if len(parts2) < 2 or not parts2[1].strip().isdigit():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入文章編號\n例：/read 3"))
                return
            try:
                mark_article_read(int(parts2[1].strip()))
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"✅ 文章 #{parts2[1].strip()} 已標記為已讀！"))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 失敗：{str(e)[:200]}"))
            return
        if cmd == "del":
            parts2 = text_raw.split(" ", 1)
            if len(parts2) < 2 or not parts2[1].strip().isdigit():
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入文章編號\n例：/del 3"))
                return
            article_id = int(parts2[1].strip())
            try:
                with engine.begin() as conn:
                    conn.execute(text("DELETE FROM articles WHERE id = :i"), {"i": article_id})
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🗑️ 文章 #{article_id} 已刪除！"))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"❌ 刪除失敗：{str(e)[:200]}"))
            return
        if cmd == "web":
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📚 龍蝦文章庫\n\nhttps://eln-bot.onrender.com/articles"))
            return
        if cmd == "dbcheck":
            parts2 = text_raw.split(" ", 1)
            bid = parts2[1].strip() if len(parts2) > 1 else "WMGS26040252"
            with engine.begin() as conn:
                row = conn.execute(text("SELECT bond_id, agent_name, detail, updated_at FROM eln_detail WHERE chat_key=:k AND bond_id ILIKE :b LIMIT 1"), {"k": ck, "b": f"%{bid}%"}).fetchone()
            if not row:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"DB 查不到 {bid}，chat_key={ck}"))
                return
            updated = str(row[3])
            detail_preview = (row[2] or "")[:500]
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"bond_id: {row[0]}\nagent: {row[1]}\nupdated_at: {updated}\n\n{detail_preview}"))
            return
        if cmd == "forget":
            try:
                with engine.begin() as conn:
                    conn.execute(text("DELETE FROM chat_history WHERE chat_key = :k"), {"k": ck})
                _bot_api.reply_message(event.reply_token, TextSendMessage(text="🧹 記憶已清除！龍蝦從頭開始囉。"))
            except Exception as e:
                _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"清除失敗：{e}"))
            return
        if any(k in text_raw for k in SPENDING_NL_KEYWORDS):
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="💳 收到！正在幫你分析消費明細，請稍候約30秒..."))
            try:
                from spending_analyzer import get_monthly_spending_report
                report = get_monthly_spending_report(days=31)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=report[:4900]))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 消費分析失敗：{str(e)[:200]}"))
            return
        reply = ai_router(text_raw, chat_key=ck)
        _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🦞 龍蝦\n\n{reply[:4700]}"))
    except Exception as e:
        print("[ERROR] handle_text_message:", e)
        print(_traceback.format_exc())
        try:
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="我收到訊息但處理時出錯了。你可以先輸入 /help。"))
        except Exception:
            pass

# ==============================
# File message handler
# ==============================
UPLOAD_DIR = Path("/tmp/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

def extract_text_from_file(file_path: str, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    text = ""
    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, start=1):
                text += f"\n--- 第 {i} 頁 ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        elif ext in (".xlsx", ".xls"):
            import pandas as pd
            xl = pd.ExcelFile(file_path)
            for sheet in xl.sheet_names:
                df = xl.parse(sheet)
                text += f"\n--- 工作表: {sheet} ---\n"
                text += df.to_string(index=False) + "\n"
    except Exception as e:
        print(f"Extract error: {e}")
        text = ""
    return text.strip()

def analyze_file_with_claude(text: str, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    type_map = {".pdf": "PDF文件", ".docx": "Word文件", ".pptx": "PowerPoint簡報", ".xlsx": "Excel試算表", ".xls": "Excel試算表"}
    file_type = type_map.get(ext, "文件")
    prompt = (f"我收到一份{file_type}，內容如下:\n\n{text[:6000]}\n\n"
              "請幫我:\n1. 用2-3句話說明這份文件的主題與目的\n"
              "2. 條列出5-8個最重要的重點\n3. 如果有數據或結論，特別標示出來\n"
              "4. 最後一句話說明這份文件的主要價值或建議行動\n\n"
              "格式規定: 不使用 Markdown 符號（禁止 ## ** --- 等），標題用 emoji，條列用 •")
    resp = claude_client.messages.create(model="claude-sonnet-4-6", max_tokens=1500, messages=[{"role": "user", "content": prompt}])
    return (resp.content[0].text or "").strip()

def analyze_image_with_claude(image_data: bytes, media_type: str) -> str:
    image_b64 = _base64.b64encode(image_data).decode("utf-8")
    resp = claude_client.messages.create(
        model="claude-sonnet-4-6", max_tokens=1500,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": image_b64}},
            {"type": "text", "text": "請分析這張圖片，幫我:\n1. 說明圖片的主要內容\n2. 如果有文字或數據，擷取重要資訊\n3. 條列出重點\n格式規定: 不使用 Markdown 符號，標題用 emoji，條列用 •"}
        ]}]
    )
    return (resp.content[0].text or "").strip()

def generate_invest_post(image_data: bytes, reason: str, targets: str) -> str:
    image_b64 = _base64.b64encode(image_data).decode("utf-8")
    user_input = ""
    if reason:
        user_input += f"投資理由：{reason}\n"
    if targets:
        user_input += f"建議標的：{targets}\n"
    prompt = f"""你是一位台灣私人銀行的投資輔銷人員，正在為高資產客戶撰寫 LINE 群組推播文。
根據上方的新聞截圖，結合以下我提供的投資觀點，生成兩個版本的推播文：
{user_input}
【規格要求】- 每個版本 100-250 字 - 繁體中文 - 不使用 Markdown - 用 emoji 當標題和分段符號 - 結尾附上建議標的
【版本一：專業版】適合傳給高資產客戶，語氣專業簡練，強調市場邏輯和風險意識。
【版本二：輕鬆版】適合一般投資群組，語氣親切，用比喻讓人容易理解，帶點觀點但不失專業。
格式：===專業版===（內容）===輕鬆版===（內容）"""
    resp = claude_client.messages.create(
        model="claude-sonnet-4-6", max_tokens=1500,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": image_b64}},
            {"type": "text", "text": prompt}
        ]}]
    )
    return (resp.content[0].text or "").strip()

@handler.add(MessageEvent, message=FileMessage)
def handle_file_message(event):
    _bot_api = getattr(_current_bot_api, "api", None) or line_bot_api
    try:
        ck = chat_key_of(event)
        filename = getattr(event.message, "file_name", "") or ""
        ext = Path(filename).suffix.lower()
        print("[FILE]", ck, filename)
        # ── 海外債群組：只收 Bond_Pricing 報價檔（.xlsx），其他檔案不理 ──
        if event.source.type in ("group", "room") and is_bond_query_group(ck):
            return  # 查價群不收任何檔案
        _bond_group = event.source.type in ("group", "room") and is_bond_group_chat(ck)
        if _bond_group and ext not in (".xlsx", ".xlsm"):
            return
        message_id = event.message.id
        content = _bot_api.get_message_content(message_id)
        tmp_path = UPLOAD_DIR / f"upload_{int(datetime.now(TZ_TAIPEI).timestamp())}{ext}"
        with open(tmp_path, "wb") as f:
            for chunk in content.iter_content():
                f.write(chunk)
        if ext in (".mp3", ".m4a", ".wav", ".ogg", ".mp4", ".webm"):
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"🎙️ 收到音檔 {filename}，正在轉逐字稿..."))
            with open(tmp_path, "rb") as f:
                audio_data = f.read()
            text_result = transcribe_audio(audio_data, filename=filename)
            if not text_result:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="❌ 無法辨識語音內容，請確認音檔有聲音。"))
                return
            summary = build_transcript_summary(text_result, chat_key=ck)
            db_set_transcript_cache(ck, text_result, summary)
            db_save_meeting_transcript(ck, filename, text_result, summary)
            push_long_message(_bot_api, ck.split(":", 1)[1], f"📝 逐字稿（前段）：\n\n{text_result[:2000]}")
            push_long_message(_bot_api, ck.split(":", 1)[1], f"📌 會議摘要：\n\n{summary}")
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="要不要把這份會議重點做成 PDF？\n\n可直接回：做成 PDF / 不用"))
            return
        invest_mode_now, _ = db_invest_get(ck)
        if invest_mode_now == "kb_upload" and ext in (".pdf", ".pptx", ".ppt", ".docx", ".doc", ".jpg", ".jpeg", ".png", ".gif", ".webp"):
            db_set_await(ck, False)
            with engine.begin() as conn:
                conn.execute(sql_text("UPDATE eln_session SET invest_mode='', await_file=FALSE WHERE chat_key=:k"), {"k": ck})
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📚 收到！正在處理 {filename} 並存入知識庫，請稍候..."))
            try:
                with open(tmp_path, "rb") as f:
                    file_bytes = f.read()
                result = knowledge.process_and_index_file(filename, file_bytes)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(
                    text=f"✅ 已存入知識庫！\n\n📄 {filename}\n📑 {result['pages']} 頁\n🔍 {result['chunks']} 個索引\n\n現在可以用 /kb 問題 來查詢了！"
                ))
            except Exception as e:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 存入知識庫失敗：{str(e)[:200]}"))
            return
        # ── 海外債報價檔：存成最新檔 + 立刻重算配息雷達 ──
        _looks_like_bond = ext in (".xlsx", ".xlsm") and ("bond" in filename.lower())
        _is_bond = False
        if ext in (".xlsx", ".xlsm") and _BOND_RADAR_OK:
            try:
                _is_bond = _is_bond_pricing_file(str(tmp_path), filename)
            except Exception as _e:
                print(f"[BondRadar] 偵測失敗：{_e}")
        print(f"[BondRadar] 模組OK={_BOND_RADAR_OK} 檔名像債券={_looks_like_bond} 偵測結果={_is_bond}")
        if _bond_group and not _is_bond:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
            return
        if _looks_like_bond and not _BOND_RADAR_OK:
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text="⚠️ 檔名看起來是海外債報價檔，但 bond_coupon_alert.py 模組沒有載入成功。\n請確認該檔案已放進 repo 根目錄並重新部署（Render Logs 搜尋 [BondRadar] 看錯誤原因）。"))
            return
        if _is_bond:
            db_set_await(ck, False)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"📥 收到海外債報價檔 {filename}，正在更新配息雷達..."))
            def _run_bond_update(tmp_path_str, chat_key, bot_api_ref, fname):
                try:
                    # ── 歷史回填:檔名日期比今天舊 3 天以上 → 只存歷史快照,不動最新報價檔、不跑雷達 ──
                    _snap0 = parse_pricing_file_date(fname)
                    _today0 = datetime.now(TZ_TAIPEI).date()
                    if _snap0 and (_today0 - _snap0).days >= 3:
                        n_hist = save_price_history(snap_date=_snap0, path=tmp_path_str)
                        with engine.begin() as conn:
                            n_days, d_min, d_max = conn.execute(sql_text(
                                "SELECT COUNT(DISTINCT snap_date), MIN(snap_date), MAX(snap_date) FROM bond_price_history")).fetchone()
                        bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(
                            text=f"📚 已補入 {_snap0:%Y/%m/%d} 歷史報價快照（{n_hist} 檔）\n"
                                 f"最新報價檔與配息雷達未變更。\n"
                                 f"目前歷史庫:{n_days} 個交易日（{d_min:%m/%d}～{d_max:%m/%d}），"
                                 f"可用 /move 30 5 之類的指令比較。"))
                        try:
                            os.remove(tmp_path_str)
                        except Exception:
                            pass
                        return
                    import shutil
                    BOND_PRICE_FILE.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy(tmp_path_str, str(BOND_PRICE_FILE))
                    write_job_log("海外債配息雷達", "file_updated", fname)
                    msg = _bond_build_alert(str(BOND_PRICE_FILE), today=datetime.now(TZ_TAIPEI).date())
                    push_long_message(bot_api_ref, chat_key.split(":", 1)[1], "✅ 報價檔已更新\n\n" + msg)
                    # 信評異動 / 新上架 / 下架
                    try:
                        diff_txt = snapshot_and_diff()
                        if diff_txt:
                            push_long_message(bot_api_ref, chat_key.split(":", 1)[1], "📋 與上一份報價檔比對\n" + diff_txt)
                    except Exception as e:
                        print(f"[BondSnapshot ERROR] {e}")
                    # 報價歷史 + 異動（vs 上一份 ≥2%；vs 7天前 ≥3%）
                    try:
                        _snap = parse_pricing_file_date(fname)
                        if _snap and _snap > datetime.now(TZ_TAIPEI).date():
                            _snap = None  # 未來日期視為解析錯誤,退回今天
                        save_price_history(snap_date=_snap)
                        for db_, th_ in ((1, 2.0), (7, 3.0)):
                            mv, _, _ = price_movers(days_back=db_, threshold_pct=th_)
                            if mv:
                                push_long_message(bot_api_ref, chat_key.split(":", 1)[1], mv)
                    except Exception as e:
                        print(f"[BondMovers ERROR] {e}")
                    # 到價通知
                    try:
                        n_hit = check_bond_alerts(bot_api_ref, source="upload")
                        print(f"[BondAlert] upload check hits={n_hit}")
                    except Exception as e:
                        print(f"[BondAlert ERROR] {e}")
                except Exception as e:
                    print(f"[BondRadar ERROR] {e}")
                    print(_traceback.format_exc())
                    bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(text=f"❌ 配息雷達更新失敗：{str(e)[:200]}"))
            import threading
            threading.Thread(target=_run_bond_update, args=(str(tmp_path), ck, _bot_api, filename), daemon=True).start()
            return
        if ext in (".xlsx", ".xls"):
            db_set_await(ck, False)
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="📊 收到 Excel！計算中，請稍候約30秒..."))
            def _run_calc(tmp_path_str, chat_key, bot_api_ref):
                try:
                    print(f"[CALC] 開始計算，chat_key={chat_key}")
                    summary, top5_lines, detail_map, agent_name_map = run_autotracking(tmp_path_str)
                    print(f"[CALC] 計算完成，共 {len(detail_map)} 筆，開始寫入 chat_key={chat_key}")
                    db_save_result(chat_key, summary, top5_lines, detail_map, agent_name_map)
                    print(f"[CALC] 寫入完成，共 {len(detail_map)} 筆，chat_key={chat_key}")
                    try:
                        from eln_storage import upload_eln_excel
                        upload_eln_excel(tmp_path_str)
                    except Exception as e:
                        print("[ELN Storage] upload failed:", e)
                    bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(text=(summary or "已收到檔案，但沒有產出內容")[:4900]))
                except Exception as e:
                    print(f"[CALC ERROR] {e}")
                    import traceback as _tb
                    print(_tb.format_exc())
                    bot_api_ref.push_message(chat_key.split(":", 1)[1], TextSendMessage(text=f"❌ 計算失敗：{str(e)[:200]}"))
            import threading
            t = threading.Thread(target=_run_calc, args=(str(tmp_path), ck, _bot_api), daemon=True)
            t.start()
            return
        if ext in (".pdf", ".docx", ".pptx"):
            _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"收到！正在分析 {filename}，請稍候..."))
            text = extract_text_from_file(str(tmp_path), filename)
            if not text:
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="檔案解析失敗，可能是掃描版 PDF 或格式不支援。"))
                return
            analysis = analyze_file_with_claude(text, filename)
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=analysis[:4900]))
            try:
                from pdf_generator import create_and_upload_pdf
                link = create_and_upload_pdf("analysis", analysis, filename)
                _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"📄 分析報告 PDF：\n{link}"))
            except Exception as e:
                print(f"PDF upload error: {e}")
            return
        _bot_api.reply_message(event.reply_token, TextSendMessage(text=f"目前支援的檔案格式: PDF、Word、PowerPoint、Excel\n收到的格式 {ext} 暫不支援。"))
    except Exception as e:
        print("[ERROR] handle_file_message:", e)
        print(_traceback.format_exc())
        try:
            db_set_await(chat_key_of(event), False)
        except Exception:
            pass
        try:
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="檔案處理時出錯了，請稍後再試。"))
        except Exception:
            pass

@handler.add(MessageEvent, message=ImageMessage)
def handle_image_message(event):
    _bot_api = getattr(_current_bot_api, "api", None) or line_bot_api
    try:
        ck = chat_key_of(event)
        print("[IMAGE]", ck)
        if event.source.type in ("group", "room") and (is_bond_group_chat(ck) or is_bond_query_group(ck)):
            return  # 海外債群組不處理圖片
        message_id = event.message.id
        content = _bot_api.get_message_content(message_id)
        image_data = b""
        for chunk in content.iter_content():
            image_data += chunk
        invest_mode, _ = db_invest_get(ck)
        if invest_mode == "await_image":
            db_invest_set(ck, "await_reason", image=image_data)
            _bot_api.reply_message(event.reply_token, TextSendMessage(
                text="✅ 收到截圖！\n\n請輸入你的投資理由和標的：\n\n理由：（你認為能投資的原因）\n標的：（股票/ETF代號）"
            ))
            return
        _bot_api.reply_message(event.reply_token, TextSendMessage(text="🖼️ 收到圖片！正在分析並儲存到文章庫..."))
        try:
            summary = save_article_image(image_data, str(message_id))
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(
                text=f"✅ 圖片已儲存到文章庫！\n\n{summary}\n\n打 /unread 可以查看所有未讀文章"
            ))
        except Exception as e:
            print(f"[IMAGE SAVE ERROR] {e}")
            analysis = analyze_image_with_claude(image_data, "image/jpeg")
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=analysis[:4900]))
    except Exception as e:
        print("[ERROR] handle_image_message:", e)
        print(_traceback.format_exc())
        try:
            _bot_api.reply_message(event.reply_token, TextSendMessage(text="圖片處理時出錯了，請稍後再試。"))
        except Exception:
            pass

def transcribe_audio(audio_data: bytes, filename: str = "audio.m4a") -> str:
    import tempfile
    if not openai_client:
        raise RuntimeError("缺少 OPENAI_API_KEY，無法使用語音轉文字")
    MAX_BYTES = 24 * 1024 * 1024
    with tempfile.TemporaryDirectory() as tmp:
        ext = os.path.splitext(filename)[1].lower() or ".m4a"
        src_path = os.path.join(tmp, f"audio_input{ext}")
        with open(src_path, "wb") as f:
            f.write(audio_data)
        if len(audio_data) > MAX_BYTES:
            try:
                from pydub import AudioSegment
                out_path = os.path.join(tmp, "audio_compressed.mp3")
                audio = AudioSegment.from_file(src_path)
                audio = audio.set_channels(1).set_frame_rate(16000)
                audio.export(out_path, format="mp3", bitrate="32k")
                send_path = out_path
            except Exception as e:
                print(f"[Audio] 壓縮失敗，嘗試直接送原檔: {e}")
                send_path = src_path
        else:
            send_path = src_path
        with open(send_path, "rb") as f:
            resp = openai_client.audio.transcriptions.create(model="whisper-1", file=(os.path.basename(send_path), f), language="zh")
        return resp.text.strip()

@handler.add(MessageEvent, message=AudioMessage)
def handle_audio_message(event, _override_bot_api=None):
    _bot_api = _override_bot_api or line_bot_api
    ck = chat_key_of(event)
    print(f"[AUDIO] {ck}")
    if event.source.type in ("group", "room") and (is_bond_group_chat(ck) or is_bond_query_group(ck)):
        return  # 海外債群組不處理語音
    try:
        _bot_api.reply_message(event.reply_token, TextSendMessage(text="🎙️ 收到語音，轉文字中，請稍候..."))
        message_id = event.message.id
        content = _bot_api.get_message_content(message_id)
        audio_data = b""
        for chunk in content.iter_content():
            audio_data += chunk
        text_result = transcribe_audio(audio_data)
        if not text_result:
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text="❌ 無法辨識語音內容，請確認音檔有聲音。"))
            return
        push_long_message(_bot_api, ck.split(":", 1)[1], f"📝 語音轉文字：\n\n{text_result}")
        reply = ai_router(text_result, chat_key=ck)
        push_long_message(_bot_api, ck.split(":", 1)[1], reply)
    except Exception as e:
        print(f"[ERROR] handle_audio_message: {e}")
        try:
            _bot_api.push_message(ck.split(":", 1)[1], TextSendMessage(text=f"❌ 語音處理失敗：{str(e)[:200]}"))
        except Exception:
            pass

# ==============================
# ELN Auto-Tracking 群組專用 handler
# ==============================
ELN_PERSONAL_CHAT_KEY = f"user:{os.getenv('LINE_USER_ID', '')}"

@eln_group_handler.add(MessageEvent, message=TextMessage)
def handle_eln_group_message(event):
    try:
        text_raw = (event.message.text or "").strip()
        tl = text_raw.lower().strip()
        uid = event.source.user_id if hasattr(event.source, "user_id") else "unknown"
        print(f"[ELN-G USER] uid={uid} msg={repr(text_raw[:50])}")
        if not (tl.startswith("/list") or tl.startswith("/detail") or tl.startswith("/nc") or tl.startswith("/內規")):
            return
        ck = ELN_PERSONAL_CHAT_KEY
        from linebot.models import TextSendMessage as TSM
        
        # ==========================================
        # 群組版 內規專屬指令攔截 (handle_eln_group_message)
        # ==========================================
        if tl.startswith("/內規"):
            actual_query = text_raw.replace("/內規", "").strip()
            if not actual_query:
                eln_group_bot_api.reply_message(event.reply_token, TSM(text="請在指令後面加上想查詢的內容喔！\n例如：/內規 Lombard lending 最高可以到幾歲？"))
                return
            try:
                file_path = Path("regulations.txt")
                if not file_path.exists():
                    eln_group_bot_api.reply_message(event.reply_token, TSM(text="❌ 找不到 regulations.txt，請確認已將法規檔案上傳至系統。"))
                    return
                regulation_text = file_path.read_text(encoding="utf-8")
                
                # 防縮排錯誤的 Prompt 組合方式
                prompt_lines = [
                    "你現在是銀行的法遵與內部規範專家。請根據以下【內部規範全文】，直接且精準回答同仁的問題。",
                    "",
                    "【嚴格限制】",
                    "1. 絕對不要輸出任何「因為文本是程式碼...」或「無法回答」的廢話警告。",
                    "2. 絕對不要在結尾補充「可以這樣跟客戶/專員說」的話術。",
                    "3. 嚴禁使用 Markdown 語法 (例如 **, ##, --- 等)，請用純文字或 Emoji 條列排版。",
                    "",
                    "【內部規範全文】",
                    regulation_text,
                    "",
                    "【同仁問題】",
                    actual_query
                ]
                prompt = "\n".join(prompt_lines)

                answer = ai_claude_long(prompt, chat_key=ck)
                final_answer = "⚠️ 僅供參考，本回覆由 AI 統整，不代表總行最終內規解釋。\n\n" + answer
                
                eln_group_bot_api.reply_message(event.reply_token, TSM(text=final_answer[:4900]))
            except Exception as e:
                eln_group_bot_api.reply_message(event.reply_token, TSM(text=f"❌ 內部規範查詢失敗：{e}"))
            return
        # ==========================================
        
        if tl.startswith("/list"):
            from collections import defaultdict
            list_parts = text_raw.split(" ", 2)
            is_detail_mode = len(list_parts) > 1 and list_parts[1].strip().lower() == "detail"
            if is_detail_mode:
                name_filter = list_parts[2].strip() if len(list_parts) > 2 else ""
            else:
                name_filter = list_parts[1].strip() if len(list_parts) > 1 else ""
            bonds = db_list_bonds(ck, limit=200)
            if not bonds:
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無資料。"))
                return
            # /list detail 姓名：顯示完整 detail，只顯示比價中
            if is_detail_mode:
                if not name_filter:
                    eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入理專名稱\n例：/list detail 小美"))
                    return
                matched_details = []
                seen = set()
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()]
                    if any(name_filter in a for a in agents) and bond_id not in seen:
                        if bond_status_tag(detail) == "":
                            matched_details.append(detail)
                            seen.add(bond_id)
                if not matched_details:
                    eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到「{name_filter}」比價中的持倉。"))
                    return
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=f"👤 {name_filter} 比價中商品（共 {len(matched_details)} 筆）"))
                for det in matched_details:
                    eln_group_bot_api.push_message(event.source.user_id, TextSendMessage(text=det[:4900]))
                return
            detail_map_status = {bond_id: bond_status_tag(detail) for bond_id, _, detail in bonds}
            if name_filter:
                matched = []
                seen = set()
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()]
                    if any(name_filter in a for a in agents) and bond_id not in seen:
                        matched.append((bond_id, detail_map_status.get(bond_id, "")))
                        seen.add(bond_id)
                if not matched:
                    eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到「{name_filter}」的持倉。"))
                    return
                lines = [f"👤 {name_filter} 的持倉（共 {len(matched)} 筆）：\n"] + [f"   • {b}{tag}" for b, tag in matched]
                full_text = "\n".join(lines)
            else:
                grouped = defaultdict(list)
                for bond_id, agent_raw, detail in bonds:
                    agents = [a.strip() for a in re.split(r"[,，、/]", agent_raw) if a.strip()] or ["未指定"]
                    for agent in agents:
                        if bond_id not in [b for b, _ in grouped[agent]]:
                            grouped[agent].append((bond_id, detail_map_status.get(bond_id, "")))
                lines = [f"📋 全部商品（共 {len(set(b for b,_,_ in bonds))} 筆，按理專排列）：\n"]
                for agent, bond_list in sorted(grouped.items()):
                    lines.append(f"👤 {agent}（{len(bond_list)} 筆）")
                    lines += [f"   • {b}{tag}" for b, tag in bond_list]
                full_text = "\n".join(lines)
            chunks = []
            current = ""
            for line in full_text.split("\n"):
                if len(current) + len(line) + 1 > 4800:
                    chunks.append(current)
                    current = line
                else:
                    current = current + "\n" + line if current else line
            if current:
                chunks.append(current)
            eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=chunks[0]))
            for chunk in chunks[1:]:
                eln_group_bot_api.push_message(event.source.user_id, TextSendMessage(text=chunk))
            return
        if tl.startswith("/detail"):
            parts = text_raw.split(" ", 1)
            if len(parts) < 2 or not parts[1].strip():
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入：/detail 商品代號"))
                return
            matched_id, detail, candidates = db_find_detail(ck, parts[1].strip())
            if detail:
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=detail[:4900]))
                return
            if candidates and matched_id is None:
                sample = "\n".join([f"• {c}" for c in candidates[:20]])
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=f"請再精準一點，候選代號：\n{sample}"[:4900]))
                return
            eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="查不到該代號。"))
            return
        if tl.startswith("/nc"):
            nc_parts = text_raw.split(" ")
            if len(nc_parts) < 2 or not nc_parts[1].strip():
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="請輸入：/nc YYYYMM\n例：/nc 202606\n或：/nc 202606 小美"))
                return
            qm_nc = nc_parts[1].strip().replace("/", "").replace("-", "")
            if len(qm_nc) != 6 or not qm_nc.isdigit():
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="格式錯誤，請輸入6位數字\n例：/nc 202606"))
                return
            yr_nc = qm_nc[:4]
            mo_nc = qm_nc[4:]
            name_filter_nc = nc_parts[2].strip() if len(nc_parts) > 2 else ""
            search_str_nc = f"{yr_nc}-{mo_nc}"
            with engine.begin() as conn:
                rows = conn.execute(text("SELECT bond_id, agent_name, detail FROM eln_detail WHERE chat_key=:k ORDER BY agent_name ASC, bond_id ASC"), {"k": ck}).fetchall()
            if not rows:
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="目前尚無資料。"))
                return
            matched_nc = []
            for bond_id, agent_name, detail in rows:
                m = re.search(r"NC閉鎖期 \(至 (\d{4}-\d{2})-\d{2}\)", detail)
                if m and m.group(1) == search_str_nc:
                    if name_filter_nc:
                        ags = [a.strip() for a in re.split(r"[,，、/]", agent_name or "") if a.strip()]
                        if not any(name_filter_nc in a for a in ags):
                            continue
                    matched_nc.append((bond_id, agent_name or "-", bond_status_tag(detail)))
            if not matched_nc:
                tip = f"「{name_filter_nc}」" if name_filter_nc else ""
                eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text=f"找不到 {yr_nc}/{mo_nc} 閉鎖期打開{tip}的商品。"))
                return
            tip2 = f"（{name_filter_nc}）" if name_filter_nc else ""
            lines_nc = [f"🔓 {yr_nc}/{mo_nc} 閉鎖期打開{tip2}（共 {len(matched_nc)} 筆）:\n"]
            lines_nc += [f"   • {bid} [{ag}]{tag}" for bid, ag, tag in matched_nc]
            eln_group_bot_api.reply_message(event.reply_token, TextSendMessage(text="\n".join(lines_nc)[:4900]))
            return
    except Exception as e:
        print(f"[ELN-GROUP ERROR] {e}")

# ==============================
# 內建排程
# ==============================
TZ_TAIPEI_PYTZ = pytz.timezone("Asia/Taipei")

def job_daily_report():
    now = datetime.now(TZ_TAIPEI_PYTZ)
    if now.weekday() >= 6:  # 只跳過週日
        write_job_log("財經日報", "skipped", "週日跳過")
        return
    write_job_log("財經日報", "started", now.strftime('%Y-%m-%d %H:%M'))
    try:
        from daily_report import main as report_main
        report_main()
        write_job_log("財經日報", "success", "推播成功")
    except Exception as e:
        write_job_log("財經日報", "error", str(e))

def job_bond_daily_report():
    now = datetime.now(TZ_TAIPEI_PYTZ)
    if now.weekday() >= 6:  # 只跳過週日
        write_job_log("債券日報", "skipped", "週日跳過")
        return
    write_job_log("債券日報", "started", now.strftime('%Y-%m-%d %H:%M'))
    try:
        from bond_daily_report import main as bond_report_main
        bond_report_main()
        write_job_log("債券日報", "success", "推播成功")
    except Exception as e:
        write_job_log("債券日報", "error", str(e))

def job_alert_monitor():
    try:
        from alert_monitor import main as alert_main
        alert_main()
    except Exception as e:
        print(f"[Scheduler] 價格警示失敗: {e}")

def job_mail_monitor():
    try:
        from mail_monitor import main as mail_main
        mail_main()
    except Exception as e:
        print(f"[Scheduler] 郵件監控失敗: {e}")

def job_article_reminder():
    try:
        rows = get_unread_articles(limit=100)
        if not rows:
            write_job_log("未讀提醒", "skipped", "沒有未讀文章")
            return
        targets = load_targets()
        target_id = targets.get("default", "")
        if not target_id:
            write_job_log("未讀提醒", "skipped", "沒有設定推播對象")
            return
        msg = f"📚 本週提醒：你有 {len(rows)} 篇文章還沒看！\n\n打 /unread 查看清單\n打 /web 開啟文章庫"
        line_bot_api.push_message(target_id, TextSendMessage(text=msg))
        write_job_log("未讀提醒", "success", f"共{len(rows)}篇未讀")
    except Exception as e:
        write_job_log("未讀提醒", "error", str(e))

def job_fund_nav_update():
    """每天早上9點自動更新基金淨值，也可 /fundnav 手動觸發"""
    from datetime import datetime as _dt
    now = _dt.now(TZ_TAIPEI_PYTZ)
    write_job_log("基金淨值更新", "started", now.strftime('%Y-%m-%d %H:%M'))
    try:
        from update_fund_nav_moneydj import run_fund_nav_update  # ← 正確的模組名稱
        def _line_push(user_id, msg):
            try:
                line_bot_api.push_message(user_id, TextSendMessage(text=msg))
            except Exception as e:
                print(f"LINE推播失敗：{e}")
        updated, skipped, failed = run_fund_nav_update(line_push_fn=_line_push)
        write_job_log("基金淨值更新", "done", f"新增{updated}筆 跳過{skipped}檔 失敗{failed}檔")
    except Exception as e:
        import traceback
        write_job_log("基金淨值更新", "error", str(e))
        user_id = os.getenv("LINE_USER_ID", "")
        if user_id:
            try:
                line_bot_api.push_message(user_id, TextSendMessage(text=f"❌ 基金淨值更新失敗：{str(e)[:200]}"))
            except:
                pass

def write_job_log(job_name: str, status: str, message: str = ""):
    try:
        with engine.begin() as conn:
            conn.execute(text("INSERT INTO eln_job_log (job_name, status, message, executed_at) VALUES (:j, :s, :m, NOW())"),
                         {"j": job_name, "s": status, "m": message[:1000]})
    except Exception as e:
        print(f"[LOG] 寫入失敗: {e}")

def job_spending_report():
    now = datetime.now(TZ_TAIPEI_PYTZ)
    import calendar
    last_day = calendar.monthrange(now.year, now.month)[1]
    if now.day != last_day:
        return
    try:
        from spending_analyzer import get_monthly_spending_report
        report = get_monthly_spending_report(days=31)
        user_id = os.getenv("LINE_USER_ID", "")
        if user_id:
            line_bot_api.push_message(user_id, TextSendMessage(text=report[:4900]))
        write_job_log("月度消費明細", "success", "已發送")
    except Exception as e:
        write_job_log("月度消費明細", "error", str(e))

def job_auto_tracking():
    now = datetime.now(TZ_TAIPEI_PYTZ)
    if now.weekday() >= 6:  # 只跳過週日
        write_job_log("ELN追蹤", "skipped", "週日跳過")
        return
    write_job_log("ELN追蹤", "started", now.strftime('%Y-%m-%d %H:%M'))
    try:
        from auto_tracking_cron import main as tracking_main
        tracking_main()
        write_job_log("ELN追蹤", "success", "追蹤完成")
    except Exception as e:
        write_job_log("ELN追蹤", "error", str(e))

def llm_json_fallback(prompt, max_tokens=4000):
    """
    要 JSON 的 LLM 呼叫，依序 Claude → OpenAI → Gemini，回傳 (parsed_json, source, errors)。
    任一成功就回傳；全失敗回 (None, "", errors)。
    """
    def _parse(raw):
        raw = re.sub(r"^```(?:json)?|```$", "", raw.strip(), flags=re.M).strip()
        return json.loads(raw)
    errs = []
    try:
        resp = claude_client.messages.create(model="claude-sonnet-4-6", max_tokens=max_tokens,
                                             messages=[{"role": "user", "content": prompt}])
        return _parse("".join(getattr(b, "text", "") for b in resp.content)), "claude", errs
    except Exception as e:
        errs.append(f"claude:{str(e)[:60]}")
    if openai_client:
        try:
            resp = openai_client.chat.completions.create(model="gpt-4.1-mini", temperature=0.3, max_tokens=max_tokens,
                                                         messages=[{"role": "user", "content": prompt}])
            return _parse(resp.choices[0].message.content), "openai", errs
        except Exception as e:
            errs.append(f"openai:{str(e)[:60]}")
    if GEMINI_API_KEY:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"
            payload = {"contents": [{"parts": [{"text": prompt}]}]}
            req = urllib.request.Request(url, data=json.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"}, method="POST")
            with urllib.request.urlopen(req, timeout=90) as r_:
                data = json.loads(r_.read().decode("utf-8"))
            return _parse(data["candidates"][0]["content"]["parts"][0]["text"]), "gemini", errs
        except Exception as e:
            errs.append(f"gemini:{str(e)[:60]}")
    return None, "", errs

def get_issuer_intros(issuers):
    """先查 DB 快取，缺的才叫 AI（Claude → OpenAI → Gemini）；同時取得英文名並自動加入信評監控名單"""
    intros = {}
    with engine.begin() as conn:
        rows = conn.execute(text("SELECT issuer, intro FROM bond_issuer_intro")).fetchall()
    cache = {r[0]: r[1] for r in rows}
    missing = [i for i in issuers if i not in cache]
    for iss in issuers:
        if iss in cache:
            intros[iss] = cache[iss]
    if missing:
        prompt = ("你是台灣銀行財富管理部門的固定收益研究員。請針對下列債券發行機構各寫一段 60～100 字的繁體中文簡介，"
                  "內容包含：主要業務、總部所在國家、在產業中的地位、信用體質重點（例如投資等級/政府支持/主要風險）。"
                  "語氣專業中性，不要投資建議。同時給出該機構在國際新聞中最常用的英文名稱（短名，例如 Apple、Verizon、Goldman Sachs、U.S. Treasury）。"
                  "只回傳 JSON 物件，key 為發行機構名稱（必須跟輸入完全一致），value 為物件 {\"intro\": 簡介, \"en\": 英文名}，不要任何其他文字或 markdown。\n\n發行機構名單：\n"
                  + "\n".join(f"- {x}" for x in missing))
        got, source, errs = llm_json_fallback(prompt)
        if got is None:
            for iss in missing:
                intros[iss] = "（簡介暫無法取得，AI 服務不可用：" + "；".join(errs)[:120] + "）"
        else:
            with engine.begin() as conn:
                for iss in missing:
                    v = got.get(iss, {})
                    if isinstance(v, str):
                        v = {"intro": v, "en": ""}
                    txt_ = str(v.get("intro", "")).strip()
                    en_ = str(v.get("en", "")).strip()
                    if txt_:
                        intros[iss] = txt_
                        conn.execute(text("""INSERT INTO bond_issuer_intro(issuer, intro, source, updated_at)
                                             VALUES (:i, :t, :s, NOW())
                                             ON CONFLICT (issuer) DO UPDATE SET intro=EXCLUDED.intro, source=EXCLUDED.source, updated_at=NOW()"""),
                                     {"i": iss, "t": txt_, "s": source})
                    else:
                        intros[iss] = "（AI 未回傳此機構簡介）"
                    if en_:
                        conn.execute(text("""INSERT INTO bond_rating_watch(issuer, en_name, added_by, active)
                                             VALUES (:i, :e, 'auto', TRUE)
                                             ON CONFLICT (issuer) DO UPDATE SET en_name=COALESCE(bond_rating_watch.en_name, EXCLUDED.en_name)"""),
                                     {"i": iss, "e": en_})
    # 已有簡介但還沒進監控名單的（例如舊快取）→ 也加進去，英文名先留空，之後 /rating fix 補
    with engine.begin() as conn:
        for iss in issuers:
            conn.execute(text("""INSERT INTO bond_rating_watch(issuer, added_by, active) VALUES (:i, 'auto', TRUE)
                                 ON CONFLICT (issuer) DO NOTHING"""), {"i": iss})
    return intros

def ensure_en_names():
    """幫監控名單裡沒有英文名的機構補英文名（一次 LLM 呼叫）"""
    with engine.begin() as conn:
        rows = conn.execute(text("SELECT issuer FROM bond_rating_watch WHERE active=TRUE AND (en_name IS NULL OR en_name='')")).fetchall()
    missing = [r[0] for r in rows]
    if not missing:
        return 0
    prompt = ("請給出下列債券發行機構在國際財經新聞中最常用的英文短名（例如 蘋果→Apple、威瑞森電信→Verizon、高盛金融→Goldman Sachs、美國公債→U.S. Treasury）。"
              "只回傳 JSON 物件，key 為中文名（與輸入完全一致），value 為英文名。\n\n" + "\n".join(f"- {x}" for x in missing))
    got, _, _ = llm_json_fallback(prompt, max_tokens=2000)
    if not got:
        return 0
    n = 0
    with engine.begin() as conn:
        for iss in missing:
            en_ = str(got.get(iss, "")).strip()
            if en_:
                conn.execute(text("UPDATE bond_rating_watch SET en_name=:e WHERE issuer=:i"), {"e": en_, "i": iss})
                n += 1
    return n

def run_rating_news_check(days=2, use_llm=True):
    """
    掃監控名單的 Google News RSS，過濾已推過的連結，（可選）用 LLM 過濾出真正的評等動作並中文摘要。
    回傳推播用文字（沒有新聞則空字串）。
    """
    from bond_rating_news import fetch_rating_news, format_news_block
    ensure_en_names()
    with engine.begin() as conn:
        rows = conn.execute(text("SELECT issuer, en_name FROM bond_rating_watch WHERE active=TRUE ORDER BY issuer")).fetchall()
        seen = {r[0] for r in conn.execute(text("SELECT link FROM bond_rating_news_seen WHERE seen_at > NOW() - INTERVAL '30 days'")).fetchall()}
    if not rows:
        return ""
    stats = {"issuers": len(rows), "raw": 0, "new": 0, "kept": 0, "hit_issuers": []}
    fresh = {}
    _AGENCY_QUERY = {  # 評等機構自己也是發行機構時,用公司主體名搜尋,避免抓到它對別人的評等動作
        "moody's": '"Moody\'s Corporation"', "moodys": '"Moody\'s Corporation"',
        "s&p": '"S&P Global Inc"', "s&p global": '"S&P Global Inc"', "standard & poor's": '"S&P Global Inc"',
        "fitch": '"Fitch Group"', "msci": '"MSCI Inc"',
    }
    for iss, en in rows:
        en_q = _AGENCY_QUERY.get((en or "").strip().lower(), en or "")
        items = fetch_rating_news(en_q, zh_name=iss if iss != (en or "") else "", days=days)
        stats["raw"] += len(items)
        items = [it for it in items if it["link"] not in seen]
        stats["new"] += len(items)
        if items:
            fresh[iss] = items
        with engine.begin() as conn:
            conn.execute(text("UPDATE bond_rating_watch SET last_checked=NOW() WHERE issuer=:i"), {"i": iss})
    if not fresh:
        globals()["_LAST_RATING_STATS"] = stats
        return ""
    # LLM 二次過濾：只留「真的評等/展望動作」，順便中文一句話
    keep = fresh
    summary_map = {}
    if use_llm:
        flat = []
        for iss, items in fresh.items():
            for k, it in enumerate(items):
                flat.append({"id": f"{iss}|{k}", "issuer": iss, "title": it["title"], "source": it["source"]})
        prompt = ("以下是新聞標題清單，每則附有我們關注的發行機構名稱(issuer)。"
                  "請只挑出『信評機構(Moody's/S&P/Fitch)對【該 issuer 本身】做出的評等、展望或觀察名單動作』的標題。"
                  "嚴格規則:被評等的對象必須就是 issuer 這家公司;若 issuer 本身是評等機構(例如穆迪、標準普爾、惠譽)，"
                  "則只收錄『它自己被評等或其公司債相關』的新聞，它對其他公司做的評等動作一律排除。"
                  "也排除產品/財報/股價等無關新聞。"
                  "回傳 JSON 物件，key 為 id，value 為 20 字內的繁體中文摘要（例如「穆迪將展望調為負向」）。沒有符合的回傳 {}。\n\n"
                  + json.dumps(flat, ensure_ascii=False))
        got, _, _ = llm_json_fallback(prompt, max_tokens=3000)
        if isinstance(got, dict):
            keep = {}
            for iss, items in fresh.items():
                sel = []
                for k, it in enumerate(items):
                    key = f"{iss}|{k}"
                    if key in got:
                        it["zh"] = str(got[key])
                        sel.append(it)
                if sel:
                    keep[iss] = sel
    stats["kept"] = sum(len(v) for v in keep.values()) if keep else 0
    stats["hit_issuers"] = list(keep.keys()) if keep else []
    globals()["_LAST_RATING_STATS"] = stats
    if not keep:
        # 全部被 LLM 判定為雜訊：仍記錄為已看過，避免明天重複評估
        with engine.begin() as conn:
            for iss, items in fresh.items():
                for it in items:
                    conn.execute(text("INSERT INTO bond_rating_news_seen(link, issuer, title) VALUES (:l,:i,:t) ON CONFLICT DO NOTHING"),
                                 {"l": it["link"], "i": iss, "t": it["title"][:300]})
        return ""
    blocks = ["🚨 外部信評異動雷達（Google News，近 %d 天）" % days]
    with engine.begin() as conn:
        for iss, items in keep.items():
            lines = [f"🏦 {iss}"]
            for it in items[:4]:
                d = it["published"].astimezone(TZ_TAIPEI).strftime("%m/%d") if it["published"] else ""
                zh = f"｜{it['zh']}" if it.get("zh") else ""
                src = f"（{it['source']}）" if it["source"] else ""
                lines.append(f"▪ {d} {it['title']}{src}{zh}\n  {it['link']}")
                conn.execute(text("INSERT INTO bond_rating_news_seen(link, issuer, title) VALUES (:l,:i,:t) ON CONFLICT DO NOTHING"),
                             {"l": it["link"], "i": iss, "t": it["title"][:300]})
            blocks.append("\n".join(lines))
        # 未過濾但被排除的也記為已看過
        for iss, items in fresh.items():
            for it in items:
                conn.execute(text("INSERT INTO bond_rating_news_seen(link, issuer, title) VALUES (:l,:i,:t) ON CONFLICT DO NOTHING"),
                             {"l": it["link"], "i": iss, "t": it["title"][:300]})
    blocks.append("※ 新聞為 AI 初篩，請點連結確認原文；報價檔的評等以總行更新為準")
    return "\n\n".join(blocks)

def job_drive_cleanup():
    """每週日 03:00 清掉 Drive「龍蝦報告」資料夾裡 30 天前的舊檔（移到垃圾桶，可救回）"""
    now = datetime.now(TZ_TAIPEI_PYTZ)
    write_job_log("Drive清理", "started", now.strftime('%Y-%m-%d %H:%M'))
    try:
        from pdf_generator import cleanup_drive_folder
        days = int(os.getenv("DRIVE_KEEP_DAYS", "30"))
        deleted, checked, _ = cleanup_drive_folder(days=days)
        write_job_log("Drive清理", "success", f"檢查{checked}個，刪除{deleted}個({days}天前)")
        user_id = os.getenv("LINE_USER_ID", "")
        if user_id and deleted:
            line_bot_api.push_message(user_id, TextSendMessage(
                text=f"🧹 Drive 定期清理：已將 {deleted} 個超過 {days} 天的舊報告移到垃圾桶（共檢查 {checked} 個）"))
    except Exception as e:
        write_job_log("Drive清理", "error", str(e))
        print(f"[DriveCleanup ERROR] {e}")

def job_econ_watch():
    """每 10 分鐘檢查一次重要經濟數據/央行會議是否剛公布(僅在合理時段內實際查詢)"""
    try:
        from econ_watch import check_econ_events
    except Exception as e:
        print(f"[EconWatch] 模組載入失敗: {e}")
        return
    _t = load_targets() or {}
    user_id = os.getenv("LINE_USER_ID", "")
    recipients = [t for t in dict.fromkeys([user_id, _t.get("bond", "")] + list(_t.get("bond_subscribers", []))) if t]
    if not recipients:
        return

    def _push(msg):
        for rid in recipients:
            try:
                push_long_message(line_bot_api, rid, msg)
            except Exception as e:
                print(f"[EconWatch] push fail {rid[:8]}...: {e}")

    try:
        hit = check_econ_events(engine, sql_text, claude_client, _push)
        if hit:
            write_job_log("經濟數據監控", "success", f"推播 {hit} 項")
    except Exception as e:
        print(f"[EconWatch ERROR] {e}")
        print(_traceback.format_exc()[:500])
        write_job_log("經濟數據監控", "error", str(e))


def job_bond_rating_news():
    """每天 07:00 外部信評新聞掃描，推給配息雷達同一批對象"""
    now = datetime.now(TZ_TAIPEI_PYTZ)
    write_job_log("信評新聞雷達", "started", now.strftime('%Y-%m-%d %H:%M'))
    user_id = os.getenv("LINE_USER_ID", "")
    _t = load_targets() or {}
    recipients = [t for t in dict.fromkeys([user_id, _t.get("bond", "")] + list(_t.get("bond_subscribers", []))) if t]
    try:
        msg = run_rating_news_check(days=2)
        if not msg:
            write_job_log("信評新聞雷達", "success", "無新評等新聞")
            return
        for rid in recipients:
            try:
                push_long_message(line_bot_api, rid, msg)
            except Exception as e:
                print(f"[RatingNews] push fail {e}")
        write_job_log("信評新聞雷達", "success", f"推播 {len(recipients)} 個對象")
    except Exception as e:
        write_job_log("信評新聞雷達", "error", str(e))

def get_issuer_ticker(issuer):
    """發行機構 → (母集團/上市主體名, 股票代碼)。先查快取,缺的問 AI。抓不到回 (None, None)"""
    with engine.begin() as conn:
        row = conn.execute(text("SELECT parent, ticker FROM bond_issuer_ticker WHERE issuer=:i"), {"i": issuer}).fetchone()
    if row:
        return row[0], row[1]
    prompt = ("債券發行機構「" + issuer + "」(台灣銀行架上海外債的發行人)。請判斷其財報應追溯到哪一家上市公司:"
              "若發行主體是子公司/SPV(如美林私人→美國銀行、高盛金融國際→高盛集團),給母集團;若本身就是上市公司,給它自己。"
              "只回傳 JSON:{\"parent\": 上市主體中文名, \"ticker\": 美股或主要市場代碼(如 AAPL、BAC、8306.T),無法對應上市公司則 ticker 給 null}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=300)
    parent = (got or {}).get("parent") or None
    ticker = (got or {}).get("ticker") or None
    if isinstance(ticker, str) and ticker.lower() in ("null", "none", ""):
        ticker = None
    with engine.begin() as conn:
        conn.execute(text("""INSERT INTO bond_issuer_ticker(issuer, parent, ticker) VALUES (:i,:p,:t)
                             ON CONFLICT (issuer) DO UPDATE SET parent=EXCLUDED.parent, ticker=EXCLUDED.ticker, updated_at=NOW()"""),
                     {"i": issuer, "p": parent, "t": ticker})
    return parent, ticker

def get_issuer_profile(issuer, parent=""):
    """
    /sheet 用的加長版簡介(150~200字):營運概況、主要營收分佈、產業地位、信用重點。
    先查快取,缺的問 AI(依模型知識,禁止捏造精確比例)。
    """
    with engine.begin() as conn:
        row = conn.execute(text("SELECT profile FROM bond_issuer_profile WHERE issuer=:i"), {"i": issuer}).fetchone()
    if row:
        return row[0]
    subject = issuer + (f"(母集團為{parent})" if parent and parent != issuer else "")
    prompt = ("你是銀行固定收益研究員。請為債券發行機構「" + subject + "」寫一段 150~200 字的繁體中文簡介,"
              "依序涵蓋:1.主要業務與營運概況 2.主要營收分佈(以部門或地區描述;依你的知識用「主要來自…、其次為…」的寫法,"
              "不確定的比例不要寫出精確數字) 3.產業地位 4.信用體質重點。"
              "語氣專業中性,不做投資建議。只回傳 JSON:{\"profile\": 簡介文字}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=800)
    profile = str((got or {}).get("profile") or "").strip()
    if profile:
        with engine.begin() as conn:
            conn.execute(text("""INSERT INTO bond_issuer_profile(issuer, profile) VALUES (:i,:p)
                                 ON CONFLICT (issuer) DO UPDATE SET profile=EXCLUDED.profile, updated_at=NOW()"""),
                         {"i": issuer, "p": profile})
        return profile
    # 退回短版簡介
    return (get_issuer_intros([issuer]) or {}).get(issuer, "")

def get_issuer_bullets(issuer, parent=""):
    """把發行機構簡介改成 3~4 條 bullet(仿商品文宣版型),存快取"""
    with engine.begin() as conn:
        row = conn.execute(text("SELECT profile FROM bond_issuer_profile WHERE issuer=:i"),
                           {"i": issuer + "|bullets"}).fetchone()
    if row:
        try:
            return json.loads(row[0])
        except Exception:
            pass
    subject = issuer + (f"(母集團為{parent})" if parent and parent != issuer else "")
    prompt = ("你是銀行固定收益研究員。請為債券發行機構「" + subject + "」寫 4 條簡介重點,"
              "每條 30~50 字繁體中文,依序為:\n"
              "1. 成立背景、總部所在地、在該產業/國家的規模地位\n"
              "2. 主要業務組成(用文字描述部門或業務線,不要編造精確的營收百分比)\n"
              "3. 市場地位或競爭優勢的具體事實\n"
              "4. 信用體質重點:財務結構特徵與主要風險\n"
              "語氣專業中性,不做投資建議,不確定的數字寧可不寫。"
              "只回傳 JSON:{\"bullets\": [第1條, 第2條, 第3條, 第4條]}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=800)
    bullets = (got or {}).get("bullets") or []
    bullets = [str(b).strip() for b in bullets if str(b).strip()][:5]
    if bullets:
        with engine.begin() as conn:
            conn.execute(text("""INSERT INTO bond_issuer_profile(issuer, profile) VALUES (:i,:p)
                                 ON CONFLICT (issuer) DO UPDATE SET profile=EXCLUDED.profile, updated_at=NOW()"""),
                         {"i": issuer + "|bullets", "p": json.dumps(bullets, ensure_ascii=False)})
    return bullets

def get_peer_caps(issuer, parent=""):
    """同業市值(給長條圖用):AI 給同業代碼,市值一律用 yfinance 實抓"""
    subject = parent or issuer
    prompt = ("公司「" + subject + "」的主要同業(相同產業、相同國家或區域優先)請給 2~3 家。"
              "只回傳 JSON:{\"peers\":[{\"name\":中文名,\"ticker\":代碼}]}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=300)
    out = []
    from bond_sheet import get_financials
    for p_ in ((got or {}).get("peers") or [])[:3]:
        tk = str(p_.get("ticker") or "").strip()
        if not tk:
            continue
        pf = get_financials(tk)
        if pf and pf.get("market_cap"):
            out.append({"name": str(p_.get("name") or tk), "ticker": tk, "market_cap": pf["market_cap"]})
    return out

def get_fin_comment(issuer, fin):
    """對五大財務比率做 2~3 句的信用角度 AI 解讀"""
    if not fin:
        return ""
    facts = {k: fin.get(k) for k in ("market_cap", "eps", "roe", "debt_ratio", "net_debt_ebitda", "currency")}
    prompt = ("以下是債券發行機構「" + issuer + "」(或其母集團)的財務指標:" + json.dumps(facts, ensure_ascii=False) +
              "。請從『債權人/信用分析』角度用 2~3 句繁體中文解讀(60~120字):獲利與規模代表什麼、"
              "槓桿(負債比、淨負債/EBITDA)在該產業屬於什麼水準、對償債能力的意義。"
              "語氣中性,不評價股票、不做買賣建議、不用果決斷言。只回傳 JSON:{\"comment\": 解讀文字}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=500)
    return str((got or {}).get("comment") or "").strip()

def get_peer_comparison(issuer, parent, fin):
    """同業比較:AI 給 2~3 家同業代碼,實際數字用 yfinance 抓,避免捏造"""
    if not fin:
        return ""
    subject = parent or issuer
    prompt = ("公司「" + subject + "」的主要同業(相同產業、規模相近的上市公司)請給 2~3 家。"
              "只回傳 JSON:{\"peers\":[{\"name\":中文名,\"ticker\":代碼}]}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=300)
    peers = (got or {}).get("peers") or []
    from bond_sheet import get_financials
    out = []
    for p_ in peers[:3]:
        tk = str(p_.get("ticker") or "").strip()
        nm = str(p_.get("name") or tk).strip()
        if not tk:
            continue
        pf = get_financials(tk)
        if not pf:
            continue
        roe = f"ROE {pf['roe']*100:.0f}%" if pf.get("roe") is not None else "ROE -"
        dr = f"負債比 {pf['debt_ratio']:.0f}%" if pf.get("debt_ratio") is not None else "負債比 -"
        nd = f"淨負債/EBITDA {pf['net_debt_ebitda']}x" if pf.get("net_debt_ebitda") is not None else ""
        out.append(f"{nm}（{tk}）{roe}、{dr}" + (f"、{nd}" if nd else ""))
    if not out:
        return ""
    return "；".join(out) + "（資料來源：公開財報，僅供比較參考）"

def get_rating_outlook(issuer):
    """評等展望/最近一次評等動作:用近 180 天的信評新聞讓 AI 歸納,查不到回空字串"""
    try:
        from bond_rating_news import fetch_rating_news
        with engine.begin() as conn:
            row = conn.execute(text("SELECT en_name FROM bond_rating_watch WHERE issuer=:i"), {"i": issuer}).fetchone()
        en = row[0] if row and row[0] else ""
        items = fetch_rating_news(en, zh_name=issuer, days=180)[:8]
        if not items:
            return ""
        titles = [{"t": it["title"], "d": it["published"].strftime("%Y/%m") if it["published"] else ""} for it in items]
        prompt = ("以下是關於「" + issuer + "」的信評相關新聞標題。請歸納出目前的評等展望與最近一次評等動作,"
                  "格式如:『展望：穩定（S&P，2026/03）』或『2026/05 穆迪確認 A2 評等，展望穩定』。"
                  "只能根據標題內容,無法判斷就回空字串。只回傳 JSON:{\"outlook\": 文字}\n\n" + json.dumps(titles, ensure_ascii=False))
        got, _, _ = llm_json_fallback(prompt, max_tokens=300)
        return str((got or {}).get("outlook") or "").strip()
    except Exception as e:
        print(f"[BondSheet outlook] {e}")
        return ""

def get_charts_comment(issuer, q):
    """對近五季四張圖做 AI 解讀:正向、建設性,但不得捏造或美化負面數字"""
    if not q:
        return ""
    data = {k: q.get(k) for k in ("labels", "revenue", "op_income", "ocf", "fcf",
                                  "debt", "cash", "debt_ebitda", "int_cover")}
    prompt = ("以下是「" + issuer + "」近五季的財報數據(單位:億,依序由舊到新):"
              + json.dumps(data, ensure_ascii=False) +
              "\n請以債券投資人(債權人)角度,用 100~150 字繁體中文解讀這四組數據:"
              "營運表現(營收/營業淨利)、現金流(營業現金流/自由現金流)、債務與現金、信用比率(債務EBITDA、利息保障倍數)。\n"
              "要求:\n"
              "1. 語氣正向且具建設性,優先點出結構性優勢與穩定性(例如營收規模穩定、現金流回升、利息保障充足)。\n"
              "2. 但不得捏造或美化:若某季出現下滑或負值(例如自由現金流為負),要據實說明並補充合理的解釋角度(例如季節性、資本支出集中),不可略過不提。\n"
              "3. 只描述數據本身與對償債能力的意義,不做投資建議、不預測股價、不用果決斷言。\n"
              "只回傳 JSON:{\"comment\": 解讀文字}")
    got, _, _ = llm_json_fallback(prompt, max_tokens=600)
    return str((got or {}).get("comment") or "").strip()

def build_issuer_hist_map(isins, days=30):
    """{isin: '｜近30日±x.x%'} 由 bond_price_history 計算"""
    out = {}
    try:
        with engine.begin() as conn:
            for isin in isins:
                rows = conn.execute(text("""SELECT snap_date, offer FROM bond_price_history
                                           WHERE isin=:i AND offer IS NOT NULL AND snap_date >= CURRENT_DATE - :d * INTERVAL '1 day'
                                           ORDER BY snap_date"""), {"i": isin, "d": days}).fetchall()
                if len(rows) >= 2 and rows[0][1]:
                    chg = (rows[-1][1] - rows[0][1]) / rows[0][1] * 100
                    out[isin] = f"｜近30日{chg:+.1f}%"
    except Exception as e:
        print(f"[BondSheet hist] {e}")
    return out

def job_bond_coupon_radar():
    """每天 06:45 推播海外債配息雷達給 Albert（週一～週五）"""
    now = datetime.now(TZ_TAIPEI_PYTZ)
    write_job_log("海外債配息雷達", "started", now.strftime('%Y-%m-%d %H:%M'))
    user_id = os.getenv("LINE_USER_ID", "")
    _t = load_targets() or {}
    bond_target = _t.get("bond", "")
    subscribers = _t.get("bond_subscribers", [])
    # 推播對象：Albert 個人 ＋ 海外債群組（/coupon settarget）＋ 個人訂閱者（/coupon subscribe）；去重
    recipients = [t for t in dict.fromkeys([user_id, bond_target] + list(subscribers)) if t]
    if not recipients:
        write_job_log("海外債配息雷達", "skipped", "缺少推播對象（LINE_USER_ID / /coupon settarget）")
        return
    try:
        if not _BOND_RADAR_OK:
            raise RuntimeError("bond_coupon_alert 模組未載入")
        if not BOND_PRICE_FILE.exists():
            if user_id:
                line_bot_api.push_message(user_id, TextSendMessage(text="📭 配息雷達：還沒有海外債報價檔，請把 Bond_Pricing Excel 傳給我。"))
            write_job_log("海外債配息雷達", "skipped", "無報價檔")
            return
        msg = _bond_build_alert(str(BOND_PRICE_FILE), today=now.date(), lookahead=14, days_ahead=3)
        mtime = datetime.fromtimestamp(BOND_PRICE_FILE.stat().st_mtime, TZ_TAIPEI_PYTZ)
        age_days = (now.date() - mtime.date()).days
        msg += f"\n📎 報價檔更新於 {mtime:%m/%d %H:%M}｜/coupon all 看全部｜/coupon table 出Excel"
        if age_days >= 3:
            msg += f"（已 {age_days} 天未更新，記得丟新檔）"
        ok_n = 0
        for rid in recipients:
            try:
                push_long_message(line_bot_api, rid, msg)
                ok_n += 1
            except Exception as e:
                print(f"[BondRadar] push to {rid[:8]}… failed: {e}")
        write_job_log("海外債配息雷達", "success", f"推播 {ok_n}/{len(recipients)} 個對象")
        try:
            check_bond_alerts(line_bot_api, source="daily")
        except Exception as e:
            print(f"[BondAlert daily ERROR] {e}")
    except Exception as e:
        write_job_log("海外債配息雷達", "error", str(e))
        try:
            if user_id:
                line_bot_api.push_message(user_id, TextSendMessage(text=f"❌ 配息雷達失敗：{str(e)[:200]}"))
        except Exception:
            pass

def start_scheduler():
    scheduler = BackgroundScheduler(timezone=TZ_TAIPEI_PYTZ)
    scheduler.add_job(job_bond_coupon_radar, CronTrigger(day_of_week="mon-fri", hour=6, minute=25, timezone=TZ_TAIPEI_PYTZ), id="bond_coupon_radar", name="海外債配息雷達")
    scheduler.add_job(job_bond_rating_news, CronTrigger(day_of_week="mon-fri", hour=6, minute=50, timezone=TZ_TAIPEI_PYTZ), id="bond_rating_news", name="信評新聞雷達")
    scheduler.add_job(job_drive_cleanup, CronTrigger(day_of_week="sun", hour=3, minute=0, timezone=TZ_TAIPEI_PYTZ), id="drive_cleanup", name="Drive清理")
    scheduler.add_job(job_econ_watch, IntervalTrigger(minutes=10, timezone=TZ_TAIPEI_PYTZ), id="econ_watch", name="經濟數據監控")
    scheduler.add_job(job_daily_report, CronTrigger(day_of_week="mon-sat", hour=6, minute=40, timezone=TZ_TAIPEI_PYTZ), id="daily_report", name="財經日報")
    scheduler.add_job(job_bond_daily_report, CronTrigger(day_of_week="mon-sat", hour=6, minute=30, timezone=TZ_TAIPEI_PYTZ), id="bond_daily_report", name="債券日報")
    scheduler.add_job(job_auto_tracking, CronTrigger(day_of_week="mon-sat", hour=7, minute=0, timezone=TZ_TAIPEI_PYTZ), id="auto_tracking", name="ELN自動追蹤")
    scheduler.add_job(job_alert_monitor, IntervalTrigger(minutes=15), id="alert_monitor", name="價格警示")
    scheduler.add_job(job_spending_report, CronTrigger(hour=9, minute=0, timezone=TZ_TAIPEI_PYTZ), id="spending_report", name="月度消費明細")
    scheduler.add_job(job_mail_monitor, IntervalTrigger(minutes=15, start_date=datetime.now(TZ_TAIPEI_PYTZ).replace(second=0, microsecond=0) + timedelta(minutes=5)), id="mail_monitor", name="郵件監控")
    scheduler.add_job(job_article_reminder, CronTrigger(day_of_week="mon", hour=8, minute=30, timezone=TZ_TAIPEI_PYTZ), id="article_reminder", name="未讀文章提醒")
    scheduler.add_job(job_fund_nav_update, CronTrigger(day_of_week="mon-fri", hour=9, minute=0, timezone=TZ_TAIPEI_PYTZ), id="fund_nav_update", name="基金淨值更新")
    scheduler.start()
    print("[Scheduler] 排程啟動完成 ✅")
    return scheduler

_scheduler = start_scheduler()

# ==============================
# 客戶建議書 PPT 生成 /generate-ppt
# ==============================
import subprocess
import tempfile
from fastapi.responses import FileResponse

@app.post("/generate-ppt")
async def generate_ppt_endpoint(request: Request):
    """
    接收來自 Streamlit 已組好的完整 JSON payload，
    直接呼叫 generate_ppt.js 生成 PPTX 並回傳。
    所有 Claude API 呼叫（市場背景、基金介紹、AI解讀）
    已移至 Streamlit 端處理，eln-bot 只負責 Node.js 生成。
    """
    try:
        data = await request.json()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            output_path = tmp.name

        js_path = Path(__file__).parent / "generate_ppt.js"
        result = subprocess.run(
            ["node", str(js_path), "--data", json.dumps(data, ensure_ascii=False), "--output", output_path],
            capture_output=True, text=True, timeout=120
        )

        if not result.stdout.strip().startswith("OK:"):
            raise Exception(f"PPT 生成失敗：{result.stdout[:500]} {result.stderr[:500]}")

        client_name = data.get("client_name", "客戶")
        report_date = data.get("report_date", datetime.now(TZ_TAIPEI).strftime("%Y%m%d"))
        filename = f"{client_name}_投資組合建議書_{report_date}.pptx"

        return FileResponse(
            path=output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename,
            background=None
        )

    except Exception as e:
        print(f"[generate-ppt] 錯誤：{e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPT 生成失敗：{str(e)[:300]}")

# ==============================
# 知識庫路由 /kb
# ==============================
@app.get("/kb")
async def kb_home():
    with open("static/kb/index.html", "r", encoding="utf-8") as f:
        return HTMLResponse(f.read())

@app.post("/kb/upload-table")
async def kb_upload_table(file: UploadFile = File(...)):
    try:
        file_bytes = await file.read()
        result = knowledge.process_and_index_file(file.filename, file_bytes, as_table=True)
        return {"success": True, **result, "message": f"✅ 已登記為表格直查圖片，查詢時自動用 Vision 看圖回答"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"處理失敗：{str(e)}")

@app.post("/kb/upload")
async def kb_upload(file: UploadFile = File(...)):
    try:
        file_bytes = await file.read()
        result = knowledge.process_and_index_file(file.filename, file_bytes)
        return {"success": True, **result, "message": f"✅ 成功處理 {result['pages']} 頁，建立 {result['chunks']} 個索引"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"處理失敗：{str(e)}")

@app.post("/kb/ask")
async def kb_ask(question: str = Form(...)):
    try:
        return knowledge.query_knowledge(question)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"查詢失敗：{str(e)}")

@app.get("/kb/page-image/{doc_id}/{page_num}")
async def kb_page_image(doc_id: str, page_num: int):
    try:
        img_data = knowledge.get_page_image_base64(doc_id, page_num)
        return {"image_base64": img_data}
    except:
        raise HTTPException(status_code=404, detail="頁面圖片不存在")

@app.post("/kb/image-to-table")
async def kb_image_to_table(file: UploadFile = File(...)):
    try:
        file_bytes = await file.read()
        suffix = Path(file.filename).suffix.lower()
        media_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}
        media_type = media_map.get(suffix, "image/png")
        img_data = _base64.b64encode(file_bytes).decode("utf-8")
        response = claude_client.messages.create(
            model="claude-sonnet-4-6", max_tokens=3000,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": img_data}},
                {"type": "text", "text": (
                    "請把這張圖片裡的所有內容完整轉換成純文字格式，規則如下：\n\n"
                    "1. 如果有表格，用 Markdown 表格格式輸出\n"
                    "2. 如果有條列式文字，保持原本的條列結構\n"
                    "3. 標題和小標題要保留\n"
                    "4. 所有數字、時間、百分比、金額一個都不能少\n"
                    "5. 不需要說明你在做什麼，直接輸出轉換後的文字\n\n直接輸出結果："
                )}
            ]}]
        )
        extracted_text = response.content[0].text
        return {"success": True, "text": extracted_text, "filename": file.filename}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"轉換失敗：{str(e)}")

@app.post("/kb/save-text")
async def kb_save_text(request: Request):
    try:
        body = await request.json()
        text = body.get("text", "").strip()
        filename = body.get("filename", "手動輸入.txt")
        if not text:
            raise HTTPException(status_code=400, detail="文字不能為空")
        file_bytes = text.encode("utf-8")
        result = knowledge.process_and_index_file(filename, file_bytes)
        return {"success": True, **result, "message": f"✅ 成功建立 {result['chunks']} 個索引"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"儲存失敗：{str(e)}")

@app.get("/kb/files")
async def kb_files_page():
    with open("static/kb/files.html", "r", encoding="utf-8") as f:
        return HTMLResponse(f.read())

@app.get("/kb/files-data")
async def kb_files_data():
    return {"files": knowledge.list_files_detail()}

@app.get("/kb/documents")
async def kb_documents():
    return {"documents": knowledge.list_documents()}

@app.delete("/kb/document/{doc_id}")
async def kb_delete(doc_id: str):
    knowledge.delete_document(doc_id)
    return {"success": True}

app.mount("/kb/static", StaticFiles(directory="static/kb"), name="kb-static")
