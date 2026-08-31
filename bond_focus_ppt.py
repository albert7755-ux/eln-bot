# -*- coding: utf-8 -*-
"""
bond_focus_ppt.py — 「債市每日聚焦 / 富邦好債報」PPTX 產生器（直式 A4，兩頁）
使用 python-pptx（與 /report 相同的套件，不需 Node）
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import tempfile

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x1F, 0x8A, 0xC0)
DEEP = RGBColor(0x1B, 0x7F, 0xA8)
TEAL = RGBColor(0x16, 0x9B, 0x9B)
GRAY = RGBColor(0x59, 0x59, 0x59)
RED = RGBColor(0xC0, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PEACH = RGBColor(0xFD, 0xF2, 0xEC)
F = "Microsoft JhengHei"

PAGE_W, PAGE_H = Cm(21.0), Cm(29.7)
M = Cm(1.15)
W = PAGE_W - 2 * M


def _txt(slide, x, y, w, h, runs, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT,
         line=1.35, anchor=MSO_ANCHOR.TOP, space_after=4):
    """runs: str 或 [(文字, {bold, color, size}), ...]；每個 tuple 為一個段落"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0.1)
    tf.margin_top = tf.margin_bottom = 0
    items = [(runs, {})] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        text, opt = (item if isinstance(item, tuple) else (item, {}))
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        para.space_after = Pt(opt.get("space_after", space_after))
        for j, seg in enumerate(text if isinstance(text, list) else [(text, {})]):
            s_text, s_opt = (seg if isinstance(seg, tuple) else (seg, {}))
            r = para.add_run()
            r.text = s_text
            r.font.name = F
            r.font.size = Pt(s_opt.get("size", opt.get("size", size)))
            r.font.bold = s_opt.get("bold", opt.get("bold", bold))
            r.font.color.rgb = s_opt.get("color", opt.get("color", color))
    return tb


def _rect(slide, x, y, w, h, fill=None, line_color=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _line(slide, x, y, w, color=BLUE, width=1.6):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def _table(slide, x, y, w, rows_data, col_w, row_h=Cm(1.0), font=12,
           head_fill=None, body_fill=None, head_color=NAVY):
    rows, cols = len(rows_data), len(rows_data[0])
    tb = slide.shapes.add_table(rows, cols, x, y, w, row_h * rows).table
    for j, cw in enumerate(col_w):
        tb.columns[j].width = cw
    for i, row in enumerate(rows_data):
        tb.rows[i].height = row_h
        for j, cell in enumerate(row):
            text, opt = (cell if isinstance(cell, tuple) else (cell, {}))
            c = tb.cell(i, j)
            c.text = str(text)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Cm(0.08)
            c.margin_top = c.margin_bottom = 0
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = F
                r.font.size = Pt(opt.get("size", font))
                r.font.bold = opt.get("bold", i == 0)
                r.font.color.rgb = opt.get("color", head_color if i == 0 else DARK)
            c.fill.solid()
            c.fill.fore_color.rgb = (head_fill or RGBColor(0xEA, 0xF3, 0xFA)) if i == 0 else (body_fill or WHITE)
    return tb


def _cjk_font():
    """找中文字型(與 bond_sheet 共用邏輯)"""
    base = os.path.dirname(os.path.abspath(__file__))
    import glob
    cands = [os.getenv("BOND_SHEET_FONT", "")]
    for n in ("NotoSansTC-Regular.ttf", "NotoSansTC.ttf", "msjh.ttf"):
        cands += [os.path.join(base, "fonts", n), os.path.join(base, n)]
    cands += ["/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
              "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"]
    cands += sorted(glob.glob(os.path.join(base, "fonts", "*.tt*")))
    for c in cands:
        try:
            if c and os.path.exists(c) and os.path.getsize(c) > 50 * 1024:
                return c
        except Exception:
            pass
    return None


def _donut_png(rm):
    """用 matplotlib 畫營收結構甜甜圈,回傳暫存 PNG 路徑(失敗回 None)"""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import font_manager
        fp = _cjk_font()
        PROP = None
        if fp:
            try:
                PROP = font_manager.FontProperties(fname=fp)
                font_manager.fontManager.addfont(fp)
                matplotlib.rcParams["font.family"] = PROP.get_name()
            except Exception:
                PROP = None
        labels = list(rm.get("labels") or [])
        values = [float(v) for v in rm.get("values") or []]
        if not values:
            return None
        colors = ["#169B9B", "#1F8AC0", "#BFBFBF", "#7F7F7F", "#4472C4"]
        fig, ax = plt.subplots(figsize=(4.2, 3.4), dpi=170)
        total = sum(values)
        wedges, _texts, autotexts = ax.pie(
            values, startangle=90, counterclock=False,
            colors=colors[:len(values)], wedgeprops=dict(width=0.42, edgecolor="white"),
            autopct=lambda pct: f"{pct:.0f}%" if pct >= 5 else "",
            pctdistance=0.78, textprops={"fontsize": 13, "color": "white",
                                         "fontproperties": PROP, "weight": "bold"})
        ax.legend(wedges, labels, loc="lower center", bbox_to_anchor=(0.5, -0.22),
                  ncol=2, frameon=False, fontsize=9,
                  prop=(PROP.copy() if PROP else None))
        if PROP:
            for lg in ax.get_legend().get_texts():
                lg.set_fontproperties(PROP)
                lg.set_fontsize(9)
        ax.set(aspect="equal")
        fig.tight_layout(pad=0.3)
        f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(f.name, bbox_inches="tight", facecolor="white", transparent=False)
        plt.close(fig)
        return f.name
    except Exception as e:
        print(f"[BondFocus] donut fail: {e}")
        return None


def _magnifier_png(size_px=220, color="#1F8AC0"):
    """畫一個放大鏡圖示,取代在 PDF/PPT 無法顯示的 🔍 emoji"""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.patches import Circle, FancyBboxPatch
        fig, ax = plt.subplots(figsize=(1.6, 1.6), dpi=size_px / 1.6)
        ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")
        ax.add_patch(Circle((4.3, 6.0), 2.6, fill=False, lw=6.5, color=color))
        ax.add_patch(Circle((4.3, 6.0), 2.1, facecolor="#DCEEF8", edgecolor="none", alpha=0.85))
        ax.plot([6.2, 8.6], [4.0, 1.7], lw=7.5, color=color, solid_capstyle="round")
        fig.tight_layout(pad=0)
        f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(f.name, transparent=True, bbox_inches="tight", pad_inches=0.02)
        plt.close(fig)
        return f.name
    except Exception as e:
        print(f"[BondFocus] magnifier fail: {e}")
        return None


def _bar_png(labels, series, title=""):
    """
    備援圖表:長條圖(可 1~2 組數列)。series = [(名稱, [值...]), ...]
    回傳暫存 PNG 路徑,失敗回 None。
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        from matplotlib import font_manager
        fp = _cjk_font()
        PROP = None
        if fp:
            try:
                PROP = font_manager.FontProperties(fname=fp)
                font_manager.fontManager.addfont(fp)
                matplotlib.rcParams["font.family"] = PROP.get_name()
            except Exception:
                PROP = None
        matplotlib.rcParams["axes.unicode_minus"] = False
        colors_ = ["#169B9B", "#1F8AC0"]
        fig, ax = plt.subplots(figsize=(4.6, 3.2), dpi=170)
        x = np.arange(len(labels))
        n = max(1, len(series))
        w = 0.7 / n
        for i, (name, vals) in enumerate(series):
            v = [0 if x_ is None else x_ for x_ in vals]
            bars = ax.bar(x + (i - (n - 1) / 2) * w, v, w, label=name, color=colors_[i % 2])
            for b_, val in zip(bars, v):
                if val:
                    ax.annotate(f"{val:,.0f}" if abs(val) >= 100 else f"{val:,.1f}",
                                (b_.get_x() + b_.get_width() / 2, val), ha="center",
                                va="bottom", fontsize=8, fontproperties=PROP)
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=9, fontproperties=PROP)
        ax.tick_params(axis="y", labelsize=8)
        ax.spines[["top", "right"]].set_visible(False)
        if len(series) > 1:
            lg = ax.legend(fontsize=8, frameon=False, ncol=2, loc="lower center",
                           bbox_to_anchor=(0.5, 1.005),
                           prop=(PROP.copy() if PROP else None))
            if PROP:
                for tx in lg.get_texts():
                    tx.set_fontproperties(PROP)
                    tx.set_fontsize(8)
        if title:
            # 有圖例時標題再往上,避免與圖例重疊
            ax.set_title(title, fontsize=10.5, color="#0B2A4A",
                         pad=(26 if len(series) > 1 else 8), fontproperties=PROP)
        fig.tight_layout(pad=0.4)
        f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(f.name, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return f.name
    except Exception as e:
        print(f"[BondFocus] bar png fail: {e}")
        return None


def _chart_for_revenue_block(D):
    """
    營收結構區塊要放哪張圖:
    1) AI 有給部門別營收 → 甜甜圈
    2) 否則用近五季營收/營業淨利 → 長條圖
    3) 再否則用架上債券到期分布 → 長條圖
    回傳 (png_path, 說明文字) 或 (None, "")
    """
    rm = D.get("revenue_mix") or {}
    if rm.get("values"):
        return _donut_png(rm), rm.get("unit_note", "")
    q = D.get("quarterly") or {}
    if q.get("labels") and any(v is not None for v in (q.get("revenue") or [])):
        series = [("營業收入", q.get("revenue") or [])]
        if any(v is not None for v in (q.get("op_income") or [])):
            series.append(("營業淨利", q["op_income"]))
        return _bar_png(q["labels"], series, "近五季營收與獲利"), "單位：該公司報表幣別之億元"
    md = D.get("maturity_dist") or {}
    if md.get("labels"):
        return (_bar_png(md["labels"], [("檔數", md["values"])], "架上債券到期分布"),
                "本行架上該發行機構之債券檔數")
    return None, ""


OUTLOOK_ZH = {
    "stable": "穩定", "positive": "正向", "negative": "負向",
    "developing": "發展中", "watch negative": "負向觀察", "watch positive": "正向觀察",
    "rating watch negative": "負向觀察", "creditwatch negative": "負向觀察",
    "under review": "評等審查中", "n/a": "--", "na": "--", "none": "--", "": "--",
}


def _zh_outlook(v):
    """評等展望統一中文化(AI 有時回 Stable/Positive)"""
    t = str(v or "").strip()
    return OUTLOOK_ZH.get(t.lower(), t or "--")


def _revenue_block_title(D):
    """區塊標題依實際畫的圖調整"""
    rm = D.get("revenue_mix") or {}
    if rm.get("values"):
        return "營收結構"
    q = D.get("quarterly") or {}
    if q.get("labels") and any(v is not None for v in (q.get("revenue") or [])):
        return "營運趨勢"
    if (D.get("maturity_dist") or {}).get("labels"):
        return "架上債券分布"
    return "營收結構"


def build_focus_pptx(out_path, D):
    prs = Presentation()
    prs.slide_width, prs.slide_height = PAGE_W, PAGE_H
    blank = prs.slide_layouts[6]

    # ================= P1 債市每日聚焦 =================
    s = prs.slides.add_slide(blank)
    mag = _magnifier_png()
    if mag:
        s.shapes.add_picture(mag, M, Cm(0.62), height=Cm(1.55))
        try:
            os.remove(mag)
        except Exception:
            pass
    _txt(s, M + Cm(1.85), Cm(0.7), Cm(12), Cm(1.5), "債市每日聚焦", size=30, bold=True, color=NAVY)
    bar = s.shapes.add_shape(1, M, Cm(2.5), W, Cm(0.78))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
    _txt(s, M, Cm(2.58), W - Cm(0.3), Cm(0.7), D.get("date_str", ""), size=13, bold=True,
         color=WHITE, align=PP_ALIGN.RIGHT)

    _txt(s, M, Cm(3.65), Cm(6), Cm(1.0), "焦點新聞", size=20, bold=True, color=BLUE)
    _line(s, M, Cm(4.65), Cm(5.6))
    _txt(s, M, Cm(5.05), W, Cm(1.8), D.get("headline", ""), size=22, bold=True, color=DARK, line=1.25)

    news = [str(x) for x in (D.get("news_bullets") or [])]
    # 依字數自動縮放:每行約 26 字(15.5pt)、行高約 0.78cm,避免長文壓到下方區塊
    n_size = 15.5
    n_chars = sum(len(x) for x in news)
    if n_chars > 300:
        n_size = 13.0
    elif n_chars > 230:
        n_size = 14.0
    per_line = int(26 * 15.5 / n_size)
    n_lines = sum(max(1, -(-len(x) // per_line)) for x in news)
    news_h = Cm(0.052 * n_size * 1.5) * n_lines + Cm(0.25) * max(0, len(news) - 1)
    news_h = min(news_h, Cm(8.6))
    bullets = [([("・", {"color": BLUE, "bold": True}), (t, {})], {"space_after": 7})
               for t in news]
    _txt(s, M, Cm(7.1), W, news_h, bullets, size=n_size, line=1.5)

    # 焦點債券區塊位置隨新聞長度下移(最低不超過 Cm(18.4),確保表格與註記放得下)
    y_bond = min(Cm(18.4), max(Cm(15.0), Cm(7.1) + news_h + Cm(1.0)))
    _txt(s, M, y_bond, Cm(6), Cm(1.0), "焦點債券", size=20, bold=True, color=BLUE)
    _line(s, M, y_bond + Cm(1.0), Cm(5.6))
    if D.get("bond_tagline"):
        _txt(s, M + Cm(6), y_bond + Cm(0.1), W - Cm(6), Cm(0.9), D["bond_tagline"], size=14,
             bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    head = [("債券代碼", {}), ("債券名稱", {}), ("票面%", {}), ("YTM%", {}), ("到期日", {})]
    body = []
    for b in D.get("bonds", []):
        body.append([(b.get("code", "-"), {}),
                     (b.get("name", "-"), {"color": TEAL, "bold": True}),
                     (str(b.get("coupon", "-")), {"color": RED, "bold": True}),
                     (str(b.get("ytm", "-")), {}),
                     (b.get("maturity", "-"), {})])
    y_tbl = y_bond + Cm(1.5)
    _table(s, M, y_tbl, W, [head] + body,
           col_w=[Cm(4.5), Cm(4.3), Cm(2.3), Cm(3.2), Cm(4.4)],
           row_h=Cm(1.05), font=13.5, head_fill=PEACH, body_fill=PEACH)
    _txt(s, M, y_tbl + Cm(1.05) * (len(body) + 1) + Cm(0.2), W, Cm(0.8),
         "※ 報價與可承作與否以本行系統為準；商品條件依產品說明書。", size=10.5, color=GRAY)

    _txt(s, M, Cm(27.8), Cm(8), Cm(0.8), "僅限內部教育訓練使用", size=12.5, bold=True, color=RED)
    _txt(s, M, Cm(27.8), W, Cm(0.8), "台北富邦銀行", size=12.5, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    # ================= P2 富邦好債報 =================
    t = prs.slides.add_slide(blank)
    _rect(t, M, Cm(0.7), Cm(6.6), Cm(1.35), fill=DEEP)
    _txt(t, M, Cm(0.86), Cm(6.6), Cm(1.0), "富 邦 好 債 報", size=20, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _txt(t, M, Cm(2.45), W, Cm(1.5), D.get("issuer", ""), size=32, bold=True, align=PP_ALIGN.CENTER)
    _txt(t, M, Cm(3.95), W, Cm(0.8), D.get("issuer_en", ""), size=14, bold=True,
         color=GRAY, align=PP_ALIGN.CENTER)

    _rect(t, M, Cm(4.75), W, Cm(3.0), fill=WHITE, line_color=BLUE)
    _txt(t, M + Cm(0.3), Cm(4.95), W - Cm(0.6), Cm(2.6), D.get("intro", ""), size=12, line=1.45)

    # 營運概況 / 營收結構
    half = (W - Cm(0.5)) / 2
    _txt(t, M, Cm(8.0), half, Cm(1.0), "營運概況", size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _line(t, M + Cm(1.2), Cm(9.0), half - Cm(2.4))
    _rect(t, M, Cm(9.25), half, Cm(6.0), fill=WHITE, line_color=BLUE)
    ops = []
    for blk in (D.get("ops_blocks") or []):
        ops.append(([(str(blk[0]) + "：", {"bold": True, "color": BLUE})], {"space_after": 2}))
        ops.append(([(str(blk[1]), {})], {"space_after": 8}))
    _txt(t, M + Cm(0.25), Cm(9.45), half - Cm(0.5), Cm(5.6), ops, size=11, line=1.4)

    x2 = M + half + Cm(0.5)
    _txt(t, x2, Cm(8.0), half, Cm(1.0), _revenue_block_title(D), size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _line(t, x2 + Cm(1.2), Cm(9.0), half - Cm(2.4))
    _rect(t, x2, Cm(9.25), half, Cm(6.0), fill=WHITE, line_color=BLUE)
    png, png_note = _chart_for_revenue_block(D)
    if png:
        # 以高度為準置中,避免圖超出卡片
        try:
            from PIL import Image as _PIL
            iw, ih = _PIL.open(png).size
            max_w, max_h = half - Cm(0.3), Cm(5.05)
            img_w, img_h = max_w, int(max_w * ih / iw)
            if img_h > max_h:
                img_h, img_w = max_h, int(max_h * iw / ih)
            t.shapes.add_picture(png, x2 + (half - img_w) // 2,
                                 Cm(9.35) + (Cm(5.2) - img_h) // 2, width=img_w, height=img_h)
        except Exception:
            t.shapes.add_picture(png, x2 + Cm(0.4), Cm(9.4), height=Cm(4.9))
        _txt(t, x2, Cm(14.55), half, Cm(0.7), png_note, size=9, color=GRAY,
             align=PP_ALIGN.CENTER)
        try:
            os.remove(png)
        except Exception:
            pass
    else:
        _txt(t, x2, Cm(11.8), half, Cm(1.0), "（暫無可用圖表資料）", size=12, color=GRAY,
             align=PP_ALIGN.CENTER)

    # 信用評等
    _txt(t, M, Cm(15.6), W, Cm(1.0), "信用評等", size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _line(t, M + Cm(6.0), Cm(16.6), W - Cm(12.0))
    R = D.get("ratings") or {}
    rows2 = [[("信用評等", {}), ("穆迪", {}), ("標普", {}), ("惠譽", {})],
             [("長期評等", {"bold": True}), (R.get("moody") or "--", {}), (R.get("sp") or "--", {}), (R.get("fitch") or "--", {})],
             [("評等展望", {"bold": True}), (_zh_outlook(R.get("moody_outlook")), {}), (_zh_outlook(R.get("sp_outlook")), {}), (_zh_outlook(R.get("fitch_outlook")), {})],
             [("最近評等動作", {"bold": True}), (R.get("moody_date") or "--", {}), (R.get("sp_date") or "--", {}), (R.get("fitch_date") or "--", {})]]
    _table(t, M, Cm(16.95), W, rows2,
           col_w=[Cm(5.0), Cm(4.5), Cm(4.5), Cm(4.7)], row_h=Cm(0.95), font=12)

    # 信評公司評析（高度依內容自動調整）
    _txt(t, M, Cm(21.1), W, Cm(1.0), "信評公司評析", size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _line(t, M + Cm(6.0), Cm(22.1), W - Cm(12.0))
    ac = []
    for c in (D.get("agency_comments") or []):
        ac.append(([(str(c[0]) + "－", {"bold": True, "color": BLUE}), (str(c[1]), {})], {"space_after": 7}))
    chars = sum(len(str(c[0])) + len(str(c[1])) for c in (D.get("agency_comments") or []))
    lines = -(-chars // 40) + len(D.get("agency_comments") or [])
    ac_h = min(Cm(5.6), max(Cm(2.4), Cm(0.62) * lines + Cm(0.7)))
    _rect(t, M, Cm(22.5), W, ac_h, fill=WHITE, line_color=BLUE)
    _txt(t, M + Cm(0.25), Cm(22.7), W - Cm(0.5), ac_h - Cm(0.4), ac, size=11, line=1.4)

    y_note = Cm(22.5) + ac_h + Cm(0.15)
    _txt(t, M, y_note, W, Cm(0.6), D.get("source_note", ""), size=8.5, color=GRAY)
    _txt(t, M, y_note + Cm(0.55), W, Cm(1.0),
         "「本內容僅供參考且不構成要約或要約引誘。特定標的之商品風險、申購之條件、限制與費用及其他相關權利義務，"
         "應依產品說明暨投資風險預告書等相關文件為準。」", size=8.5, color=GRAY, line=1.25)

    _txt(t, M, Cm(28.3), Cm(8), Cm(0.8), "僅限內部教育訓練使用", size=12, bold=True, color=RED)
    _txt(t, M, Cm(28.3), W, Cm(0.8), "台北富邦銀行", size=12, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    prs.save(out_path)
    return out_path


# ================= PDF 版（reportlab 直接產，不依賴 Drive 轉檔）=================
def build_focus_pdf(out_path, D):
    """與 PPTX 相同版型的 PDF（直式 A4 兩頁）。成功回傳路徑。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, PageBreak, Image as RLImage)
    from reportlab.lib.styles import ParagraphStyle

    # 字型
    fp = _cjk_font()
    FN = "MSung-Light"
    if fp:
        try:
            pdfmetrics.registerFont(TTFont("CJK", fp, subfontIndex=0))
            FN = "CJK"
        except Exception as e:
            print(f"[BondFocus] PDF font fail: {e}")
            pdfmetrics.registerFont(UnicodeCIDFont("MSung-Light"))
    else:
        pdfmetrics.registerFont(UnicodeCIDFont("MSung-Light"))

    NAVY_ = colors.HexColor("#1F4E79")
    BLUE_ = colors.HexColor("#1F8AC0")
    DEEP_ = colors.HexColor("#1B7FA8")
    TEAL_ = colors.HexColor("#169B9B")
    GRAY_ = colors.HexColor("#595959")
    RED_ = colors.HexColor("#C00000")
    PEACH_ = colors.HexColor("#FDF2EC")

    W_ = 18.0 * cm
    st_h1 = ParagraphStyle("h1", fontName=FN, fontSize=26, leading=32, textColor=NAVY_)
    st_sec = ParagraphStyle("sec", fontName=FN, fontSize=17, leading=22, textColor=BLUE_,
                            spaceBefore=10, spaceAfter=5)
    st_secc = ParagraphStyle("secc", fontName=FN, fontSize=14, leading=18, textColor=NAVY_,
                             alignment=1, spaceBefore=6, spaceAfter=3)
    st_head = ParagraphStyle("hd", fontName=FN, fontSize=20, leading=27, textColor=colors.HexColor("#222222"))
    st_body = ParagraphStyle("bd", fontName=FN, fontSize=14.5, leading=24,
                             textColor=colors.HexColor("#333333"), spaceAfter=10)
    st_small = ParagraphStyle("sm", fontName=FN, fontSize=8.5, leading=12, textColor=GRAY_)
    st_cell = ParagraphStyle("cl", fontName=FN, fontSize=10.5, leading=14,
                             textColor=colors.HexColor("#333333"))
    st_white = ParagraphStyle("wt", fontName=FN, fontSize=16, leading=20,
                              textColor=colors.white, alignment=1)
    st_center = ParagraphStyle("ct", fontName=FN, fontSize=17, leading=22, alignment=1)
    st_center_s = ParagraphStyle("cts", fontName=FN, fontSize=10.5, leading=14,
                                 alignment=1, textColor=GRAY_)

    def sec_bar(title, color=BLUE_, center=False):
        t = Table([[Paragraph(f"<b>{title}</b>", st_secc if center else st_sec)]], colWidths=[W_])
        t.setStyle(TableStyle([("LINEBELOW", (0, 0), (-1, -1), 1.2, color),
                               ("LEFTPADDING", (0, 0), (-1, -1), 0),
                               ("BOTTOMPADDING", (0, 0), (-1, -1), 1)]))
        return t

    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            leftMargin=1.5 * cm, rightMargin=1.5 * cm,
                            topMargin=1.2 * cm, bottomMargin=1.0 * cm)
    el = []

    # ---------- P1 ----------
    mag = _magnifier_png()
    if mag:
        from reportlab.platypus import Image as _RLImg
        head_tbl = Table([[_RLImg(mag, width=1.05 * cm, height=1.05 * cm),
                           Paragraph("<b>債市每日聚焦</b>", st_h1)]],
                         colWidths=[1.35 * cm, W_ - 1.35 * cm])
        head_tbl.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                      ("LEFTPADDING", (0, 0), (-1, -1), 0),
                                      ("BOTTOMPADDING", (0, 0), (-1, -1), 0)]))
        el.append(head_tbl)
    else:
        el.append(Paragraph("<b>債市每日聚焦</b>", st_h1))
    band = Table([[Paragraph(f"<b>{D.get('date_str','')}</b>",
                             ParagraphStyle("bn", fontName=FN, fontSize=13, leading=17,
                                            textColor=colors.white, alignment=2))]], colWidths=[W_])
    band.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), BLUE_),
                              ("TOPPADDING", (0, 0), (-1, -1), 4),
                              ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                              ("RIGHTPADDING", (0, 0), (-1, -1), 6)]))
    el.append(Spacer(1, 0.15 * cm))
    el.append(band)
    el.append(Spacer(1, 0.4 * cm))

    el.append(sec_bar("焦點新聞"))
    el.append(Paragraph(f"<b>{D.get('headline','')}</b>", st_head))
    el.append(Spacer(1, 0.3 * cm))
    for t in (D.get("news_bullets") or []):
        el.append(Paragraph(f'<font color="#1F8AC0">・</font>{t}', st_body))

    el.append(Spacer(1, 0.35 * cm))
    tag = D.get("bond_tagline") or ""
    hdr_row = Table([[Paragraph("<b>焦點債券</b>", st_sec),
                      Paragraph(f"<b>{tag}</b>",
                                ParagraphStyle("tg", fontName=FN, fontSize=14, leading=19,
                                               textColor=NAVY_, alignment=2))]],
                    colWidths=[W_ * 0.42, W_ * 0.58])
    hdr_row.setStyle(TableStyle([("LINEBELOW", (0, 0), (0, 0), 1.2, BLUE_),
                                 ("LEFTPADDING", (0, 0), (-1, -1), 0),
                                 ("VALIGN", (0, 0), (-1, -1), "BOTTOM")]))
    el.append(hdr_row)
    el.append(Spacer(1, 0.2 * cm))

    rows = [["債券代碼", "債券名稱", "票面%", "YTM%", "到期日"]]
    for b in (D.get("bonds") or []):
        rows.append([b.get("code", "-"), b.get("name", "-"), str(b.get("coupon", "-")),
                     str(b.get("ytm", "-")), b.get("maturity", "-")])
    tb = Table(rows, colWidths=[4.5 * cm, 4.3 * cm, 2.3 * cm, 3.0 * cm, 3.9 * cm], repeatRows=1)
    tb_style = [("FONTNAME", (0, 0), (-1, -1), FN), ("FONTSIZE", (0, 0), (-1, -1), 13.5),
                ("BACKGROUND", (0, 0), (-1, -1), PEACH_),
                ("TEXTCOLOR", (0, 0), (-1, 0), NAVY_),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#E0D3CC")),
                ("TOPPADDING", (0, 0), (-1, -1), 10), ("BOTTOMPADDING", (0, 0), (-1, -1), 10)]
    for i in range(1, len(rows)):
        tb_style += [("TEXTCOLOR", (1, i), (1, i), TEAL_), ("TEXTCOLOR", (2, i), (2, i), RED_)]
    tb.setStyle(TableStyle(tb_style))
    el.append(tb)
    el.append(Spacer(1, 0.15 * cm))
    el.append(Paragraph("※ 報價與可承作與否以本行系統為準；商品條件依產品說明書。",
                        ParagraphStyle("sm2", fontName=FN, fontSize=10, leading=14, textColor=GRAY_)))

    el.append(PageBreak())

    # ---------- P2 ----------
    ttl = Table([[Paragraph("<b>富 邦 好 債 報</b>", st_white)]], colWidths=[6.4 * cm])
    ttl.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), DEEP_),
                             ("TOPPADDING", (0, 0), (-1, -1), 7),
                             ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                             ("ALIGN", (0, 0), (-1, -1), "CENTER")]))
    hold = Table([[ttl]], colWidths=[W_])
    hold.setStyle(TableStyle([("ALIGN", (0, 0), (-1, -1), "LEFT"),
                              ("LEFTPADDING", (0, 0), (-1, -1), 0)]))
    el.append(hold)
    el.append(Spacer(1, 0.45 * cm))
    el.append(Paragraph(f"<b>{D.get('issuer','')}</b>",
                        ParagraphStyle("iss", fontName=FN, fontSize=26, leading=32, alignment=1)))
    el.append(Paragraph(f"<b>{D.get('issuer_en','')}</b>", st_center_s))
    el.append(Spacer(1, 0.3 * cm))

    box = Table([[Paragraph(D.get("intro", ""), st_cell)]], colWidths=[W_])
    box.setStyle(TableStyle([("BOX", (0, 0), (-1, -1), 0.8, BLUE_),
                             ("TOPPADDING", (0, 0), (-1, -1), 8),
                             ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                             ("LEFTPADDING", (0, 0), (-1, -1), 8),
                             ("RIGHTPADDING", (0, 0), (-1, -1), 8)]))
    el.append(box)
    el.append(Spacer(1, 0.3 * cm))

    # 營運概況 + 營收結構（並排）
    ops_flow = []
    for blk in (D.get("ops_blocks") or []):
        ops_flow.append(Paragraph(f'<font color="#1F8AC0"><b>{blk[0]}：</b></font>', st_cell))
        ops_flow.append(Paragraph(str(blk[1]), st_cell))
    png, png_note = _chart_for_revenue_block(D)
    right_flow = []
    if png:
        try:
            from PIL import Image as _PIL
            iw, ih = _PIL.open(png).size
            w_ = 8.5 * cm
            h_ = w_ * ih / iw
            if h_ > 5.6 * cm:
                h_ = 5.6 * cm
                w_ = h_ * iw / ih
            _img = RLImage(png, width=w_, height=h_)
            _img.hAlign = "CENTER"
            right_flow.append(_img)
            if png_note:
                right_flow.append(Paragraph(png_note, st_center_s))
        except Exception as e:
            print(f"[BondFocus] pdf donut embed: {e}")
    else:
        right_flow.append(Paragraph("（暫無可用圖表資料）", st_center_s))

    heads = Table([[Paragraph("<b>營運概況</b>", st_secc),
                    Paragraph(f"<b>{_revenue_block_title(D)}</b>", st_secc)]],
                  colWidths=[W_ / 2, W_ / 2])
    heads.setStyle(TableStyle([("LINEBELOW", (0, 0), (-1, -1), 1.2, BLUE_),
                               ("LEFTPADDING", (0, 0), (-1, -1), 0)]))
    el.append(heads)
    two = Table([[ops_flow, right_flow]], colWidths=[W_ / 2 - 0.2 * cm, W_ / 2 - 0.2 * cm])
    two.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"),
                             ("BOX", (0, 0), (0, 0), 0.8, BLUE_),
                             ("BOX", (1, 0), (1, 0), 0.8, BLUE_),
                             ("TOPPADDING", (0, 0), (-1, -1), 8),
                             ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                             ("LEFTPADDING", (0, 0), (-1, -1), 7),
                             ("RIGHTPADDING", (0, 0), (-1, -1), 7)]))
    el.append(two)

    el.append(sec_bar("信用評等", center=True))
    R = D.get("ratings") or {}
    r2 = [["信用評等", "穆迪", "標普", "惠譽"],
          ["長期評等", R.get("moody") or "--", R.get("sp") or "--", R.get("fitch") or "--"],
          ["評等展望", _zh_outlook(R.get("moody_outlook")), _zh_outlook(R.get("sp_outlook")), _zh_outlook(R.get("fitch_outlook"))],
          ["最近評等動作", R.get("moody_date") or "--", R.get("sp_date") or "--", R.get("fitch_date") or "--"]]
    t2 = Table(r2, colWidths=[5.0 * cm, 4.3 * cm, 4.3 * cm, 4.4 * cm])
    t2.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), FN), ("FONTSIZE", (0, 0), (-1, -1), 10.5),
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EAF3FA")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), NAVY_),
                            ("ALIGN", (0, 0), (-1, -1), "CENTER"), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D9D9D9")),
                            ("TOPPADDING", (0, 0), (-1, -1), 6), ("BOTTOMPADDING", (0, 0), (-1, -1), 6)]))
    el.append(t2)

    el.append(sec_bar("信評公司評析", center=True))
    ac_flow = []
    for c in (D.get("agency_comments") or []):
        ac_flow.append(Paragraph(f'<font color="#1F8AC0"><b>{c[0]}－</b></font>{c[1]}', st_cell))
    acbox = Table([[ac_flow]], colWidths=[W_])
    acbox.setStyle(TableStyle([("BOX", (0, 0), (-1, -1), 0.8, BLUE_),
                               ("TOPPADDING", (0, 0), (-1, -1), 8),
                               ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                               ("LEFTPADDING", (0, 0), (-1, -1), 8),
                               ("RIGHTPADDING", (0, 0), (-1, -1), 8)]))
    el.append(acbox)
    el.append(Spacer(1, 0.2 * cm))
    el.append(Paragraph(D.get("source_note", ""), st_small))
    el.append(Paragraph("「本內容僅供參考且不構成要約或要約引誘。特定標的之商品風險、申購之條件、限制與費用及"
                        "其他相關權利義務，應依產品說明暨投資風險預告書等相關文件為準。」", st_small))

    def _footer(canv, doc_):
        canv.saveState()
        canv.setFont(FN, 10)
        canv.setFillColor(RED_)
        canv.drawString(1.5 * cm, 1.0 * cm, "僅限內部教育訓練使用")
        canv.setFillColor(NAVY_)
        canv.drawRightString(A4[0] - 1.5 * cm, 1.0 * cm, "台北富邦銀行")
        canv.restoreState()

    doc.build(el, onFirstPage=_footer, onLaterPages=_footer)
    if mag:
        try:
            os.remove(mag)
        except Exception:
            pass
    if png:
        try:
            os.remove(png)   # 需等 build 完成後再刪,reportlab 延後讀檔
        except Exception:
            pass
    return out_path
